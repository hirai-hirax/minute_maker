"""
Document parser module for extracting text from various file formats.
Supports: Word (.docx), Excel (.xlsx), PowerPoint (.pptx), PDF (.pdf), and plain text (.txt)

All libraries used are MIT-licensed:
- python-docx: MIT License
- openpyxl: MIT License
- python-pptx: MIT License
- PyPDF2: MIT License
"""

import logging
from io import BytesIO
from pathlib import Path
from typing import Optional

from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation
from PyPDF2 import PdfReader

logger = logging.getLogger(__name__)


def extract_text_from_docx(file_bytes: bytes) -> str:
    """
    Extract text from a Word (.docx) document.
    
    Args:
        file_bytes: Raw bytes of the .docx file
        
    Returns:
        Extracted text as a string
    """
    try:
        doc = DocxDocument(BytesIO(file_bytes))
        paragraphs = [paragraph.text for paragraph in doc.paragraphs if paragraph.text.strip()]
        
        # Also extract text from tables
        table_texts = []
        for table in doc.tables:
            for row in table.rows:
                row_text = ' | '.join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    table_texts.append(row_text)
        
        all_text = '\n'.join(paragraphs)
        if table_texts:
            all_text += '\n\n' + '\n'.join(table_texts)
        
        logger.info(f"Extracted {len(all_text)} characters from Word document")
        return all_text
    except Exception as e:
        logger.error(f"Failed to extract text from Word document: {e}")
        raise ValueError(f"Word document parsing failed: {str(e)}")


def extract_text_from_xlsx(file_bytes: bytes) -> str:
    """
    Extract text from an Excel (.xlsx) spreadsheet.
    
    Args:
        file_bytes: Raw bytes of the .xlsx file
        
    Returns:
        Extracted text as a string
    """
    try:
        workbook = load_workbook(BytesIO(file_bytes), read_only=True, data_only=True)
        all_text = []
        
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            all_text.append(f"=== Sheet: {sheet_name} ===")
            
            for row in sheet.iter_rows(values_only=True):
                # Filter out None values and convert to strings
                row_values = [str(cell) for cell in row if cell is not None and str(cell).strip()]
                if row_values:
                    all_text.append(' | '.join(row_values))
        
        result = '\n'.join(all_text)
        logger.info(f"Extracted {len(result)} characters from Excel document ({len(workbook.sheetnames)} sheets)")
        return result
    except Exception as e:
        logger.error(f"Failed to extract text from Excel document: {e}")
        raise ValueError(f"Excel document parsing failed: {str(e)}")


def extract_text_from_pptx(file_bytes: bytes) -> str:
    """
    Extract text from a PowerPoint (.pptx) presentation.
    
    Args:
        file_bytes: Raw bytes of the .pptx file
        
    Returns:
        Extracted text as a string
    """
    try:
        presentation = Presentation(BytesIO(file_bytes))
        all_text = []
        
        for slide_num, slide in enumerate(presentation.slides, start=1):
            slide_texts = []
            
            # Extract text from all shapes in the slide
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())
            
            if slide_texts:
                all_text.append(f"=== Slide {slide_num} ===")
                all_text.extend(slide_texts)
        
        result = '\n'.join(all_text)
        logger.info(f"Extracted {len(result)} characters from PowerPoint ({len(presentation.slides)} slides)")
        return result
    except Exception as e:
        logger.error(f"Failed to extract text from PowerPoint: {e}")
        raise ValueError(f"PowerPoint parsing failed: {str(e)}")


def extract_text_from_pdf(file_bytes: bytes) -> str:
    """
    Extract text from a PDF document.
    
    Args:
        file_bytes: Raw bytes of the PDF file
        
    Returns:
        Extracted text as a string
    """
    try:
        reader = PdfReader(BytesIO(file_bytes))
        all_text = []
        
        for page_num, page in enumerate(reader.pages, start=1):
            text = page.extract_text()
            if text.strip():
                all_text.append(f"=== Page {page_num} ===")
                all_text.append(text.strip())
        
        result = '\n'.join(all_text)
        logger.info(f"Extracted {len(result)} characters from PDF ({len(reader.pages)} pages)")
        return result
    except Exception as e:
        logger.error(f"Failed to extract text from PDF: {e}")
        raise ValueError(f"PDF parsing failed: {str(e)}")


def extract_text_from_txt(file_bytes: bytes) -> str:
    """
    Extract text from a plain text (.txt) file.
    
    Args:
        file_bytes: Raw bytes of the .txt file
        
    Returns:
        Extracted text as a string
    """
    try:
        # Try UTF-8 first, then fallback to other encodings
        encodings = ['utf-8', 'shift-jis', 'cp932', 'iso-2022-jp', 'euc-jp', 'latin-1']
        
        for encoding in encodings:
            try:
                text = file_bytes.decode(encoding)
                logger.info(f"Extracted {len(text)} characters from text file (encoding: {encoding})")
                return text
            except UnicodeDecodeError:
                continue
        
        # If all encodings fail, use utf-8 with error handling
        text = file_bytes.decode('utf-8', errors='replace')
        logger.warning(f"Used fallback decoding for text file with character replacement")
        return text
    except Exception as e:
        logger.error(f"Failed to extract text from text file: {e}")
        raise ValueError(f"Text file reading failed: {str(e)}")


def extract_text_from_file(filename: str, file_bytes: bytes) -> str:
    """
    Extract text from a file based on its extension.
    
    Args:
        filename: Name of the file (used to determine file type)
        file_bytes: Raw bytes of the file
        
    Returns:
        Extracted text as a string
        
    Raises:
        ValueError: If file type is not supported or parsing fails
    """
    extension = Path(filename).suffix.lower()
    
    extractors = {
        '.docx': extract_text_from_docx,
        '.xlsx': extract_text_from_xlsx,
        '.pptx': extract_text_from_pptx,
        '.pdf': extract_text_from_pdf,
        '.txt': extract_text_from_txt,
    }
    
    extractor = extractors.get(extension)
    if not extractor:
        supported = ', '.join(extractors.keys())
        raise ValueError(f"Unsupported file type: {extension}. Supported types: {supported}")
    
    logger.info(f"Extracting text from {filename} (type: {extension})")
    return extractor(file_bytes)
