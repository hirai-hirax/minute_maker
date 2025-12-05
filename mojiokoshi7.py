import os
os.environ["KMP_DUPLICATE_LIB_OK"] = "TRUE"
import streamlit as st
from openai import AzureOpenAI
import tempfile
from io import BytesIO
from pydub import AudioSegment
import fitz
import torch
import pandas as pd
from dotenv import load_dotenv
from resemblyzer import VoiceEncoder, preprocess_wav
import numpy as np
import zipfile
from datetime import timedelta, datetime
from concurrent.futures import ThreadPoolExecutor, as_completed
import threading
import json
import base64
import uuid
import subprocess
from docx import Document as DocxDocument
from pptx import Presentation
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_openai import AzureOpenAIEmbeddings
from langchain_community.vectorstores import FAISS
from langchain.schema import Document
from pathlib import Path
from streamlit.components.v1 import html as components_html

torch.classes.__path__ = []

# .envファイルから環境変数を読み込む
load_dotenv()

# 環境変数から設定を取得（Azure OpenAI のエンドポイント・API キーを設定してください）
AZURE_OPENAI_ENDPOINT = os.environ.get("AZURE_OPENAI_ENDPOINT")
AZURE_OPENAI_API_KEY = os.environ.get("AZURE_OPENAI_API_KEY")
API_VERSION = "2025-03-01-preview"  # gpt-4o-transcribe モデルに対応したバージョン

# RAGデータベースの保存先フォルダ設定
DEFAULT_RAGDB_FOLDER = ""  # 空文字列の場合はカレントディレクトリ（デフォルト）
# 例: DEFAULT_RAGDB_FOLDER = "C:/Users/username/Documents/ragdb"
# 例: DEFAULT_RAGDB_FOLDER = "./data/ragdb"

PROMPTS_DATA = {
    "rag_proofreading": {
        "presets": {
            "standard": {
                "name": "標準校正",
                "description": "バランスの取れた標準的な校正",
                "system_prompt_with_context": ("""
        あなたは議事録校正の専門家です。以下の関連資料を参照して、議事録を校正してください。

        【参照資料】
        {context}

        【校正指示】
        1. 誤字脱字の修正
        2. 文法の改善
        3. 専門用語の正確性確認
        4. 文脈に基づく内容の補完
        5. 読みやすさの向上

        参照資料を活用しながら、正確で読みやすい議事録に校正してください。
"""
               ),
                "system_prompt_without_context": (
                    "あなたは議事録校正の専門家です。以下の議事録を校正してください。\n\n"
                    "【校正指示】\n"
                    "1. 誤字脱字の修正\n"
                    "2. 文法の改善\n"
                    "3. 読みやすさの向上\n\n"
                    "正確で読みやすい議事録に校正してください。"
                ),
            },
            "detailed": {
                "name": "詳細校正",
                "description": "より詳細な分析と改善提案を含む校正",
                "system_prompt_with_context": (
                    "あなたは議事録校正の専門家です。以下の関連資料を参照して、議事録を詳細に校正してください。\n\n"
                    "【参照資料】\n"
                    "{context}\n\n"
                    "【詳細校正指示】\n"
                    "1. 誤字脱字の修正と詳細な説明\n"
                    "2. 文法の改善と代替表現の提案\n"
                    "3. 専門用語の正確性確認と定義の補足\n"
                    "4. 文脈に基づく内容の補完と背景情報の追加\n"
                    "5. 読みやすさの向上と文章構造の最適化\n"
                    "6. 曖昧な表現の明確化\n"
                    "7. 重要なポイントの強調\n\n"
                    "参照資料を十分に活用し、より詳細で分かりやすい議事録に校正してください。\n"
                    "必要に応じて、改善点の説明も含めてください。"
                ),
                "system_prompt_without_context": (
                    "あなたは議事録校正の専門家です。以下の議事録を詳細に校正してください。\n\n"
                    "【詳細校正指示】\n"
                    "1. 誤字脱字の修正と詳細な説明\n"
                    "2. 文法の改善と代替表現の提案\n"
                    "3. 読みやすさの向上と文章構造の最適化\n"
                    "4. 曖昧な表現の明確化\n"
                    "5. 重要なポイントの強調\n\n"
                    "より詳細で分かりやすい議事録に校正してください。"
                ),
            },
            "simple": {
                "name": "簡潔校正",
                "description": "必要最小限の修正のみを行う簡潔な校正",
                "system_prompt_with_context": (
                    "あなたは議事録校正の専門家です。以下の関連資料を参照して、議事録を簡潔に校正してください。\n\n"
                    "【参照資料】\n"
                    "{context}\n\n"
                    "【簡潔校正指示】\n"
                    "1. 明らかな誤字脱字のみ修正\n"
                    "2. 重大な文法エラーのみ修正\n"
                    "3. 元の文章をできるだけ維持\n\n"
                    "参照資料を活用しながら、最小限の修正で読みやすい議事録に校正してください。"
                ),
                "system_prompt_without_context": (
                    "あなたは議事録校正の専門家です。以下の議事録を簡潔に校正してください。\n\n"
                    "【簡潔校正指示】\n"
                    "1. 明らかな誤字脱字のみ修正\n"
                    "2. 重大な文法エラーのみ修正\n"
                    "3. 元の文章をできるだけ維持\n\n"
                    "最小限の修正で読みやすい議事録に校正してください。"
                ),
            },
            "formal": {
                "name": "フォーマル校正",
                "description": "よりフォーマルで丁寧な表現に校正",
                "system_prompt_with_context": (
                    "あなたは議事録校正の専門家です。以下の関連資料を参照して、議事録をよりフォーマルな表現に校正してください。\n\n"
                    "【参照資料】\n"
                    "{context}\n\n"
                    "【フォーマル校正指示】\n"
                    "1. 誤字脱字の修正\n"
                    "2. 文法の改善\n"
                    "3. カジュアルな表現をフォーマルな表現に変換\n"
                    "4. 敬語表現の統一と適切な使用\n"
                    "5. ビジネス文書として適切な語彙の使用\n"
                    "6. 専門用語の正確性確認\n\n"
                    "参照資料を活用しながら、フォーマルで丁寧な印象の議事録に校正してください。"
                ),
                "system_prompt_without_context": (
                    "あなたは議事録校正の専門家です。以下の議事録をよりフォーマルな表現に校正してください。\n\n"
                    "【フォーマル校正指示】\n"
                    "1. 誤字脱字の修正\n"
                    "2. 文法の改善\n"
                    "3. カジュアルな表現をフォーマルな表現に変換\n"
                    "4. 敬語表現の統一と適切な使用\n"
                    "5. ビジネス文書として適切な語彙の使用\n\n"
                    "フォーマルで丁寧な印象の議事録に校正してください。"
                ),
            },
        },
    },
    "legacy": {
        "initial_transcription_prompt": (
            "\"こんにちは。\\n\\nはい、こんにちは。\\n\\nお元気ですか？\\n\\nはい、元気です。\\n\\nそれは何よりです。では早速始めましょう。\\n\\nはい、よろしくお願いいたします。\""
        ),
        "summarizing_prompt1": (
            "ユーザーからテキストを渡されます。当該のテキストの内容を読んだ上で、150文字程度の要約を生成してください。"
        ),
    },
}

DEFAULT_MEETING_TYPES = [
    {
        "name": "経営会議",
        "id": "executive_meeting",
        "embeddings_folder": "speaker_embeddings/executive",
        "description": "月例経営会議（役員メンバー固定）",
    },
    {
        "name": "開発チーム定例",
        "id": "dev_team_meeting",
        "embeddings_folder": "speaker_embeddings/dev_team",
        "description": "週次開発チーム会議",
    },
    {
        "name": "営業定例",
        "id": "sales_meeting",
        "embeddings_folder": "speaker_embeddings/sales",
        "description": "月次営業会議",
    },
    {
        "name": "カスタム（手動選択）",
        "id": "custom",
        "embeddings_folder": "",
        "description": "会議タイプを指定せず、手動で話者埋め込みを選択",
    },
]


class _InlinePromptLoader:
    def __init__(self, data: dict):
        self._cache = data

    def get_prompt(self, *keys: str) -> str:
        data = self._cache
        for key in keys:
            if not isinstance(data, dict) or key not in data:
                raise KeyError(f"プロンプトが見つかりません: {' -> '.join(keys)}")
            data = data[key]
        if not isinstance(data, str):
            raise ValueError(f"指定されたキーパスはプロンプト文字列ではありません: {' -> '.join(keys)}")
        return data

    def get_all_prompts(self) -> dict:
        return json.loads(json.dumps(self._cache))

    def get_presets(self, category: str) -> dict:
        category_data = self._cache.get(category, {})
        return category_data.get("presets", {}) if isinstance(category_data, dict) else {}

    def get_preset_list(self, category: str) -> list:
        presets = self.get_presets(category)
        return [
            {
                "id": preset_id,
                "name": preset_data.get("name", preset_id),
                "description": preset_data.get("description", ""),
            }
            for preset_id, preset_data in presets.items()
        ]

    def get_prompt_from_preset(self, category: str, preset_id: str, prompt_type: str) -> str:
        presets = self.get_presets(category)
        if preset_id not in presets:
            raise KeyError(f"プロンプトプリセットが見つかりません: {category} -> {preset_id}")
        preset_data = presets[preset_id]
        if prompt_type not in preset_data:
            raise KeyError(f"プロンプトタイプが見つかりません: {category} -> {preset_id} -> {prompt_type}")
        value = preset_data[prompt_type]
        if not isinstance(value, str):
            raise ValueError(f"指定されたプリセット値は文字列ではありません: {category} -> {preset_id} -> {prompt_type}")
        return value


_DEFAULT_PROMPT_LOADER: _InlinePromptLoader | None = None


def get_default_loader() -> _InlinePromptLoader:
    global _DEFAULT_PROMPT_LOADER
    if _DEFAULT_PROMPT_LOADER is None:
        _DEFAULT_PROMPT_LOADER = _InlinePromptLoader(PROMPTS_DATA)
    return _DEFAULT_PROMPT_LOADER


def get_prompt(*keys: str) -> str:
    return get_default_loader().get_prompt(*keys)

def _extract_pdf(file: BytesIO) -> str:
    file_bytes = file.read()
    file.seek(0)
    pdf_document = fitz.open(stream=file_bytes, filetype='pdf')
    text = "".join(page.get_text() for page in pdf_document)
    pdf_document.close()
    return text

def _extract_docx(file: BytesIO) -> str:
    doc = DocxDocument(file)
    text = "\n".join(p.text for p in doc.paragraphs)
    for table in doc.tables:
        text += "\n" + "\n".join(" ".join(cell.text for cell in row.cells) for row in table.rows)
    return text.strip()

def _extract_pptx(file: BytesIO) -> str:
    presentation = Presentation(file)
    text_parts = []
    for slide_num, slide in enumerate(presentation.slides, 1):
        text_parts.append(f"\n--- スライド {slide_num} ---\n")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text_parts.append(shape.text + "\n")
            if shape.has_table:
                for row in shape.table.rows:
                    text_parts.append(" ".join(cell.text for cell in row.cells) + "\n")
    return "".join(text_parts).strip()

def _extract_msg(file: BytesIO) -> str:
    try:
        import extract_msg
    except ImportError:
        raise Exception("MSGファイルの処理には extract-msg ライブラリが必要です。pip install extract-msg を実行してください。")

    with tempfile.NamedTemporaryFile(delete=False, suffix=".msg") as temp_file:
        temp_file.write(file.read())
        temp_path = temp_file.name

    try:
        msg = extract_msg.Message(temp_path)
        parts = []
        if msg.subject:
            parts.append(f"件名: {msg.subject}\n")
        if msg.sender:
            parts.append(f"送信者: {msg.sender}\n")
        if msg.to:
            parts.append(f"宛先: {msg.to}\n")
        if msg.body:
            parts.append(f"\n{msg.body}")
        msg.close()
        return "\n".join(parts).strip()
    finally:
        if os.path.exists(temp_path):
            os.unlink(temp_path)

def _extract_txt(file: BytesIO) -> str:
    file_content = file.read()
    for encoding in ['utf-8', 'cp932', 'shift_jis', 'utf-16', 'iso-2022-jp']:
        try:
            return file_content.decode(encoding).strip()
        except UnicodeDecodeError:
            continue
    raise Exception("テキストファイルのエンコーディングを判定できませんでした")

FILE_EXTRACTORS = {
    'pdf': _extract_pdf,
    'docx': _extract_docx,
    'pptx': _extract_pptx,
    'txt': _extract_txt,
    'msg': _extract_msg,
}

def extract_text_from_file(file, file_extension):
    """ファイル形式に応じてテキストを抽出"""
    extractor = FILE_EXTRACTORS.get(file_extension.lower())
    if not extractor:
        raise Exception(f"サポートされていないファイル形式: {file_extension}")
    try:
        return extractor(file)
    except Exception as e:
        raise Exception(f"{file_extension.upper()}読み込みエラー: {e}")

def _create_azure_client():
    """Azure OpenAIクライアントを作成"""
    return AzureOpenAI(
        azure_endpoint=AZURE_OPENAI_ENDPOINT,
        api_key=AZURE_OPENAI_API_KEY,
        api_version=API_VERSION,
    )

def generate_summary(model, prompt, text):
    """テキストを要約"""
    client = _create_azure_client()
    response = client.chat.completions.create(
        model=model,
        messages=[
            {"role": "system", "content": prompt},
            {"role": "user", "content": text},
        ],
    )
    print(f"Response: {response}")
    return response.choices[0].message.content

from contextlib import contextmanager

@contextmanager
def temp_file_path(data: bytes, suffix: str):
    """一時ファイルを作成し、パスを返すコンテキストマネージャ"""
    with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
        tmp.write(data)
        tmp_path = tmp.name
    try:
        yield tmp_path
    finally:
        if os.path.exists(tmp_path):
            os.unlink(tmp_path)


def trigger_auto_download(data: bytes, file_name: str, key: str | None, mime: str = "application/octet-stream"):
    """Streamlitコンポーネントを使って即時ダウンロードを開始"""
    if not data:
        return

    encoded = base64.b64encode(data).decode()
    mime_js = json.dumps(mime)
    file_name_js = json.dumps(file_name)
    element_id_js = json.dumps(key or f"dl_{uuid.uuid4().hex}")

    components_html(
        f"""
        <div id={element_id_js}></div>
        <script>
            (function() {{
                const mimeType = {mime_js};
                const fileName = {file_name_js};
                const link = document.createElement('a');
                link.href = "data:" + mimeType + ";base64,{encoded}";
                link.download = fileName;
                link.style.display = 'none';
                document.body.appendChild(link);
                link.click();
                document.body.removeChild(link);
                const element = document.getElementById({element_id_js});
                if (element) {{
                    element.remove();
                }}
            }})();
        </script>
        """,
        height=0,
    )


def transcribe_audio_to_dataframe(uploaded_file: BytesIO, reference_file: BytesIO = None, model: str = "gpt-4o-transcribe-diarize"):
    """音声ファイル全体を文字起こし（モデル選択可能、25MB制限対応）

    Args:
        uploaded_file: 音声ファイル
        reference_file: 参考資料（Whisperのみサポート）
        model: 使用するモデル ("gpt-4o-transcribe-diarize" または "whisper")
    """
    # ファイルサイズをチェック（25MB = 26,214,400 bytes）
    MAX_FILE_SIZE = 25 * 1024 * 1024
    uploaded_file.seek(0)
    file_size = len(uploaded_file.getvalue())
    uploaded_file.seek(0)

    # 25MB以上の場合は分割処理
    if file_size > MAX_FILE_SIZE:
        st.warning(f"⚠️ ファイルサイズが {file_size / (1024*1024):.1f}MB です。25MB制限のため、音声を分割して処理します。")
        return _transcribe_large_audio_chunked(uploaded_file, reference_file, model)

    # 25MB未満の場合は通常処理
    return _transcribe_audio_single(uploaded_file, reference_file, model)

def _transcribe_large_audio_chunked(uploaded_file: BytesIO, reference_file: BytesIO = None, model: str = "gpt-4o-transcribe-diarize"):
    """大きな音声ファイルを分割して文字起こし

    Args:
        uploaded_file: 音声ファイル
        reference_file: 参考資料
        model: 使用するモデル

    Returns:
        pd.DataFrame: 文字起こし結果
    """
    try:
        suffix = f".{uploaded_file.name.split('.')[-1]}"

        # 音声ファイルを読み込み
        with temp_file_path(uploaded_file.getvalue(), suffix) as tmp_path:
            audio = AudioSegment.from_file(tmp_path)

        # 音声の長さ（ミリ秒）
        total_duration_ms = len(audio)
        total_duration_sec = total_duration_ms / 1000

        # チャンクサイズを決定（10分 = 600秒 = 600,000ミリ秒）
        CHUNK_DURATION_MS = 10 * 60 * 1000  # 10分

        # チャンク数を計算
        num_chunks = (total_duration_ms + CHUNK_DURATION_MS - 1) // CHUNK_DURATION_MS

        st.info(f"📊 音声長: {total_duration_sec/60:.1f}分、{num_chunks}個のチャンクに分割して処理します")

        # 進捗表示
        progress_bar = st.progress(0)
        status_text = st.empty()

        all_segments = []

        for i in range(num_chunks):
            start_ms = i * CHUNK_DURATION_MS
            end_ms = min((i + 1) * CHUNK_DURATION_MS, total_duration_ms)

            status_text.write(f"🎤 チャンク {i+1}/{num_chunks} を処理中... ({start_ms/1000:.1f}秒 ～ {end_ms/1000:.1f}秒)")

            # チャンクを抽出
            chunk = audio[start_ms:end_ms]

            # 一時ファイルとして保存
            chunk_io = BytesIO()
            chunk.export(chunk_io, format="mp3", bitrate="192k")
            chunk_io.seek(0)
            chunk_io.name = f"chunk_{i}.mp3"

            # チャンクを文字起こし
            try:
                chunk_df = _transcribe_audio_single(chunk_io, reference_file, model)

                # タイムスタンプをオフセット調整
                if not chunk_df.empty and 'start' in chunk_df.columns and 'end' in chunk_df.columns:
                    offset_sec = start_ms / 1000
                    chunk_df['start'] = chunk_df['start'] + offset_sec
                    chunk_df['end'] = chunk_df['end'] + offset_sec

                all_segments.append(chunk_df)

            except Exception as e:
                st.error(f"❌ チャンク {i+1} の処理中にエラー: {e}")

            progress_bar.progress((i + 1) / num_chunks)

        # すべてのチャンクを結合
        if all_segments:
            result_df = pd.concat(all_segments, ignore_index=True)
            status_text.write(f"✅ 分割処理完了！合計 {len(result_df)} 個のセグメント")
            return result_df
        else:
            st.error("❌ すべてのチャンクの処理に失敗しました")
            return pd.DataFrame(columns=["start", "end", "speaker", "text"])

    except Exception as e:
        st.error(f"❌ 分割処理中にエラーが発生しました: {e}")
        import traceback
        st.error(traceback.format_exc())
        return pd.DataFrame(columns=["start", "end", "speaker", "text"])

def _transcribe_audio_single(uploaded_file: BytesIO, reference_file: BytesIO = None, model: str = "gpt-4o-transcribe-diarize"):
    """単一の音声ファイルを文字起こし（内部用）

    Args:
        uploaded_file: 音声ファイル
        reference_file: 参考資料（Whisperのみサポート）
        model: 使用するモデル ("gpt-4o-transcribe-diarize" または "whisper")
    """
    # モデルに応じたAPIバージョンを選択
    if model == "whisper":
        api_version = "2024-06-01"  # Whisper用の安定版APIバージョン
        st.info(f"🔧 Whisperモデル用にAPIバージョン {api_version} を使用します")
    else:
        api_version = "2025-03-01-preview"  # gpt-4o-transcribe-diarize用

    # モデル専用のクライアントを作成
    client = AzureOpenAI(
        azure_endpoint=AZURE_OPENAI_ENDPOINT,
        api_key=AZURE_OPENAI_API_KEY,
        api_version=api_version,
    )

    try:
        suffix = f".{uploaded_file.name.split('.')[-1]}"

        if model == "gpt-4o-transcribe-diarize":
            # gpt-4o-transcribe-diarizeモデル（話者識別付き）
            if reference_file:
                st.warning("⚠️ gpt-4o-transcribe-diarizeモデルは参考資料（prompt）をサポートしていません。参考資料は無視されます。")

            st.info("文字起こしを開始します（gpt-4o-transcribe-diarize、話者識別付き）...")

            with temp_file_path(uploaded_file.getvalue(), suffix) as tmp_path:
                with open(tmp_path, "rb") as audio_file:
                    transcript = client.audio.transcriptions.create(
                        model="gpt-4o-transcribe-diarize",
                        file=(uploaded_file.name, audio_file, f"audio/{uploaded_file.name.split('.')[-1]}"),
                        response_format="diarized_json",
                        chunking_strategy="auto"
                    )

            # レスポンスからセグメントを取得
            transcript_dict = transcript.model_dump()
            segments = transcript_dict.get("segments", [])

            if segments:
                # セグメントをデータフレームに変換
                seg_list = []
                for seg in segments:
                    seg_list.append({
                        "start": seg.get("start", 0),
                        "end": seg.get("end", 0),
                        "speaker": seg.get("speaker", ""),
                        "text": seg.get("text", "")
                    })

                st.success(f"文字起こしが完了しました！（{len(seg_list)}個のセグメント、話者識別付き）")
                seg_df = pd.DataFrame(seg_list)
                return seg_df
            else:
                # segmentsがない場合はテキスト全体を取得
                text = transcript_dict.get("text", "")
                if text:
                    st.warning("⚠️ セグメント情報が取得できませんでした。全体を1つのセグメントとして扱います。")
                    seg_df = pd.DataFrame([{
                        "start": 0,
                        "end": 0,
                        "speaker": "",
                        "text": text
                    }])
                    return seg_df
                else:
                    st.error("❌ 文字起こし結果が空でした")
                    return pd.DataFrame(columns=["start", "end", "speaker", "text"])

        elif model == "whisper":
            # Whisperモデル（話者識別なし）
            st.info("文字起こしを開始します（Whisper、話者識別なし）...")

            with temp_file_path(uploaded_file.getvalue(), suffix) as tmp_path:
                with open(tmp_path, "rb") as audio_file:
                    # promptパラメータの準備
                    kwargs = {
                        "model": "whisper",
                        "file": (uploaded_file.name, audio_file, f"audio/{uploaded_file.name.split('.')[-1]}"),
                        "response_format": "verbose_json"
                    }

                    # 参考資料がある場合はpromptとして使用
                    if reference_file:
                        try:
                            file_extension = reference_file.name.split(".")[-1].lower()
                            reference_text = extract_text_from_file(BytesIO(reference_file.read()), file_extension)
                            if reference_text:
                                # promptは最大224トークン程度に制限（約1000文字）
                                kwargs["prompt"] = reference_text[:1000]
                                st.info("✅ 参考資料をpromptとして使用します")
                        except Exception as e:
                            st.warning(f"⚠️ 参考資料の読み込みに失敗しました: {e}")

                    transcript = client.audio.transcriptions.create(**kwargs)

            # レスポンスからセグメントを取得
            transcript_dict = transcript.model_dump()
            segments = transcript_dict.get("segments", [])

            if segments:
                # セグメントをデータフレームに変換（話者情報は空）
                seg_list = []
                for seg in segments:
                    seg_list.append({
                        "start": seg.get("start", 0),
                        "end": seg.get("end", 0),
                        "speaker": "",  # Whisperは話者識別をサポートしていない
                        "text": seg.get("text", "")
                    })

                st.success(f"文字起こしが完了しました！（{len(seg_list)}個のセグメント、Whisper使用）")
                seg_df = pd.DataFrame(seg_list)
                return seg_df
            else:
                # segmentsがない場合はテキスト全体を取得
                text = transcript_dict.get("text", "")
                if text:
                    st.warning("⚠️ セグメント情報が取得できませんでした。全体を1つのセグメントとして扱います。")
                    seg_df = pd.DataFrame([{
                        "start": 0,
                        "end": 0,
                        "speaker": "",
                        "text": text
                    }])
                    return seg_df
                else:
                    st.error("❌ 文字起こし結果が空でした")
                    return pd.DataFrame(columns=["start", "end", "speaker", "text"])
        else:
            st.error(f"❌ サポートされていないモデル: {model}")
            return pd.DataFrame(columns=["start", "end", "speaker", "text"])

    except Exception as e:
        st.error(f"❌ 文字起こし中にエラーが発生しました")
        st.error(f"**エラー詳細**: {str(e)}")
        st.error(f"**使用モデル**: {model}")
        st.error(f"**APIバージョン**: {api_version}")

        import traceback
        error_details = traceback.format_exc()

        with st.expander("🔍 詳細なエラー情報（デバッグ用）"):
            st.code(error_details, language="python")

        # エラーをコンソールにも出力（デバッグ用）
        print(f"===== Transcription Error =====")
        print(f"Model: {model}")
        print(f"API Version: {api_version}")
        print(f"Error: {e}")
        print(error_details)
        print(f"==============================")

        return pd.DataFrame(columns=["start", "end", "speaker", "text"])

@st.cache_resource
def load_voice_encoder():
    """Caches the VoiceEncoder model."""
    return VoiceEncoder()

def extract_embedding(audio_content):
    """Extracts embedding from audio content."""
    with temp_file_path(audio_content.read(), ".wav") as wav_path:
        wav = preprocess_wav(wav_path)
        encoder = load_voice_encoder()
        return encoder.embed_utterance(wav)

def load_speaker_embeddings_from_files(uploaded_files):
    """Loads known speaker embeddings from uploaded files.

    Args:
        uploaded_files: リストまたはタプル。各要素は以下のいずれか：
            - Streamlit UploadedFile オブジェクト（file_uploaderから）
            - 辞書 {'name': filename, 'data': file_bytes}（フォルダから読み込んだ場合）

    Returns:
        dict: {speaker_name: embedding_array}
    """
    if not uploaded_files:
        st.warning("話者埋め込みファイルがアップロードされていません。")
        return {}

    speaker_embeddings = {}
    for uploaded_file in uploaded_files:
        try:
            # 辞書形式（フォルダから読み込んだ場合）かUploadedFileオブジェクトかを判定
            if isinstance(uploaded_file, dict):
                # 辞書形式: {'name': filename, 'data': file_bytes}
                filename = uploaded_file['name']
                file_data = BytesIO(uploaded_file['data'])
            else:
                # UploadedFile オブジェクト
                filename = uploaded_file.name
                file_data = uploaded_file

            filename_stem = Path(filename).stem
            # ファイル名の区切り文字で話者名を抽出
            for delimiter in ['‗', '_']:
                if delimiter in filename_stem:
                    filename_stem = filename_stem.split(delimiter)[0]
                    break
            speaker_name = filename_stem.strip()

            if not speaker_name:
                speaker_name = Path(filename).stem

            speaker_embeddings[speaker_name] = np.load(file_data)
        except Exception as e:
            filename_str = uploaded_file.get('name') if isinstance(uploaded_file, dict) else getattr(uploaded_file, 'name', 'unknown')
            st.error(f"埋め込みファイルの読み込み中にエラーが発生しました {filename_str}: {e}")
    return speaker_embeddings

def _calculate_similarity(embedding1, embedding2):
    """2つの埋め込みベクトルの類似度を計算"""
    return np.dot(embedding1, embedding2) / (np.linalg.norm(embedding1) * np.linalg.norm(embedding2))

def _identify_speaker(segment_embedding, known_embeddings, threshold):
    """セグメント埋め込みから話者を識別"""
    best_speaker, best_similarity = None, -1
    for speaker_name, known_embedding in known_embeddings.items():
        similarity = _calculate_similarity(segment_embedding, known_embedding)
        print(f"Comparing segment with {speaker_name}: similarity = {similarity:.4f}")
        if similarity > best_similarity:
            best_similarity, best_speaker = similarity, speaker_name
    return best_speaker if best_similarity >= threshold else ""

def identify_speakers_in_dataframe(audio_file, df: pd.DataFrame, uploaded_embedding_files, similarity_threshold: float) -> pd.DataFrame:
    known_embeddings = load_speaker_embeddings_from_files(uploaded_embedding_files)
    if not known_embeddings:
        st.warning("既知の話者埋め込みが見つかりませんでした。識別を実行できません。")
        df['speaker'] = None
        return df

    st.info(f"Loaded embeddings for speakers: {list(known_embeddings.keys())}")

    with temp_file_path(audio_file.getvalue(), ".wav") as audio_path:
        try:
            audio = AudioSegment.from_file(audio_path)
            df['speaker'] = None
            encoder = load_voice_encoder()
            progress_bar, status_text = st.progress(0), st.empty()

            for index, row in df.iterrows():
                segment = audio[row['start'] * 1000:row['end'] * 1000]

                with temp_file_path(segment.export(format="wav").read(), ".wav") as segment_path:
                    try:
                        wav = preprocess_wav(segment_path)
                        segment_embedding = encoder.embed_utterance(wav)
                        speaker = _identify_speaker(segment_embedding, known_embeddings, similarity_threshold)
                        df.at[index, 'speaker'] = speaker
                        status = f"Identified as {speaker}" if speaker else "Similarity below threshold"
                        status_text.text(f"Processed segment {index + 1}/{len(df)}: {status}")
                    except Exception as e:
                        st.error(f"Error processing segment {row['start']}-{row['end']}s: {e}")
                        df.at[index, 'speaker'] = "Error"

                progress_bar.progress((index + 1) / len(df))

            status_text.text("Speaker identification complete.")
            return df

        except Exception as e:
            st.error(f"Error loading or processing audio file: {e}")
            return df

def build_meeting_text_from_dataframe(df: pd.DataFrame) -> str:
    """Generate combined meeting text in (speaker) utterance format from transcription data."""
    if df is None or df.empty or 'text' not in df.columns:
        return ""

    if 'speaker' in df.columns:
        df_copy = df.copy()
        df_copy['speaker_filled'] = df_copy['speaker'].replace('', pd.NA)
        df_copy['speaker_filled'] = df_copy['speaker_filled'].ffill()
        df_copy['group_id'] = (df_copy['speaker_filled'] != df_copy['speaker_filled'].shift()).cumsum()
        df_merged = df_copy.groupby('group_id').agg(
            speaker=('speaker_filled', 'first'),
            text=('text', lambda values: ' '.join(map(str, values)))
        ).reset_index(drop=True)

        lines = []
        for _, row in df_merged.iterrows():
            speaker = row.get('speaker')
            text = row.get('text', '')
            speaker_str = speaker if isinstance(speaker, str) and speaker else '不明'
            lines.append(f"（{speaker_str}）{text}")
        return "\n".join(lines)

    # speaker列がない場合は不明話者として扱う
    lines = []
    for text in df['text'].astype(str):
        lines.append(f"（不明）{text}")
    return "\n".join(lines)

def format_time(seconds):
    """Formats seconds into HH:MM:SS."""
    td = timedelta(seconds=seconds)
    hours, remainder = divmod(td.seconds, 3600)
    minutes, seconds = divmod(remainder, 60)
    return f"{hours:02}:{minutes:02}:{seconds:02}"

def parse_time_to_seconds(time_str):
    """Converts HH:MM:SS or seconds string to total seconds."""
    if ':' in time_str:
        parts = list(map(int, time_str.split(':')))
        if len(parts) == 3:
            return parts[0] * 3600 + parts[1] * 60 + parts[2]
        elif len(parts) == 2:
            return parts[0] * 60 + parts[1]
        else:
            raise ValueError("Invalid time format. Use HH:MM:SS or MM:SS.")
    else:
        return int(time_str)

def split_text_by_lines(text, n_parts):
    """文字列を改行位置で適切にn分割"""
    if n_parts <= 0:
        raise ValueError("分割数は1以上である必要があります")
    if n_parts == 1:
        return [text]

    lines = text.split('\n')
    total_lines = len(lines)

    if total_lines <= n_parts:
        return lines + [""] * (n_parts - total_lines)

    lines_per_part = total_lines // n_parts
    remainder = total_lines % n_parts

    result, start = [], 0
    for i in range(n_parts):
        count = lines_per_part + (1 if i < remainder else 0)
        result.append('\n'.join(lines[start:start + count]))
        start += count

    return result

class RAGProofreadingSystem:
    """LangChainベースのRAG校正システム"""

    # 定数定義
    DEFAULT_CHUNK_SIZE = 500
    DEFAULT_CHUNK_OVERLAP = 100
    DEFAULT_TOP_K = 6
    DEFAULT_TEMPERATURE = 0.3
    EMBEDDING_MODEL = "text-embedding-3-large"

    def __init__(self, azure_endpoint, azure_api_key, api_version):
        self.azure_endpoint = azure_endpoint
        self.azure_api_key = azure_api_key
        self.api_version = api_version
        self.client = AzureOpenAI(
            azure_endpoint=azure_endpoint,
            api_key=azure_api_key,
            api_version=api_version
        )

        # LangChain embedding 初期化
        self.embeddings = AzureOpenAIEmbeddings(
            azure_deployment=self.EMBEDDING_MODEL,
            openai_api_version=api_version,
            azure_endpoint=azure_endpoint,
            api_key=azure_api_key
        )

        self.vectorstore = None
        self.documents = []
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=self.DEFAULT_CHUNK_SIZE,
            chunk_overlap=self.DEFAULT_CHUNK_OVERLAP,
            length_function=len,
        )

    def create_knowledge_base(self, documents_text_list, mode="add", documents_metadata=None):
        """ナレッジベースを構築（LangChain使用）

        Args:
            documents_text_list: ドキュメントテキストのリスト
            mode: "new" (新規作成) または "add" (追加構築)
            documents_metadata: 各ドキュメントのメタデータ辞書のリスト（オプション）
        """
        try:
            # テキストをLangChain Documentオブジェクトに変換
            langchain_docs = []
            for i, text in enumerate(documents_text_list):
                # メタデータを設定
                if documents_metadata and i < len(documents_metadata):
                    # カスタムメタデータを使用
                    metadata = documents_metadata[i].copy()
                else:
                    # デフォルトメタデータ
                    metadata = {"source": f"document_{i+1}"}

                doc = Document(
                    page_content=text,
                    metadata=metadata
                )
                langchain_docs.append(doc)

            # テキスト分割
            split_docs = self.text_splitter.split_documents(langchain_docs)

            # モードに応じて処理
            if mode == "new" or self.vectorstore is None:
                # 新規作成（または既存がない場合）
                self.vectorstore = FAISS.from_documents(split_docs, self.embeddings)
                self.documents = documents_text_list
            else:
                # 追加構築
                self.vectorstore.add_documents(split_docs)
                self.documents.extend(documents_text_list)

            return True
        except Exception as e:
            st.error(f"ナレッジベース構築エラー: {e}")
            return False

    def retrieve_relevant_context(self, query, search_type="similarity", top_k=None):
        """関連文脈を検索（LangChain使用）"""
        if not self.vectorstore:
            return ""

        if top_k is None:
            top_k = self.DEFAULT_TOP_K

        try:
            if search_type == "similarity":
                # 類似度検索
                docs = self.vectorstore.similarity_search(query, k=top_k)
            elif search_type == "mmr":
                # MMR検索（多様性を考慮）
                docs = self.vectorstore.max_marginal_relevance_search(query, k=top_k)
            else:
                docs = self.vectorstore.similarity_search(query, k=top_k)

            # 検索結果を結合
            context = "\n\n".join([doc.page_content for doc in docs])
            return context
        except Exception as e:
            st.error(f"文脈検索エラー: {e}")
            return ""

    def rag_enhanced_proofread(self, text, model="gpt-4o", search_type="similarity", top_k=None, prompt_preset="standard"):
        """
        RAG拡張校正（LangChain使用）

        Args:
            text: 校正対象のテキスト
            model: 使用するモデル
            search_type: 検索タイプ
            top_k: 検索結果の上位K件
            prompt_preset: プロンプトプリセットID（例: 'standard', 'detailed', 'simple', 'formal'）
        """
        try:
            if top_k is None:
                top_k = self.DEFAULT_TOP_K

            # 関連文脈を検索
            relevant_context = self.retrieve_relevant_context(text, search_type=search_type, top_k=top_k)

            # プロンプト構築（外部ファイルから読み込み）
            loader = get_default_loader()
            if relevant_context:
                system_prompt = loader.get_prompt_from_preset(
                    "rag_proofreading",
                    prompt_preset,
                    "system_prompt_with_context"
                ).format(context=relevant_context)
            else:
                system_prompt = loader.get_prompt_from_preset(
                    "rag_proofreading",
                    prompt_preset,
                    "system_prompt_without_context"
                )

            # Azure OpenAI APIで校正実行
            response = self.client.chat.completions.create(
                model=model,
                messages=[
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": text}
                ],
                temperature=self.DEFAULT_TEMPERATURE,
            )

            return response.choices[0].message.content

        except Exception as e:
            st.error(f"校正実行エラー: {e}")
            return None

    def save_knowledge_base(self, output_path):
        """ナレッジベースを保存（LangChain FAISS使用）"""
        try:
            if not self.vectorstore:
                return False, "保存するナレッジベースがありません"

            # 一時ディレクトリを作成
            with tempfile.TemporaryDirectory() as tmpdir:
                # FAISSインデックスを保存
                faiss_path = os.path.join(tmpdir, "faiss_index")
                self.vectorstore.save_local(faiss_path)

                # メタデータを保存
                metadata = {
                    "documents_count": len(self.documents),
                    "timestamp": datetime.now().isoformat()
                }
                metadata_path = os.path.join(tmpdir, "metadata.json")
                with open(metadata_path, "w", encoding="utf-8") as f:
                    json.dump(metadata, f, ensure_ascii=False, indent=2)

                # ZIPファイルに圧縮
                with zipfile.ZipFile(output_path, 'w', zipfile.ZIP_DEFLATED) as zipf:
                    # FAISSファイルを追加
                    for root, dirs, files in os.walk(faiss_path):
                        for file in files:
                            file_path = os.path.join(root, file)
                            arcname = os.path.join("faiss_index", file)
                            zipf.write(file_path, arcname)

                    # メタデータを追加
                    zipf.write(metadata_path, "metadata.json")

            return True, "保存成功"

        except Exception as e:
            return False, f"保存エラー: {e}"

    def load_knowledge_base(self, input_path):
        """ナレッジベースを読み込み（LangChain FAISS使用）"""
        try:
            # 一時ディレクトリを作成
            with tempfile.TemporaryDirectory() as tmpdir:
                # ZIPファイルを展開
                with zipfile.ZipFile(input_path, 'r') as zipf:
                    zipf.extractall(tmpdir)

                # FAISSインデックスを読み込み
                faiss_path = os.path.join(tmpdir, "faiss_index")
                self.vectorstore = FAISS.load_local(
                    faiss_path,
                    self.embeddings,
                    allow_dangerous_deserialization=True
                )

                # メタデータを読み込み
                metadata_path = os.path.join(tmpdir, "metadata.json")
                if os.path.exists(metadata_path):
                    with open(metadata_path, "r", encoding="utf-8") as f:
                        metadata = json.load(f)
                else:
                    metadata = {}

            return True, "読み込み成功", metadata

        except Exception as e:
            return False, f"読み込みエラー: {e}", {}

    def get_database_info(self):
        """データベース情報を取得"""
        has_data = self.vectorstore is not None
        return {
            "has_data": has_data,
            "documents_count": len(self.documents),
            "is_indexed": has_data,
            "total_chunks": self.vectorstore.index.ntotal if has_data else 0,
            "vector_files": 1 if has_data else 0,
            "output_files": 1 if has_data else 0,
            "search_types": ["similarity", "mmr"] if has_data else []
        }

    def get_chunks_detail(self):
        """チャンクの詳細情報を取得

        Returns:
            list: チャンク情報のリスト（各要素は辞書）
        """
        if not self.vectorstore:
            return []

        try:
            chunks_info = []
            # FAISSのdocstoreから全ドキュメントを取得
            docstore = self.vectorstore.docstore
            index_to_docstore_id = self.vectorstore.index_to_docstore_id

            for i in range(self.vectorstore.index.ntotal):
                doc_id = index_to_docstore_id[i]
                doc = docstore.search(doc_id)

                if doc:
                    chunks_info.append({
                        "chunk_id": i,
                        "doc_id": doc_id,
                        "content": doc.page_content,
                        "content_length": len(doc.page_content),
                        "source": doc.metadata.get("source", "不明"),
                        "metadata": doc.metadata
                    })

            return chunks_info
        except Exception as e:
            st.error(f"チャンク詳細取得エラー: {e}")
            return []

    def clear_knowledge_base(self):
        """ナレッジベースをクリア"""
        self.vectorstore = None
        self.documents = []

# ========================================
# 共通ヘルパー関数
# ========================================

def _init_rag_system():
    """RAGシステムの初期化（共通処理）"""
    if 'global_rag_system' not in st.session_state:
        st.session_state.global_rag_system = RAGProofreadingSystem(
            azure_endpoint=AZURE_OPENAI_ENDPOINT,
            azure_api_key=AZURE_OPENAI_API_KEY,
            api_version=API_VERSION
        )

    if 'global_db_info' not in st.session_state:
        st.session_state.global_db_info = st.session_state.global_rag_system.get_database_info()

    return st.session_state.global_rag_system

def _render_database_status(db_status, show_output_files=False):
    """データベース状態表示（共通処理）"""
    metrics = [
        ("チャンク数", db_status.get('total_chunks', 0)),
        ("インデックス状態", "✅ 構築済み" if db_status.get('is_indexed', False) else "❌ 未構築")
    ]

    if show_output_files:
        metrics.insert(1, ("出力ファイル数", db_status.get('output_files', 0)))

    for col, (label, value) in zip(st.columns(len(metrics)), metrics):
        with col:
            st.metric(label, value)

def _render_database_operations(rag_system, key_prefix="", show_save=True):
    """データベース操作UI（共通処理）"""
    st.subheader("🔧 データベース操作")

    has_data = st.session_state.global_db_info.get('has_data', False)
    cols = st.columns(3)

    # === 保存グループ ===
    with cols[0]:
        with st.container():
            st.markdown("##### 💾 保存")
            st.markdown("---")

            if show_save and has_data:
                st.caption("ナレッジベースをファイルとして保存")

                # デフォルトDBとして保存
                if st.button("📌 デフォルトDBに保存", key=f"{key_prefix}_save_default_db", use_container_width=True, type="primary"):
                    output_path = get_default_ragdb_path()
                    success, message = rag_system.save_knowledge_base(output_path)
                    if success:
                        st.success(f"✅ 保存完了")
                    else:
                        st.error(f"❌ {message}")

                st.markdown("")  # スペーサー

                # 別名で保存
                with st.expander("📥 別名でダウンロード", expanded=False):
                    custom_filename = st.text_input(
                        "ファイル名",
                        value="custom_knowledge_base",
                        key=f"{key_prefix}_custom_filename",
                        placeholder="ファイル名を入力"
                    )
                    if st.button("作成", key=f"{key_prefix}_save_custom_db", use_container_width=True):
                        # 拡張子を追加
                        filename_with_ext = custom_filename if custom_filename.endswith('.ragdb') else f"{custom_filename}.ragdb"

                        # 一時ファイルに保存
                        with tempfile.NamedTemporaryFile(delete=False, suffix='.ragdb') as tmp_file:
                            tmp_path = tmp_file.name

                        try:
                            success, message = rag_system.save_knowledge_base(tmp_path)
                            if success:
                                # ファイルを読み込んでダウンロードボタン表示
                                with open(tmp_path, 'rb') as f:
                                    db_bytes = f.read()

                                st.success(f"✅ 作成完了")
                                st.download_button(
                                    label="📥 ダウンロード",
                                    data=db_bytes,
                                    file_name=filename_with_ext,
                                    mime="application/octet-stream",
                                    key=f"{key_prefix}_download_custom_db",
                                    use_container_width=True
                                )
                            else:
                                st.error(f"❌ {message}")
                        finally:
                            # 一時ファイルを削除
                            if os.path.exists(tmp_path):
                                os.remove(tmp_path)
            else:
                st.caption("保存するデータがありません")
                st.markdown("<br>", unsafe_allow_html=True)
                st.button("📌 デフォルトDBに保存", key=f"{key_prefix}_save_default_db_disabled", use_container_width=True, disabled=True)

    # === 読み込みグループ ===
    with cols[1]:
        with st.container():
            st.markdown("##### 📂 読み込み")
            st.markdown("---")
            st.caption("保存済みのナレッジベースを読み込み")

            uploaded_db = st.file_uploader(
                "RAGDBファイルを選択",
                type=['ragdb'],
                key=f"{key_prefix}_load_db",
                label_visibility="collapsed"
            )
            if st.button("📂 読み込み実行", key=f"{key_prefix}_load_btn", use_container_width=True, type="primary", disabled=uploaded_db is None):
                with temp_file_path(uploaded_db.getvalue(), '.ragdb') as tmp_path:
                    success, message, metadata = rag_system.load_knowledge_base(tmp_path)
                    if success:
                        st.success("✅ 読み込み完了")
                        st.session_state.global_db_info = rag_system.get_database_info()
                        st.rerun()
                    else:
                        st.error(message)

    # === クリアグループ ===
    with cols[2]:
        with st.container():
            st.markdown("##### 🗑️ クリア")
            st.markdown("---")

            if has_data:
                st.caption("現在のナレッジベースを削除")
                st.markdown("<br>", unsafe_allow_html=True)
                if st.button("🗑️ データベースをクリア", key=f"{key_prefix}_clear_btn", use_container_width=True, type="primary"):
                    rag_system.clear_knowledge_base()
                    st.session_state.global_db_info = rag_system.get_database_info()
                    st.success("✅ クリア完了")
                    st.rerun()
            else:
                st.caption("クリアするデータがありません")
                st.markdown("<br>", unsafe_allow_html=True)
                st.button("🗑️ データベースをクリア", key=f"{key_prefix}_clear_btn_disabled", use_container_width=True, disabled=True)

# ========================================
# ページ関数
# ========================================

def knowledge_base_management():
    st.title("📚 ナレッジベース管理")
    st.write("議事録校正用のナレッジベースを構築・管理します。")

    st.sidebar.markdown("""
    ### 📚 ナレッジベース管理

    **概要**
    LangChainとFAISSを使用した高度なRAGシステムです。

    **主な機能**
    - ドキュメント追加: PDF、Word、PowerPoint対応
    - ベクトル検索: 高速セマンティック検索
    - DB保存: .ragdb形式で保存・読み込み
    - 検索タイプ: similarity、mmr対応

    💡 ドキュメント追加で既存データと自動マージ
    """)

    # RAGシステムの初期化
    rag_system = _init_rag_system()

    # デフォルトのRAGDBファイルを自動読み込み（初回のみ）
    if 'kb_default_db_loaded' not in st.session_state:
        default_ragdb_path = get_default_ragdb_path()
        if os.path.exists(default_ragdb_path):
            try:
                success, message, metadata = rag_system.load_knowledge_base(default_ragdb_path)
                if success:
                    st.success(f"✅ デフォルトナレッジベース '{default_ragdb_path}' を自動読み込みしました")
                    st.session_state.global_db_info = rag_system.get_database_info()
                else:
                    st.warning(f"⚠️ デフォルトナレッジベースの読み込みに失敗しました: {message}")
            except Exception as e:
                st.warning(f"⚠️ デフォルトナレッジベースの読み込みに失敗しました: {e}")
        st.session_state.kb_default_db_loaded = True

    # ナレッジベース状態表示
    st.subheader("📊 現在の状態")
    _render_database_status(st.session_state.global_db_info)

    # データベース操作
    _render_database_operations(rag_system, key_prefix="kb", show_save=True)

    # ナレッジベース詳細表示
    if st.session_state.global_db_info.get("has_data", False):
        with st.expander("🔍 ナレッジベース詳細", expanded=False):
            # チャンク詳細を取得
            chunks_detail = rag_system.get_chunks_detail()

            if chunks_detail:
                # フィルタリングオプション
                filter_col1, filter_col2 = st.columns([2, 1])
                with filter_col1:
                    search_text = st.text_input(
                        "🔎 チャンク内容で検索",
                        key="chunk_search",
                        placeholder="検索キーワードを入力..."
                    )
                with filter_col2:
                    unique_sources = sorted(set(chunk["source"] for chunk in chunks_detail))
                    selected_source = st.selectbox(
                        "📁 ソースでフィルタ",
                        options=["すべて"] + unique_sources,
                        key="chunk_source_filter"
                    )

                # フィルタリング適用
                filtered_chunks = chunks_detail
                if search_text:
                    filtered_chunks = [
                        chunk for chunk in filtered_chunks
                        if search_text.lower() in chunk["content"].lower()
                    ]
                if selected_source != "すべて":
                    filtered_chunks = [
                        chunk for chunk in filtered_chunks
                        if chunk["source"] == selected_source
                    ]

                st.caption(f"表示中: {len(filtered_chunks)} / {len(chunks_detail)} チャンク")

                # データフレーム形式で表示
                if filtered_chunks:
                    display_data = []
                    for chunk in filtered_chunks:
                        # メタデータから情報を取得
                        metadata = chunk["metadata"]
                        upload_datetime = metadata.get("upload_datetime", "-")
                        file_type = metadata.get("file_type", "-")

                        # コンテンツを適度に表示
                        content_preview = chunk["content"][:200] + "..." if len(chunk["content"]) > 200 else chunk["content"]

                        display_data.append({
                            "ソース": chunk["source"],
                            "投入日時": upload_datetime,
                            "形式": file_type,
                            "文字数": chunk["content_length"],
                            "内容": content_preview
                        })

                    # テーブル表示
                    st.dataframe(
                        display_data,
                        use_container_width=True,
                        hide_index=True,
                        height=400
                    )
                else:
                    st.warning("フィルタ条件に一致するチャンクがありません。")
            else:
                st.warning("チャンク情報を取得できませんでした。")

    # ナレッジベース構築
    st.subheader("🏗️ ナレッジベース構築")

    # 構築モード選択
    build_mode = st.radio(
        "構築モードを選択してください",
        options=["追加構築", "新規構築"],
        index=0,  # デフォルト: 追加構築
        horizontal=True,
        help="追加構築: 既存のナレッジベースに追加 | 新規構築: 既存を削除して新規作成"
    )

    # ドキュメントアップロード
    doc_files = st.file_uploader(
        "ドキュメントファイルを選択してください（複数選択可）",
        type=["pdf", "txt", "docx", "pptx"],
        accept_multiple_files=True,
        key="kb_doc_files",
        help="選択したファイルからナレッジベースを構築します"
    )

    if doc_files and st.button("🚀 ナレッジベースを構築", key="kb_build_btn", type="primary"):
        try:
            from datetime import datetime

            documents = []
            documents_metadata = []
            processed_files = []
            upload_datetime = datetime.now().strftime("%Y-%m-%d %H:%M:%S")

            with st.spinner("ドキュメントを処理中..."):
                for doc_file in doc_files:
                    try:
                        file_extension = doc_file.name.split(".")[-1].lower()
                        doc_content = extract_text_from_file(BytesIO(doc_file.read()), file_extension)

                        if doc_content.strip():
                            documents.append(doc_content)
                            # メタデータを追加
                            documents_metadata.append({
                                'source': doc_file.name,  # ファイル名をソースとして保存
                                'upload_datetime': upload_datetime,  # 投入日時
                                'file_type': file_extension.upper(),
                                'content_length': len(doc_content)
                            })
                            processed_files.append({
                                'name': doc_file.name,
                                'type': file_extension.upper(),
                                'size': len(doc_content),
                                'status': 'success'
                            })
                            st.success(f"✅ {doc_file.name} ({file_extension.upper()}) - {len(doc_content):,}文字")
                        else:
                            processed_files.append({
                                'name': doc_file.name,
                                'type': file_extension.upper(),
                                'size': 0,
                                'status': 'empty'
                            })
                            st.warning(f"⚠️ {doc_file.name} - テキストが抽出できませんでした")
                    except Exception as e:
                        processed_files.append({
                            'name': doc_file.name,
                            'type': 'ERROR',
                            'size': 0,
                            'status': 'error',
                            'error': str(e)
                        })
                        st.error(f"❌ {doc_file.name} の処理に失敗: {e}")

            if documents:
                with st.spinner("ナレッジベースを構築中..."):
                    # 構築モードを変換
                    mode = "add" if build_mode == "追加構築" else "new"

                    # ナレッジベース構築
                    if mode == "new":
                        st.info("既存データを削除して新規構築します")
                    else:
                        st.info("既存データに追加構築します")

                    success = st.session_state.global_rag_system.create_knowledge_base(
                        documents,
                        mode=mode,
                        documents_metadata=documents_metadata
                    )

                    if success:
                        st.session_state.global_db_info = st.session_state.global_rag_system.get_database_info()
                    else:
                        raise Exception("ベクターストア構築に失敗しました")

                    # 構築結果表示
                    st.success(f"✅ ナレッジベースの構築が完了しました！")

                    # 詳細統計
                    with st.expander("📈 構築結果詳細"):
                        new_total_chars = sum(len(doc) for doc in documents)
                        st.metric("新規追加文字数", f"{new_total_chars:,}")
                        st.metric("ドキュメント数", st.session_state.global_db_info['documents_count'])
                        st.metric("ベクターファイル数", st.session_state.global_db_info.get('vector_files', 0))
                        st.metric("検索タイプ数", len(st.session_state.global_db_info.get('search_types', [])))

                        # 利用可能な検索タイプを表示
                        search_types = st.session_state.global_db_info.get('search_types', [])
                        if search_types:
                            st.write(f"**利用可能な検索タイプ**: {', '.join(search_types)}")

                        # ファイル処理結果テーブル
                        st.write("**ファイル処理結果**")
                        for file_info in processed_files:
                            status_icon = {
                                'success': '✅',
                                'empty': '⚠️',
                                'error': '❌'
                            }.get(file_info['status'], '❓')

                            st.write(f"{status_icon} **{file_info['name']}** ({file_info['type']}) - {file_info['size']:,}文字")

                    st.rerun()
            else:
                st.error("処理可能なドキュメントがありませんでした。")

        except Exception as e:
            st.error(f"ナレッジベース構築中にエラーが発生しました: {e}")

def proofread_meeting_minutes():
    st.title("📝 議事録校正")
    st.write("ナレッジベースを活用した高精度な議事録校正を行います。")

    st.sidebar.markdown("""
    ### 📝 議事録校正システム

    **概要**
    ナレッジベースを活用した高精度な議事録校正を提供します。

    **システムの特徴**
    - 類似度検索: セマンティック検索
    - チャンク分割: 最適サイズで処理
    - 高速検索: 効率的ベクトル検索
    - シンプル設計: 保守しやすい構成

    **使用手順**
    1. ナレッジベースを確認
    2. 議事録テキストを入力
    3. 検索タイプを選択
    4. 校正を実行

    💡 未構築の場合は「ナレッジベース管理」から
    """)

    # RAGシステムの初期化
    rag_system = _init_rag_system()

    # デフォルトのRAGDBファイルを自動読み込み（初回のみ）
    if 'default_db_loaded' not in st.session_state:
        default_ragdb_path = get_default_ragdb_path()
        if os.path.exists(default_ragdb_path):
            try:
                success, message, metadata = rag_system.load_knowledge_base(default_ragdb_path)
                if success:
                    st.success(f"✅ デフォルトナレッジベース '{default_ragdb_path}' を自動読み込みしました")
                    st.session_state.global_db_info = rag_system.get_database_info()
                else:
                    st.warning(f"⚠️ デフォルトナレッジベースの読み込みに失敗しました: {message}")
            except Exception as e:
                st.warning(f"⚠️ デフォルトナレッジベースの読み込みに失敗しました: {e}")
        st.session_state.default_db_loaded = True

    # ナレッジベース状態表示
    st.subheader("📚 ナレッジベース状態")
    db_status = st.session_state.global_db_info
    _render_database_status(db_status, show_output_files=True)

    has_data = db_status.get('has_data', False)
    if not has_data:
        st.warning("⚠️ ナレッジベースが構築されていません。「ナレッジベース管理」で事前に構築してください。")
        st.info("📖 ナレッジベースなしでも基本的な校正は実行できますが、高度な文脈参照機能は利用できません。")

    # データベース操作
    with st.expander("🔧 データベース操作"):
        _render_database_operations(rag_system, key_prefix="proofreading", show_save=False)

    # 検索結果数の設定
    top_k = st.selectbox(
        "検索する関連文脈の数",
        options=[1, 2, 3, 4, 5, 6, 7, 8, 9, 10],
        index=5,  # デフォルト: 6
        help="より多くの文脈を検索すると精度が向上しますが、処理時間が増加します"
    )

    # テキスト入力方法の選択
    input_method = st.radio(
        "議事録テキストの入力方法を選択してください",
        ["テキストファイル(.txt/.docx)をアップロード", "テキストボックスに直接入力"],
        key="rag_input_method_selector"
    )

    transcript_text = ""

    if input_method == "テキストファイル(.txt/.docx)をアップロード":
        uploaded_text_file = st.file_uploader(
            "議事録テキストファイルを選択してください",
            type=["txt", "docx"],
            key="rag_upload_text_file"
        )

        if uploaded_text_file is not None:
            try:
                file_extension = uploaded_text_file.name.lower().split('.')[-1]

                if file_extension == 'txt':
                    transcript_text = uploaded_text_file.read().decode('utf-8')
                elif file_extension == 'docx':
                    # Word文書からテキストを抽出
                    transcript_text = get_text_from_docx(BytesIO(uploaded_text_file.read()))

                st.success(f"{file_extension.upper()}ファイルが正常に読み込まれました。")
                st.text_area("読み込まれたテキスト（プレビュー）",
                           transcript_text[:500] + "..." if len(transcript_text) > 500 else transcript_text,
                           height=150, key="rag_text_preview")
            except Exception as e:
                st.error(f"テキストファイルの読み込み中にエラーが発生しました: {e}")

    else:  # テキストボックスに直接入力
        transcript_text = st.text_area(
            "議事録テキストを入力してください",
            height=300,
            key="rag_direct_text_input",
            placeholder="ここに議事録テキストを貼り付けてください..."
        )

    # 校正設定
    st.subheader("🎛️ LangChain校正設定")

    config_col1, config_col2 = st.columns(2)

    with config_col1:
        # 検索タイプ選択
        search_types = st.session_state.global_db_info.get('search_types', [])
        if search_types:
            search_type = st.selectbox(
                "検索タイプ",
                search_types,
                key="search_type_selection",
                help="similarity: ベクトル類似度検索、mmr: 多様性を考慮した検索"
            )

            if search_type == 'similarity':
                st.info("🎯 **類似度検索**: ベクトル類似度による最も関連性の高い文書検索")
            elif search_type == 'mmr':
                st.info("🔄 **MMR検索**: 関連性と多様性を両立した文書検索")
        else:
            search_type = "similarity"
            st.warning("ベクターストアが初期化されていません。デフォルトで類似度検索を使用します。")

    with config_col2:
        # 校正モデル選択
        model_choice = st.selectbox(
            "LLMモデル",
            ["gpt-4o", "gpt-4o-mini"],
            key="proofreading_model",
            help="校正に使用するAIモデルを選択"
        )

    # プロンプトプリセット選択
    st.markdown("### 🎨 プロンプトプリセット")

    # プロンプトプリセット一覧を取得
    loader = get_default_loader()
    presets = loader.get_preset_list("rag_proofreading")

    if presets:
        # プリセットの選択肢を作成
        preset_options = {f"{p['name']} - {p['description']}": p['id'] for p in presets}

        # セレクトボックスで表示
        selected_preset_label = st.selectbox(
            "校正プロンプトを選択",
            options=list(preset_options.keys()),
            index=0,  # デフォルトは最初のプリセット（standard）
            key="prompt_preset_selection",
            help="校正の方針を選択してください"
        )

        # 選択されたプリセットのIDを取得
        selected_preset = preset_options[selected_preset_label]

        # 選択されたプリセットの詳細を表示
        selected_preset_info = next((p for p in presets if p['id'] == selected_preset), None)
        if selected_preset_info:
            st.info(f"📝 **{selected_preset_info['name']}**: {selected_preset_info['description']}")
    else:
        selected_preset = "standard"
        st.warning("プロンプトプリセットが見つかりません。デフォルトプリセットを使用します。")
    ##2025.9.8 修正：分割処理の追加
    # 分割処理設定の追加
    st.subheader("📏 テキスト分割設定")

    split_col1, split_col2 = st.columns(2)

    with split_col1:
        n_length = st.number_input(
            "分割閾値 (文字数)",
            min_value=100,
            max_value=10000,
            value=1000,
            step=100,
            key="text_split_threshold",
            help="この文字数を超える場合、テキストを分割して処理します"
        )

    with split_col2:
        if transcript_text.strip():
            current_length = len(transcript_text)
            estimated_parts = max(1, (current_length // n_length) + (1 if current_length % n_length > 0 else 0))
            st.metric("現在のテキスト長", f"{current_length:,}文字")
            st.metric("推定分割数", f"{estimated_parts}部分")

    # LangChain校正実行
    if transcript_text.strip() and st.button("📝 RAG校正を実行", key="execute_rag_proofreading", type="primary"):
        if not transcript_text.strip():
            st.error("議事録テキストが入力されていません。")
            return

        try:
            ##2025.9.8 修正：分割処理の追加
            # テキストの分割判定と処理
            current_length = len(transcript_text)

            if current_length > n_length:
                # 分割が必要な場合
                n_parts = max(1, (current_length // n_length) + (1 if current_length % n_length > 0 else 0))
                st.info(f"📏 テキストが長いため、{n_parts}部分に分割して処理します（閾値: {n_length:,}文字、実際: {current_length:,}文字）")

                # テキストを分割
                text_parts = split_text_by_lines(transcript_text, n_parts)

                # 各パートを順次処理
                proofread_parts = []
                progress_bar = st.progress(0)
                status_text = st.empty()

                with st.spinner(f"LangChainベース校正を実行中（分割処理: {n_parts}部分）..."):
                    current_db_status = st.session_state.global_rag_system.get_database_info()
                    if current_db_status['has_data']:
                        st.info(f"🔍 RAGナレッジベース使用中（文書数: {current_db_status['documents_count']}, 検索タイプ: {search_type}）")
                    else:
                        st.info("📖 RAGナレッジベースなしで基本校正を実行します")

                    for i, part in enumerate(text_parts, 1):
                        print("文字列:", part)
                        if part.strip():  # 空のパートはスキップ
                            status_text.text(f"Part {i}/{n_parts} を処理中... ({len(part)}文字)")

                            # 各パートを校正
                            part_result = st.session_state.global_rag_system.rag_enhanced_proofread(
                                part,
                                model=model_choice,
                                search_type=search_type if current_db_status['has_data'] else "similarity",
                                top_k=top_k,
                                prompt_preset=selected_preset
                            )

                            if part_result:
                                proofread_parts.append(part_result)

                        progress_bar.progress(i / n_parts)

                    # 結果を連結
                    proofread_result = "\n\n".join(proofread_parts) if proofread_parts else ""
                    status_text.empty()
                    progress_bar.empty()

            else:
                # 分割不要な場合（従来の処理）
                with st.spinner(f"LangChainベース校正を実行中 ({search_type}検索使用)..."):
                    # 現在のLangChain RAGナレッジベース状況を表示
                    current_db_status = st.session_state.global_rag_system.get_database_info()
                    if current_db_status['has_data']:
                        st.info(f"🔍 RAGナレッジベース使用中（文書数: {current_db_status['documents_count']}, 検索タイプ: {search_type}）")
                    else:
                        st.info("📖 RAGナレッジベースなしで基本校正を実行します")

                    # LangChain校正の実行
                    proofread_result = st.session_state.global_rag_system.rag_enhanced_proofread(
                        transcript_text,
                        model=model_choice,
                        search_type=search_type if current_db_status['has_data'] else "similarity",
                        top_k=top_k,
                        prompt_preset=selected_preset
                    )

                if proofread_result:
                    st.success(f"✅ RAG校正が完了しました！ ({search_type}検索使用)")

                    # 結果表示
                    st.subheader("📄 校正結果")
                    st.text_area("校正された議事録", proofread_result, height=400, key="final_rag_result")

                    # 使用されたLangChain RAG文脈の表示
                    if current_db_status['has_data']:
                        with st.expander(f"🔍 RAG検索結果 ({search_type}検索)"):
                            relevant_context = st.session_state.global_rag_system.retrieve_relevant_context(
                                transcript_text,
                                search_type=search_type,
                                top_k=top_k
                            )
                            if relevant_context:
                                st.text_area("RAG検索結果", relevant_context, height=200, key="final_context_display")

                                # 検索タイプの詳細説明
                                if search_type == 'similarity':
                                    st.info("🎯 類似度検索: ベクトル類似度による関連文書から情報を抽出")
                                elif search_type == 'mmr':
                                    st.info("🔄 MMR検索: 関連性と多様性を考慮した文書から文脈を取得")
                            else:
                                st.write("関連文脈が見つかりませんでした。")

                    # ダウンロードと統計
                    col_dl, col_stats = st.columns([1, 1])

                    with col_dl:
                        # Word文書として保存
                        doc = DocxDocument()
                        doc.add_heading('RAG校正済み議事録', 0)

                        # テキストを段落に分割して追加
                        for line in proofread_result.split('\n'):
                            if line.strip():
                                doc.add_paragraph(line)
                            else:
                                doc.add_paragraph('')  # 空行を保持

                        docx_buffer = BytesIO()
                        doc.save(docx_buffer)
                        docx_buffer.seek(0)

                        st.download_button(
                            label="📥 RAG校正済み議事録をダウンロード",
                            data=docx_buffer.getvalue(),
                            file_name=f"rag_proofread_{search_type}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx",
                            mime='application/vnd.openxmlformats-officedocument.wordprocessingml.document',
                            key="download_final_rag_result"
                        )

                    with col_stats:
                        # 統計情報の表示
                        with st.expander("📊 校正統計"):
                            original_length = len(transcript_text)
                            proofread_length = len(proofread_result)
                            change_ratio = ((proofread_length - original_length) / original_length * 100) if original_length > 0 else 0

                            st.metric("元の文字数", f"{original_length:,}")
                            st.metric("校正後文字数", f"{proofread_length:,}")
                            st.metric("変化率", f"{change_ratio:.1f}%")

                            if current_db_status['has_data']:
                                st.metric("参照チャンク数", current_db_status['total_chunks'])
                                st.metric("検索範囲", f"上位{top_k}件")

                else:
                    st.error("❌ 校正処理に失敗しました。Azure OpenAIの設定を確認してください。")

        except Exception as e:
            st.error(f"❌ RAG校正処理中にエラーが発生しました: {e}")
            st.info("💡 Azure OpenAIの設定とネットワーク接続を確認してください。")

    elif not transcript_text.strip():
        st.info("💭 議事録テキストを入力してからRAG校正ボタンを押してください。")

def batch_processing_pipeline():
    """一括処理パイプライン: 1本のメディアから複数セグメントを段階的に処理"""
    st.title("🚀 一括処理パイプライン")
    st.write("アップロードした動画/音声ファイルを区間ごとに切り出し、文字起こしから校正までをまとめて実行します。")
    st.write("文字起こし → 個別話者識別 → RAG校正 → 出力のパイプラインを順番に進めるだけで完了します。")

    st.sidebar.markdown("""
    ### 🚀 一括処理パイプライン

    **概要**
    動画・音声ファイルから必要な区間を切り出し、一括で文字起こし・話者識別・RAG校正を実行します。

    **処理フロー**
    1. メディアファイルのアップロードと区間設定
    2. モデル選択と文字起こし実行
    3. 各ファイルごとの話者識別（個別設定可能）
    4. RAG校正実行
    5. 処理結果の確認とダウンロード

    **対応形式:** MP4, MOV, AVI, MKV, WebM, MP3, WAV, M4A等
    """)

    # セッション状態の初期化
    _init_session_state({
        'batch_uploaded_video': None,
        'batch_extracted_files': [],
        'batch_processing_results': {},
        'batch_processing_status': {},
        'batch_current_step': 1,
        'batch_rag_system': None,
        'batch_transcribe_model': 'whisper',
        'batch_reference_file': None,
        'batch_current_speaker_file_index': 0,
        'batch_db_info': None,
        'batch_meeting_type': None,
        'batch_default_embeddings': [],
        'batch_file_embeddings': {},  # 各ファイルごとの埋め込み設定
        'batch_embedding_states': {},  # 話者埋め込み作成UI用の状態
        'batch_temp_dir': None,  # 一時ファイルディレクトリ
        'batch_cut_settings_df': None  # 切り出し設定DataFrame
    })

    # Step 1: 動画アップロードと音声切り出し設定
    st.subheader("Step 1: メディアアップロードと音声切り出し")
    st.write("動画・音声ファイルをアップロードし、処理したい区間を設定してください。複数のセグメントを指定できます。")
    st.caption("※ 切り出しは任意です。ファイル全体を使う場合はそのまま次のステップへ進めます。")

    uploaded_video = st.file_uploader(
        "メディアファイルを選択してください",
        type=["wav", "mp3", "mp4", "mov", "avi", "mkv", "webm", "m4a"],
        key="batch_video_upload"
    )

    if uploaded_video is not None:
        st.success(f"✅ メディアファイル '{uploaded_video.name}' を読み込みました")

        # メディアプレビュー
        video_ext = os.path.splitext(uploaded_video.name)[1].lower()
        if video_ext in ['.mp4', '.mov', '.webm', '.avi', '.mkv']:
            st.video(uploaded_video)
        elif video_ext in ['.mp3', '.wav', '.m4a']:
            st.audio(uploaded_video)

        st.subheader("切り出し区間の設定")

        # デフォルトのdata_editorデータ
        if st.session_state.batch_cut_settings_df is None:
            default_data = pd.DataFrame([
                {"開始時間": "00:00:00", "終了時間": "00:00:30", "出力ファイル名": f"{os.path.splitext(uploaded_video.name)[0]}_"}
            ])
        else:
            default_data = st.session_state.batch_cut_settings_df

        edited_df = st.data_editor(
            default_data,
            num_rows="dynamic",
            use_container_width=True,
            column_config={
                "開始時間": st.column_config.TextColumn(
                    "開始時間 (HH:MM:SS or seconds)",
                    help="切り出し開始時間 (例: 00:00:10 または 10)",
                    default="00:00:00"
                ),
                "終了時間": st.column_config.TextColumn(
                    "終了時間 (HH:MM:SS or seconds)",
                    help="切り出し終了時間 (例: 00:00:30 または 30)",
                    default="00:00:30"
                ),
                "出力ファイル名": st.column_config.TextColumn(
                    "出力ファイル名 (.mp3)",
                    help="この区間のMP3出力ファイル名を入力してください (例: meeting1.mp3)。空欄の場合、自動で連番が振られます。",
                    default=f"{os.path.splitext(uploaded_video.name)[0]}_"
                )
            },
            key="batch_cut_settings_editor"
        )

        # 設定を保存
        st.session_state.batch_cut_settings_df = edited_df

        # 音声切り出し実行
        if st.button("🎵 音声を切り出してStep 2へ進む", key="batch_execute_cut", type="primary"):
            if edited_df.empty:
                st.warning("切り出し区間が設定されていません。")
            else:
                with st.spinner("音声の切り出しとMP3への変換中..."):
                    try:
                        # 一時ディレクトリを作成
                        if st.session_state.batch_temp_dir is None:
                            st.session_state.batch_temp_dir = tempfile.mkdtemp()

                        temp_dir = st.session_state.batch_temp_dir
                        temp_video_path = None
                        extracted_files = []

                        # 動画ファイルを一時保存
                        uploaded_video.seek(0)
                        temp_video_path = os.path.join(temp_dir, f"source_video{os.path.splitext(uploaded_video.name)[1]}")
                        with open(temp_video_path, 'wb') as f:
                            f.write(uploaded_video.read())

                        # 各区間を切り出し
                        for index, row in edited_df.iterrows():
                            start_time_str = str(row["開始時間"])
                            end_time_str = str(row["終了時間"])
                            output_filename_raw = str(row["出力ファイル名"]).strip()

                            try:
                                start_seconds = parse_time_to_seconds(start_time_str)
                                end_seconds = parse_time_to_seconds(end_time_str)

                                if start_seconds >= end_seconds:
                                    st.error(f"区間 {index+1}: 開始時間 ({start_time_str}) は終了時間 ({end_time_str}) より前に設定してください。この区間はスキップされます。")
                                    continue

                                # 出力ファイル名を決定
                                base_name_from_video = os.path.splitext(uploaded_video.name)[0]
                                if not output_filename_raw or output_filename_raw.upper() == "AUTO_GENERATE":
                                    output_filename = f"{base_name_from_video}_segment_{index+1}.mp3"
                                else:
                                    output_filename = output_filename_raw

                                # .mp3拡張子を確保
                                if not output_filename.lower().endswith(".mp3"):
                                    output_filename += ".mp3"

                                output_audio_path = os.path.join(temp_dir, output_filename)

                                # FFmpegで切り出し
                                command = [
                                    "ffmpeg",
                                    "-i", temp_video_path,
                                    "-ss", format_time(start_seconds),
                                    "-to", format_time(end_seconds),
                                    "-vn",  # No video
                                    "-ab", "192k",  # Audio bitrate
                                    "-map_metadata", "-1",  # Remove metadata
                                    "-y",  # Overwrite output files without asking
                                    output_audio_path
                                ]

                                process = subprocess.run(command, capture_output=True, text=True, encoding="utf-8", check=True)

                                # 生成されたファイルをメモリに読み込み
                                with open(output_audio_path, 'rb') as f:
                                    file_data = f.read()
                                    file_io = BytesIO(file_data)
                                    file_io.name = output_filename
                                    extracted_files.append({
                                        'name': output_filename,
                                        'data': file_io,
                                        'size': len(file_data)
                                    })

                                st.success(f"✅ 区間 {index+1}: {output_filename} の切り出しが完了しました")

                            except subprocess.CalledProcessError as e:
                                st.error(f"❌ 区間 {index+1}: FFmpegの実行中にエラーが発生しました: {e}")
                                st.code(e.stderr)
                            except ValueError as e:
                                st.error(f"❌ 区間 {index+1}: 時間形式エラー: {e}")
                            except Exception as e:
                                st.error(f"❌ 区間 {index+1}: 処理中にエラーが発生しました: {e}")

                        if extracted_files:
                            # 抽出されたファイルをsession stateに保存
                            st.session_state.batch_extracted_files = extracted_files
                            st.session_state.batch_current_step = 2

                            # 各ファイルの処理状態を初期化
                            for file_info in extracted_files:
                                st.session_state.batch_processing_status[file_info['name']] = {
                                    'transcription': 'pending',
                                    'speaker_id': 'pending',
                                    'rag_proofread': 'pending'
                                }
                                st.session_state.batch_processing_results[file_info['name']] = {}

                            st.success(f"✅ {len(extracted_files)} 個の音声ファイルの切り出しが完了しました")
                            st.rerun()
                        else:
                            st.warning("⚠️ 音声ファイルが1つも生成されませんでした。設定を確認してください。")

                    except Exception as e:
                        st.error(f"❌ 音声切り出し処理中にエラーが発生しました: {e}")
                        import traceback
                        st.error(traceback.format_exc())

        if st.button("⏭ 切り出しをスキップしてStep 2へ進む", key="batch_skip_cut"):
            uploaded_video.seek(0)
            file_bytes = uploaded_video.read()

            if not file_bytes:
                st.error("❌ ファイルの読み込みに失敗しました。もう一度アップロードしてください。")
            else:
                file_buffer = BytesIO(file_bytes)
                file_buffer.name = uploaded_video.name

                st.session_state.batch_extracted_files = [{
                    'name': uploaded_video.name,
                    'data': file_buffer,
                    'size': len(file_bytes)
                }]

                st.session_state.batch_processing_status = {
                    uploaded_video.name: {
                        'transcription': 'pending',
                        'speaker_id': 'pending',
                        'rag_proofread': 'pending'
                    }
                }
                st.session_state.batch_processing_results = {
                    uploaded_video.name: {}
                }
                st.session_state.batch_current_step = 2
                st.success("✅ 切り出しをスキップしました。元のメディアファイルをStep 2でそのまま文字起こしします。")
                st.rerun()

    # Step 2: モデル選択と文字起こし
    if st.session_state.batch_current_step >= 2 and len(st.session_state.batch_extracted_files) > 0:
        st.subheader("Step 2: モデル選択と文字起こし")

        # ファイル一覧を表示
        st.write(f"**検出されたファイル数:** {len(st.session_state.batch_extracted_files)}")

        file_df = pd.DataFrame([
            {
                'ファイル名': f['name'],
                'サイズ (KB)': f'{f["size"] / 1024:.1f}'
            }
            for f in st.session_state.batch_extracted_files
        ])
        st.dataframe(file_df, use_container_width=True, hide_index=True)

        st.divider()

        reference_file = None
        meeting_types = load_meeting_type_config()
        col_settings, col_meeting = st.columns(2, gap="large")

        with col_settings:
            st.write("**文字起こし設定**")

            transcribe_model = st.selectbox(
                "文字起こしモデル",
                options=["whisper", "gpt-4o-transcribe-diarize"],
                index=0,
                key="batch_transcribe_model_select",
                help="whisper: 参考資料対応 | gpt-4o-transcribe-diarize: 自動話者識別付き"
            )

            if transcribe_model == "whisper":
                reference_file = st.file_uploader(
                    "参考資料（オプション）",
                    type=["pdf", "docx", "pptx", "txt", "msg"],
                    key="batch_reference_file_upload",
                    help="全ファイルの文字起こし精度向上に使用"
                )

        with col_meeting:
            st.write("**会議タイプ設定（オプション）**")
            st.write("会議タイプを選択すると、事前に設定された話者埋め込みファイルを自動的に読み込みます。")

            if meeting_types:
                meeting_type_options = {mt['id']: f"{mt['name']} - {mt['description']}" for mt in meeting_types}

                selected_meeting_type_id = st.selectbox(
                    "会議タイプ",
                    options=list(meeting_type_options.keys()),
                    format_func=lambda x: meeting_type_options[x],
                    key="batch_meeting_type_select"
                )

                selected_meeting_type = next((mt for mt in meeting_types if mt['id'] == selected_meeting_type_id), None)

                if selected_meeting_type and selected_meeting_type['embeddings_folder']:
                    if st.session_state.batch_meeting_type != selected_meeting_type_id:
                        st.session_state.batch_meeting_type = selected_meeting_type_id
                        st.session_state.batch_default_embeddings = load_embeddings_from_folder(
                            selected_meeting_type['embeddings_folder']
                        )

                        if st.session_state.batch_default_embeddings:
                            st.success(f"✅ {len(st.session_state.batch_default_embeddings)}個の話者埋め込みファイルを読み込みました")
                            embedding_names = [emb['name'] for emb in st.session_state.batch_default_embeddings]
                            st.info(f"📁 読み込まれたファイル: {', '.join(embedding_names)}")
                elif selected_meeting_type_id == 'custom':
                    st.session_state.batch_meeting_type = 'custom'
                    st.session_state.batch_default_embeddings = []
                    st.info("💡 カスタムモード: Step 3で個別に話者埋め込みファイルを指定してください")
            else:
                st.warning("⚠️ 会議タイプのプリセットが定義されていません。手動で話者埋め込みを指定してください。")

        st.divider()

        # 文字起こし実行
        if st.button("📝 文字起こしを開始", key="batch_start_transcription", type="primary"):
            # 進捗表示用のプレースホルダー
            progress_bar = st.progress(0)
            status_text = st.empty()

            total_files = len(st.session_state.batch_extracted_files)

            for idx, file_info in enumerate(st.session_state.batch_extracted_files):
                file_name = file_info['name']
                file_data = file_info['data']

                status_text.write(f"**文字起こし中: {file_name}** ({idx + 1}/{total_files})")

                try:
                    # 文字起こし実行
                    st.session_state.batch_processing_status[file_name]['transcription'] = 'processing'

                    file_data.seek(0)
                    seg_df = transcribe_audio_to_dataframe(
                        file_data,
                        reference_file=reference_file,
                        model=transcribe_model
                    )

                    st.session_state.batch_processing_results[file_name]['transcription_df'] = seg_df
                    st.session_state.batch_processing_status[file_name]['transcription'] = 'completed'
                    st.success(f"✅ 文字起こし完了: {file_name} ({len(seg_df)}行)")

                except Exception as e:
                    st.error(f"❌ エラー発生: {file_name} - {e}")
                    st.session_state.batch_processing_status[file_name]['transcription'] = 'error'
                    import traceback
                    st.error(traceback.format_exc())

                # 進捗バー更新
                progress_bar.progress((idx + 1) / total_files)

            status_text.write("✅ **すべてのファイルの文字起こしが完了しました！**")
            st.session_state.batch_current_step = 3
            st.balloons()
            st.rerun()

    # Step 3: 個別ファイルの話者識別
    if st.session_state.batch_current_step >= 3 and len(st.session_state.batch_extracted_files) > 0:
        st.subheader("Step 3: 個別ファイルの話者識別")
        st.write("各ファイルごとに話者識別を実行できます。必要に応じてファイルごとに異なる話者埋め込みを使用できます。")

        # 処理状況サマリー
        transcription_completed = sum(1 for status in st.session_state.batch_processing_status.values() if status.get('transcription') == 'completed')
        speaker_id_completed = sum(1 for status in st.session_state.batch_processing_status.values() if status.get('speaker_id') == 'completed')

        col1, col2, col3 = st.columns(3)
        with col1:
            st.metric("文字起こし完了", f"{transcription_completed}/{len(st.session_state.batch_extracted_files)}")
        with col2:
            st.metric("話者識別完了", f"{speaker_id_completed}/{len(st.session_state.batch_extracted_files)}")
        with col3:
            if st.button("🔄 話者識別をスキップして次へ", key="skip_speaker_id"):
                st.session_state.batch_current_step = 4
                st.rerun()

        st.divider()

        # ファイルごとにタブを作成
        file_names = [f['name'] for f in st.session_state.batch_extracted_files]
        file_tabs = st.tabs(file_names)

        for tab_idx, (file_tab, file_info) in enumerate(zip(file_tabs, st.session_state.batch_extracted_files)):
            with file_tab:
                selected_file_name = file_info['name']
                selected_file_info = file_info

                if selected_file_name in st.session_state.batch_processing_results:
                    result = st.session_state.batch_processing_results[selected_file_name]

                    # 文字起こし結果のプレビュー
                    if 'transcription_df' in result:
                        transcription_df = result['transcription_df']
                        with st.expander("📝 文字起こし結果プレビュー", expanded=True):
                            if 'speaker' in transcription_df.columns:
                                disabled_columns = [col for col in transcription_df.columns if col != 'speaker']
                                edited_transcription = st.data_editor(
                                    transcription_df,
                                    hide_index=True,
                                    use_container_width=True,
                                    num_rows="dynamic",
                                    disabled=disabled_columns,
                                    column_config={
                                        "speaker": st.column_config.TextColumn(
                                            "話者",
                                            help="必要に応じて話者名を直接編集してください"
                                        )
                                    },
                                    key=f"transcription_editor_{selected_file_name}_{tab_idx}"
                                )

                                if isinstance(edited_transcription, pd.DataFrame):
                                    updated_transcription_df = edited_transcription.copy()
                                else:
                                    updated_transcription_df = pd.DataFrame(edited_transcription)

                                updated_transcription_df = updated_transcription_df.reset_index(drop=True)
                                transcription_df_normalized = transcription_df.reset_index(drop=True)

                                # Ensure original columns exist in editor output
                                for column in transcription_df_normalized.columns:
                                    if column not in updated_transcription_df.columns:
                                        updated_transcription_df[column] = transcription_df_normalized[column]
                                updated_transcription_df = updated_transcription_df[transcription_df_normalized.columns]

                                original_speaker = transcription_df_normalized['speaker'].astype(str).fillna("")
                                updated_speaker = updated_transcription_df['speaker'].astype(str).fillna("")

                                if not updated_speaker.equals(original_speaker):
                                    st.session_state.batch_processing_results[selected_file_name]['transcription_df'] = updated_transcription_df
                                    transcription_df = updated_transcription_df
                                    st.session_state.batch_processing_results[selected_file_name].pop('identified_df', None)
                                    st.session_state.batch_processing_results[selected_file_name].pop('meeting_text', None)
                                    st.session_state.batch_processing_results[selected_file_name].pop('proofread_text', None)

                                    status_entry = st.session_state.batch_processing_status.get(selected_file_name, {})
                                    status_entry['speaker_id'] = 'pending'
                                    status_entry['rag_proofread'] = 'pending'
                                    st.session_state.batch_processing_status[selected_file_name] = status_entry

                                    st.success("✏️ 話者列の変更を保存しました。再度話者識別やRAG校正を実行してください。")
                            else:
                                st.dataframe(transcription_df.head(10), use_container_width=True)

                        tab_identify, tab_embed = st.tabs(["話者識別", "話者埋め込み作成"])

                        with tab_identify:
                            st.write("**話者識別設定**")

                            # デフォルトの埋め込みファイルがあれば表示
                            if st.session_state.batch_default_embeddings:
                                st.info(f"✅ 会議タイプから{len(st.session_state.batch_default_embeddings)}個の話者埋め込みを読み込み済み")
                                default_names = [emb['name'] for emb in st.session_state.batch_default_embeddings]
                                with st.expander("読み込み済み話者埋め込みファイル", expanded=True):
                                    st.write(", ".join(default_names))

                            col1, col2 = st.columns(2)

                            with col1:
                                # 追加の埋め込みファイルをアップロード
                                additional_embeddings = st.file_uploader(
                                    "追加の話者埋め込みファイル（.npy）",
                                    type=["npy"],
                                    accept_multiple_files=True,
                                    key=f"batch_additional_embeddings_{selected_file_name}_{tab_idx}",
                                    help="会議タイプで読み込んだファイルに加えて、追加で埋め込みファイルを指定できます"
                                )

                            with col2:
                                similarity_threshold = st.slider(
                                    "類似度閾値",
                                    min_value=0.0,
                                    max_value=1.0,
                                    value=0.7,
                                    step=0.01,
                                    key=f"batch_similarity_threshold_{selected_file_name}_{tab_idx}"
                                )

                            # 使用する埋め込みファイルを結合
                            # デフォルト埋め込み + 追加アップロード
                            all_embeddings = list(st.session_state.batch_default_embeddings)
                            if additional_embeddings:
                                all_embeddings.extend(additional_embeddings)

                            # 埋め込みファイル数を表示
                            total_embeddings = len(all_embeddings)
                            if total_embeddings > 0:
                                st.success(f"📊 合計 {total_embeddings}個の話者埋め込みファイルを使用します")
                            else:
                                st.warning("⚠️ 話者埋め込みファイルがありません。Step 2で会議タイプを選択するか、上記で追加してください。")

                            # 話者識別実行
                            if all_embeddings and st.button(f"🎤 {selected_file_name} の話者識別を実行", key=f"execute_speaker_id_{selected_file_name}_{tab_idx}", type="primary"):
                                with st.spinner(f"🎤 話者識別中: {selected_file_name}"):
                                    try:
                                        file_data = selected_file_info['data']
                                        file_data.seek(0)

                                        identified_df = identify_speakers_in_dataframe(
                                            file_data,
                                            transcription_df,
                                            all_embeddings,
                                            similarity_threshold
                                        )

                                        st.session_state.batch_processing_results[selected_file_name]['identified_df'] = identified_df
                                        st.session_state.batch_processing_status[selected_file_name]['speaker_id'] = 'completed'
                                        st.success(f"✅ 話者識別完了: {selected_file_name}")
                                        st.rerun()

                                    except Exception as e:
                                        st.error(f"❌ エラー発生: {e}")
                                        st.session_state.batch_processing_status[selected_file_name]['speaker_id'] = 'error'
                                        import traceback
                                        st.error(traceback.format_exc())

                            if 'identified_df' in result and not result['identified_df'].empty:
                                st.divider()
                                st.write("**話者識別結果の手動修正**")
                                st.caption("話者列を直接編集してラベルを調整できます。修正後はRAG校正やエクスポートを再実行してください。")

                                current_identified_df = result['identified_df'].copy()
                                if 'speaker' not in current_identified_df.columns:
                                    st.warning("話者列が見つかりません。識別結果を再実行してください。")
                                else:
                                    disabled_cols = [col for col in current_identified_df.columns if col != 'speaker']
                                    edited_df = st.data_editor(
                                        current_identified_df,
                                        column_config={
                                            "speaker": st.column_config.TextColumn(
                                                "話者",
                                                help="必要に応じて話者名を直接入力してください"
                                            )
                                        },
                                        disabled=disabled_cols,
                                        hide_index=True,
                                        use_container_width=True,
                                        key=f"identified_editor_{selected_file_name}_{tab_idx}"
                                    )

                                    if isinstance(edited_df, pd.DataFrame):
                                        updated_df = edited_df.copy()
                                    else:
                                        updated_df = pd.DataFrame(edited_df)

                                    updated_df = updated_df.reset_index(drop=True)
                                    current_normalized = current_identified_df.reset_index(drop=True)

                                    # Ensure column order matches original
                                    updated_df = updated_df[current_normalized.columns]

                                    if not updated_df.equals(current_normalized):
                                        st.session_state.batch_processing_results[selected_file_name]['identified_df'] = updated_df
                                        st.session_state.batch_processing_results[selected_file_name].pop('meeting_text', None)
                                        st.session_state.batch_processing_results[selected_file_name].pop('proofread_text', None)

                                        status_entry = st.session_state.batch_processing_status.get(selected_file_name, {})
                                        status_entry['speaker_id'] = 'completed'
                                        status_entry['rag_proofread'] = 'pending'
                                        st.session_state.batch_processing_status[selected_file_name] = status_entry

                                        st.success("✏️ 話者ラベルの変更を保存しました。")
                                        st.info("変更内容を反映するには、必要に応じてStep 4以降を再実行してください。")

                        with tab_embed:
                            st.write("音声セグメントを選択し、プレビュー確認してから話者埋め込みファイルを作成できます。")

                            target_df = result.get('identified_df', transcription_df).copy()

                            if len(target_df) == 0:
                                st.warning("文字起こし結果がありません。まず文字起こしを実行してください。")
                            else:
                                embedding_states = st.session_state.batch_embedding_states
                                if selected_file_name not in embedding_states:
                                    embedding_states[selected_file_name] = {
                                        'selected_rows': set(),
                                        'preview_audio': None,
                                        'show_preview': False
                                    }
                                embedding_state = embedding_states[selected_file_name]

                                selection_mode = st.radio(
                                    "選択方式を選んでください",
                                    options=["範囲指定モード", "チェックボックスモード"],
                                    horizontal=True,
                                    help="範囲指定: 連続した行を素早く選択 | チェックボックス: 飛び飛びの行を自由に選択",
                                    key=f"batch_embedding_selection_mode_{selected_file_name}_{tab_idx}"
                                )

                                st.divider()

                                st.subheader("1️⃣ 音声セグメントを選択")

                                if selection_mode == "範囲指定モード":
                                    embedding_row_labels = _create_row_labels(target_df)
                                    col1, col2 = st.columns(2)
                                    with col1:
                                        start_row = st.selectbox(
                                            "開始行を選択",
                                            options=range(len(target_df)),
                                            format_func=lambda x: embedding_row_labels[x],
                                            key=f"batch_embedding_start_row_{selected_file_name}_{tab_idx}"
                                        )
                                    with col2:
                                        end_row = st.selectbox(
                                            "終了行を選択",
                                            options=range(len(target_df)),
                                            format_func=lambda x: embedding_row_labels[x],
                                            index=len(target_df) - 1 if len(target_df) > 0 else 0,
                                            key=f"batch_embedding_end_row_{selected_file_name}_{tab_idx}"
                                        )

                                    if start_row > end_row:
                                        st.error("⚠️ 開始行は終了行以前を選択してください")
                                        embedding_state['selected_rows'] = set()
                                    else:
                                        embedding_state['selected_rows'] = set(range(start_row, end_row + 1))
                                        st.success(f"✅ 行 {start_row} ～ {end_row} を選択しました")
                                else:
                                    display_df = target_df.copy()
                                    display_df.insert(0, "選択", False)

                                    if embedding_state['selected_rows']:
                                        for idx in embedding_state['selected_rows']:
                                            if idx in display_df.index:
                                                display_df.at[idx, "選択"] = True

                                    disabled_columns = [col for col in ["start", "end", "speaker", "text"] if col in display_df.columns]

                                    edited_display = st.data_editor(
                                        display_df,
                                        column_config={
                                            "選択": st.column_config.CheckboxColumn(
                                                "選択",
                                                help="埋め込み作成対象の行をチェック",
                                                default=False
                                            )
                                        },
                                        disabled=disabled_columns,
                                        use_container_width=True,
                                        hide_index=True,
                                        key=f"batch_embedding_data_editor_{selected_file_name}_{tab_idx}"
                                    )

                                    embedding_state['selected_rows'] = set(
                                        edited_display[edited_display["選択"] == True].index.tolist()
                                    )

                                selected_rows = embedding_state['selected_rows']

                                if selected_rows:
                                    st.subheader("2️⃣ 選択情報の確認")

                                    selection_summary = get_selection_summary(target_df, selected_rows)

                                    info_cols = st.columns(4)
                                    with info_cols[0]:
                                        st.metric("選択行数", selection_summary['count'])
                                    with info_cols[1]:
                                        st.metric("開始時刻", format_time(selection_summary['start_time']))
                                    with info_cols[2]:
                                        st.metric("終了時刻", format_time(selection_summary['end_time']))
                                    with info_cols[3]:
                                        st.metric("音声長", f"{selection_summary['duration']:.1f}秒")

                                    if selection_summary['speakers']:
                                        speakers_text = "、".join([s if s else "不明" for s in selection_summary['speakers']])
                                        st.info(f"🎤 選択範囲に含まれる話者: {speakers_text}")

                                    st.subheader("3️⃣ 音声プレビュー")

                                    col_preview, col_clear = st.columns([3, 1])

                                    with col_preview:
                                        if st.button("🔊 プレビュー音声を生成", key=f"batch_generate_preview_{selected_file_name}_{tab_idx}"):
                                            try:
                                                with st.spinner("プレビュー音声を生成中..."):
                                                    audio_io = BytesIO(selected_file_info['data'].getvalue())
                                                    preview_bytes, duration = prepare_embedding_preview_audio(
                                                        audio_io,
                                                        target_df,
                                                        sorted(selected_rows)
                                                    )
                                                    embedding_state['preview_audio'] = preview_bytes
                                                    embedding_state['show_preview'] = True
                                                st.success(f"プレビュー音声を生成しました（{duration:.1f}秒）")
                                            except Exception as e:
                                                st.error(f"プレビュー生成エラー: {e}")

                                    with col_clear:
                                        if embedding_state.get('show_preview'):
                                            if st.button("❌ プレビューをクリア", key=f"batch_clear_preview_{selected_file_name}_{tab_idx}"):
                                                embedding_state['preview_audio'] = None
                                                embedding_state['show_preview'] = False
                                                st.rerun()

                                    if embedding_state.get('show_preview') and embedding_state.get('preview_audio'):
                                        st.audio(embedding_state['preview_audio'], format="audio/wav")

                                    st.subheader("4️⃣ ファイル名設定")

                                    embedding_filename = st.text_input(
                                        "ファイル名（.npy拡張子は自動追加）",
                                        value="speaker_embedding",
                                        key=f"batch_embedding_filename_{selected_file_name}_{tab_idx}",
                                        help="作成する話者埋め込みファイルの名前を指定してください"
                                    )

                                    st.subheader("5️⃣ 埋め込み作成")

                                    if st.button("✨ 話者埋め込みを作成してダウンロード", key=f"batch_create_embedding_{selected_file_name}_{tab_idx}", type="primary"):
                                        with st.spinner("話者埋め込みを作成中..."):
                                            try:
                                                audio_io = BytesIO(selected_file_info['data'].getvalue())
                                                embedding, duration = extract_audio_segment_for_embedding(
                                                    audio_io,
                                                    target_df,
                                                    sorted(selected_rows)
                                                )

                                                filename_with_ext = embedding_filename if embedding_filename.endswith('.npy') else f"{embedding_filename}.npy"

                                                embedding_io = BytesIO()
                                                np.save(embedding_io, embedding)
                                                embedding_io.seek(0)

                                                embedding_bytes = embedding_io.getvalue()

                                                st.success(f"✅ 話者埋め込みの作成が完了しました（音声長: {duration:.1f}秒）")
                                                trigger_auto_download(
                                                    embedding_bytes,
                                                    filename_with_ext,
                                                    key=f"batch_download_embedding_{selected_file_name}_{tab_idx}",
                                                    mime="application/octet-stream"
                                                )

                                            except Exception as e:
                                                st.error(f"❌ 話者埋め込みの作成中にエラーが発生しました: {e}")
                                                import traceback
                                                with st.expander("🔍 詳細なエラー情報"):
                                                    st.code(traceback.format_exc())
                                else:
                                    if selection_mode == "範囲指定モード":
                                        st.info("💡 上記のselectboxで開始行と終了行を選択してください")
                                    else:
                                        st.info("💡 データエディタで埋め込み作成対象の行をチェックしてください")

                    # 話者識別結果の確認
                    if 'identified_df' in result:
                        st.success("✅ このファイルの話者識別は完了しています")
                else:
                    st.warning("⚠️ このファイルはまだ文字起こしが完了していません")

        st.divider()

        # 全ファイルの処理状況一覧
        st.write("**全ファイルの処理状況**")
        status_df = pd.DataFrame([
            {
                'ファイル名': f['name'],
                '文字起こし': '✅' if st.session_state.batch_processing_status.get(f['name'], {}).get('transcription') == 'completed' else '❌',
                '話者識別': '✅' if st.session_state.batch_processing_status.get(f['name'], {}).get('speaker_id') == 'completed' else '⏭️'
            }
            for f in st.session_state.batch_extracted_files
        ])
        st.dataframe(status_df, use_container_width=True, hide_index=True)

        st.divider()

        # 次のステップへ
        if st.button("➡️ RAG校正へ進む", key="proceed_to_rag", type="primary"):
            st.session_state.batch_current_step = 4
            st.rerun()

    # Step 4: RAG校正
    if st.session_state.batch_current_step >= 4:
        st.subheader("Step 4: RAG校正")
        st.write("全ファイルの文字起こし結果（または話者識別結果）に対してRAG校正を実行できます。")

        # RAG校正設定
        col1, col2 = st.columns(2)

        with col1:
            # RAGシステムの初期化
            if st.session_state.batch_rag_system is None:
                st.session_state.batch_rag_system = _init_rag_system()

            rag_system = st.session_state.batch_rag_system

            # ナレッジベース読み込み
            st.write("**ナレッジベース設定**")

            # デフォルトDBの自動読み込み
            default_ragdb_path = get_default_ragdb_path()
            if os.path.exists(default_ragdb_path) and rag_system.vectorstore is None:
                with st.spinner("デフォルトナレッジベースを読み込み中..."):
                    try:
                        rag_system.load_knowledge_base(str(default_ragdb_path))
                        # データベース情報を更新
                        st.session_state.batch_db_info = rag_system.get_database_info()
                        st.success(f"✅ デフォルトDBを読み込みました")
                    except Exception as e:
                        st.warning(f"デフォルトDBの読み込みに失敗: {e}")

            # データベース情報を常に最新の状態に更新
            st.session_state.batch_db_info = rag_system.get_database_info()

            # データベース状態表示
            _render_database_status(st.session_state.batch_db_info)

            # 別のRAGDBファイルを読み込む
            uploaded_ragdb = st.file_uploader(
                "または別の.ragdbファイルを読み込み",
                type=["ragdb"],
                key="batch_ragdb_upload_step4"
            )

            if uploaded_ragdb and st.button("RAGDBを読み込む", key="batch_load_ragdb_step4"):
                with st.spinner("ナレッジベースを読み込み中..."):
                    try:
                        with tempfile.NamedTemporaryFile(delete=False, suffix=".ragdb") as temp_file:
                            temp_file.write(uploaded_ragdb.read())
                            temp_ragdb_path = temp_file.name

                        rag_system.load_knowledge_base(temp_ragdb_path)
                        os.unlink(temp_ragdb_path)
                        # データベース情報を更新
                        st.session_state.batch_db_info = rag_system.get_database_info()
                        st.success("✅ ナレッジベースを読み込みました")
                        st.rerun()
                    except Exception as e:
                        st.error(f"❌ 読み込みエラー: {e}")

        with col2:
            # RAG校正パラメータ
            st.write("**RAG校正パラメータ**")

            # プロンプトプリセット選択
            prompt_loader = get_default_loader()
            presets = prompt_loader.get_preset_list('rag_proofreading')
            preset_options = {p['id']: f"{p['name']} - {p['description']}" for p in presets}

            selected_preset = st.selectbox(
                "プロンプトプリセット",
                options=list(preset_options.keys()),
                format_func=lambda x: preset_options[x],
                index=0,
                key="batch_rag_preset_step4"
            )

            search_type = st.selectbox(
                "検索タイプ",
                options=["similarity", "mmr"],
                index=0,
                key="batch_search_type_step4",
                help="similarity: 類似度検索 | mmr: 多様性考慮検索"
            )

            llm_model = st.selectbox(
                "LLMモデル",
                options=["gpt-4o", "gpt-4o-mini"],
                index=1,
                key="batch_llm_model_step4"
            )

            top_k = st.slider(
                "検索する関連文脈の数",
                min_value=1,
                max_value=10,
                value=6,
                key="batch_top_k_step4"
            )

        st.divider()

        # RAG校正実行
        if rag_system.vectorstore is not None:
            col1, col2 = st.columns(2)
            with col1:
                if st.button("📚 RAG校正を開始", key="batch_start_rag", type="primary"):
                    progress_bar = st.progress(0)
                    status_text = st.empty()

                    total_files = len(st.session_state.batch_extracted_files)
                    processed = 0
                    errors = 0

                    for idx, file_info in enumerate(st.session_state.batch_extracted_files):
                        file_name = file_info['name']
                        result = st.session_state.batch_processing_results.get(file_name, {})

                        status_text.write(f"**RAG校正中: {file_name}** ({idx + 1}/{total_files})")

                        try:
                            # 話者識別済みならそれを使用、なければ文字起こし結果を使用
                            if 'identified_df' in result:
                                df = result['identified_df']
                            elif 'transcription_df' in result:
                                df = result['transcription_df']
                            else:
                                st.warning(f"⏭️ スキップ: {file_name}（文字起こし結果がありません）")
                                continue

                            meeting_text = build_meeting_text_from_dataframe(df)
                            st.session_state.batch_processing_results.setdefault(file_name, {})['meeting_text'] = meeting_text

                            # RAG校正実行
                            st.session_state.batch_processing_status[file_name]['rag_proofread'] = 'processing'

                            proofread_result = rag_system.rag_enhanced_proofread(
                                meeting_text,
                                search_type=search_type,
                                top_k=top_k,
                                model=llm_model,
                                prompt_preset=selected_preset
                            )

                            st.session_state.batch_processing_results[file_name]['proofread_text'] = proofread_result
                            st.session_state.batch_processing_status[file_name]['rag_proofread'] = 'completed'
                            st.success(f"✅ RAG校正完了: {file_name}")
                            processed += 1

                        except Exception as e:
                            st.error(f"❌ エラー発生: {file_name} - {e}")
                            st.session_state.batch_processing_status[file_name]['rag_proofread'] = 'error'
                            errors += 1
                            import traceback
                            st.error(traceback.format_exc())

                        progress_bar.progress((idx + 1) / total_files)

                    # 完了メッセージ
                    if errors == 0:
                        status_text.write(f"✅ **RAG校正完了！ ({processed}/{total_files}ファイル処理)**")
                        st.session_state.batch_current_step = 5
                        st.balloons()
                        st.rerun()
                    else:
                        status_text.write(f"⚠️ **RAG校正完了（エラーあり）: 成功 {processed}件 / エラー {errors}件**")
                        st.warning("⚠️ エラーが発生したファイルがあります。上記のエラーメッセージを確認してください。")
                        st.info("💡 修正後、再度「RAG校正を開始」ボタンを押すか、「RAG校正をスキップして次へ」で進んでください。")

            with col2:
                if st.button("🔄 RAG校正をスキップして次へ", key="skip_rag"):
                    st.session_state.batch_current_step = 5
                    st.rerun()
        else:
            st.warning("⚠️ ナレッジベースが読み込まれていません。RAG校正をスキップする場合は下のボタンをクリックしてください。")
            if st.button("🔄 RAG校正をスキップして次へ", key="skip_rag_no_kb"):
                st.session_state.batch_current_step = 5
                st.rerun()

    # Step 5: 処理結果の確認とダウンロード
    if st.session_state.batch_current_step >= 5:
        st.subheader("Step 5: 処理結果の確認とダウンロード")

        # 処理状態サマリー
        st.write("**処理状態サマリー**")
        status_summary = []
        for file_name, status in st.session_state.batch_processing_status.items():
            status_summary.append({
                'ファイル名': file_name,
                '文字起こし': '✅' if status['transcription'] == 'completed' else '❌' if status['transcription'] == 'error' else '⏭️',
                '話者識別': '✅' if status['speaker_id'] == 'completed' else '❌' if status['speaker_id'] == 'error' else '⏭️',
                'RAG校正': '✅' if status['rag_proofread'] == 'completed' else '❌' if status['rag_proofread'] == 'error' else '⏭️'
            })

        st.dataframe(pd.DataFrame(status_summary), use_container_width=True, hide_index=True)

        st.divider()

        # 個別ファイルの結果表示
        st.write("**個別ファイルの結果**")

        # ファイルごとにタブを作成
        file_names = [f['name'] for f in st.session_state.batch_extracted_files]
        file_tabs = st.tabs(file_names)

        for tab_idx, (file_tab, file_info) in enumerate(zip(file_tabs, st.session_state.batch_extracted_files)):
            with file_tab:
                selected_file = file_info['name']

                if selected_file in st.session_state.batch_processing_results:
                    result = st.session_state.batch_processing_results[selected_file]

                    content_tab1, content_tab2, content_tab3 = st.tabs(["📝 文字起こし結果", "🎤 話者識別結果", "📚 RAG校正結果"])

                    with content_tab1:
                        if 'transcription_df' in result:
                            st.dataframe(result['transcription_df'], use_container_width=True)

                            # Excel形式でダウンロード
                            excel_buffer = BytesIO()
                            with pd.ExcelWriter(excel_buffer, engine='openpyxl') as writer:
                                result['transcription_df'].to_excel(writer, index=False, sheet_name='文字起こし')
                            excel_buffer.seek(0)

                            st.download_button(
                                label="📥 Excel形式でダウンロード",
                                data=excel_buffer,
                                file_name=f"{os.path.splitext(selected_file)[0]}_transcription.xlsx",
                                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                                key=f"download_excel_{selected_file}_{tab_idx}"
                            )
                        else:
                            st.info("文字起こし結果がありません")

                    with content_tab2:
                        meeting_text = result.get('meeting_text')
                        if not meeting_text:
                            source_df = result.get('identified_df') or result.get('transcription_df')
                            meeting_text = build_meeting_text_from_dataframe(source_df) if source_df is not None else ""
                            if meeting_text:
                                st.session_state.batch_processing_results[selected_file]['meeting_text'] = meeting_text

                        if meeting_text:
                            st.text_area(
                                "議事録形式テキスト",
                                meeting_text,
                                height=400,
                                key=f"meeting_text_view_{selected_file}_{tab_idx}"
                            )

                            doc = DocxDocument()
                            doc.add_heading(f'議事録: {selected_file}', level=1)
                            for line in meeting_text.splitlines():
                                doc.add_paragraph(line)

                            docx_buffer = BytesIO()
                            doc.save(docx_buffer)
                            docx_buffer.seek(0)

                            st.download_button(
                                label="📥 議事録（Word）をダウンロード",
                                data=docx_buffer,
                                file_name=f"{os.path.splitext(selected_file)[0]}_minutes.docx",
                                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                                key=f"download_docx_{selected_file}_{tab_idx}"
                            )
                        else:
                            st.info("議事録形式テキストがありません（話者識別をスキップした場合はStep 2の結果を確認してください）")

                    with content_tab3:
                        if 'proofread_text' in result:
                            st.text_area(
                                "校正後テキスト",
                                value=result['proofread_text'],
                                height=400,
                                key=f"proofread_view_{selected_file}_{tab_idx}"
                            )

                            # Word形式でダウンロード
                            doc = DocxDocument()
                            doc.add_heading(f'校正済み議事録: {selected_file}', level=1)
                            doc.add_paragraph(result['proofread_text'])

                            docx_buffer = BytesIO()
                            doc.save(docx_buffer)
                            docx_buffer.seek(0)

                            st.download_button(
                                label="📥 校正済み議事録（Word）をダウンロード",
                                data=docx_buffer,
                                file_name=f"{os.path.splitext(selected_file)[0]}_proofread.docx",
                                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                                key=f"download_proofread_{selected_file}_{tab_idx}"
                            )
                        else:
                            st.info("RAG校正結果がありません（スキップされたか、エラーが発生しました）")
                else:
                    st.info("このファイルの処理結果がありません")

        st.divider()

        # 一括ダウンロード機能
        st.write("**一括ダウンロード**")
        st.write("すべての処理結果をまとめてダウンロードできます。")

        col1, col2, col3 = st.columns(3)

        with col1:
            if st.button("📦 ZIPを作成してダウンロード", key="batch_download_all", use_container_width=True):
                with st.spinner("ZIPファイルを作成中..."):
                    try:
                        zip_buffer = BytesIO()

                        with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zip_file:
                            for file_name, result in st.session_state.batch_processing_results.items():
                                base_name = os.path.splitext(file_name)[0]

                                # 文字起こし結果（Excel）
                                if 'transcription_df' in result:
                                    excel_buffer = BytesIO()
                                    with pd.ExcelWriter(excel_buffer, engine='openpyxl') as writer:
                                        result['transcription_df'].to_excel(writer, index=False, sheet_name='文字起こし')
                                    excel_buffer.seek(0)
                                    zip_file.writestr(f"{base_name}_transcription.xlsx", excel_buffer.read())

                                # 話者識別結果（Word）
                                if 'identified_df' in result:
                                    doc = DocxDocument()
                                    doc.add_heading(f'議事録: {file_name}', level=1)

                                    # 連続する同じ話者の発言を結合
                                    df = result['identified_df'].copy()
                                    df['speaker_filled'] = df['speaker'].replace('', pd.NA)
                                    df['speaker_filled'] = df['speaker_filled'].ffill()
                                    df['group_id'] = (df['speaker_filled'] != df['speaker_filled'].shift()).cumsum()
                                    df_merged = df.groupby('group_id').agg(
                                        speaker=('speaker_filled', 'first'),
                                        text=('text', ' '.join)
                                    ).reset_index(drop=True)

                                    for _, row in df_merged.iterrows():
                                        speaker = row.get('speaker', '不明')
                                        text = row.get('text', '')
                                        speaker_str = speaker if speaker else '不明'
                                        doc.add_paragraph(f"（{speaker_str}）{text}")

                                    docx_buffer = BytesIO()
                                    doc.save(docx_buffer)
                                    docx_buffer.seek(0)
                                    zip_file.writestr(f"{base_name}_minutes.docx", docx_buffer.read())

                                # RAG校正結果（Word）
                                if 'proofread_text' in result:
                                    doc = DocxDocument()
                                    doc.add_heading(f'校正済み議事録: {file_name}', level=1)
                                    doc.add_paragraph(result['proofread_text'])
                                    docx_buffer = BytesIO()
                                    doc.save(docx_buffer)
                                    docx_buffer.seek(0)
                                    zip_file.writestr(f"{base_name}_proofread.docx", docx_buffer.read())

                        zip_buffer.seek(0)

                        trigger_auto_download(
                            zip_buffer.getvalue(),
                            file_name=f"batch_results_{datetime.now().strftime('%Y%m%d_%H%M%S')}.zip",
                            key="batch_final_download",
                            mime="application/zip"
                        )

                        st.success("✅ ZIPファイルの作成が完了しました！")

                    except Exception as e:
                        st.error(f"❌ ZIPファイルの作成中にエラーが発生しました: {e}")
                        import traceback
                        st.error(traceback.format_exc())

        with col2:
            if st.button("📝 文字起こし結果Wordを作成してダウンロード", key="batch_download_transcription_word", use_container_width=True):
                with st.spinner("Wordファイルを作成中..."):
                    try:
                        # 1つのWordドキュメントを作成
                        doc = DocxDocument()
                        doc.add_heading('文字起こし結果（一括）', level=0)
                        doc.add_paragraph(f'作成日時: {datetime.now().strftime("%Y年%m月%d日 %H:%M:%S")}')
                        doc.add_paragraph('')

                        # 各ファイルの結果を追加
                        for idx, (file_name, result) in enumerate(st.session_state.batch_processing_results.items(), 1):
                            # ファイル見出し
                            doc.add_heading(f'{idx}. {file_name}', level=1)
                            doc.add_paragraph('')

                            base_df = result.get('identified_df')
                            if base_df is None or base_df.empty:
                                fallback_df = result.get('transcription_df')
                                base_df = fallback_df if fallback_df is not None else None
                            meeting_text = result.get('meeting_text')
                            if meeting_text is None and base_df is not None and not base_df.empty:
                                meeting_text = build_meeting_text_from_dataframe(base_df)
                                if meeting_text:
                                    st.session_state.batch_processing_results[file_name]['meeting_text'] = meeting_text

                            if meeting_text:
                                doc.add_heading('議事録（話者識別結果）', level=2)
                                for line in meeting_text.splitlines():
                                    if line.strip():
                                        doc.add_paragraph(line)
                            else:
                                doc.add_paragraph('（文字起こし結果なし）')

                            doc.add_paragraph('')

                            # ファイル間の区切り
                            if idx < len(st.session_state.batch_processing_results):
                                doc.add_page_break()

                        # Wordファイルを保存
                        docx_buffer = BytesIO()
                        doc.save(docx_buffer)
                        docx_buffer.seek(0)

                        trigger_auto_download(
                            docx_buffer.getvalue(),
                            file_name=f"batch_transcription_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx",
                            key="batch_final_download_transcription_word",
                            mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                        )

                        st.success("✅ Wordファイルの作成が完了しました！")

                    except Exception as e:
                        st.error(f"❌ Wordファイルの作成中にエラーが発生しました: {e}")
                        import traceback
                        st.error(traceback.format_exc())

        with col3:
            if st.button("📚 RAG校正結果Wordを作成してダウンロード", key="batch_download_rag_word", use_container_width=True):
                with st.spinner("Wordファイルを作成中..."):
                    try:
                        # 1つのWordドキュメントを作成
                        doc = DocxDocument()
                        doc.add_heading('RAG校正結果（一括）', level=0)
                        doc.add_paragraph(f'作成日時: {datetime.now().strftime("%Y年%m月%d日 %H:%M:%S")}')
                        doc.add_paragraph('')

                        # RAG校正結果があるファイルのみ処理
                        rag_count = 0
                        for idx, (file_name, result) in enumerate(st.session_state.batch_processing_results.items(), 1):
                            if 'proofread_text' in result:
                                rag_count += 1
                                # ファイル見出し
                                doc.add_heading(f'{rag_count}. {file_name}', level=1)
                                doc.add_paragraph('')

                                # RAG校正結果
                                doc.add_heading('校正済み議事録', level=2)
                                # 段落ごとに分割して追加
                                for paragraph in result['proofread_text'].split('\n'):
                                    if paragraph.strip():
                                        doc.add_paragraph(paragraph)
                                doc.add_paragraph('')

                                # ファイル間の区切り（最後のファイル以外）
                                if rag_count < sum(1 for r in st.session_state.batch_processing_results.values() if 'proofread_text' in r):
                                    doc.add_page_break()

                        if rag_count == 0:
                            doc.add_paragraph('（RAG校正結果がありません）')

                        # Wordファイルを保存
                        docx_buffer = BytesIO()
                        doc.save(docx_buffer)
                        docx_buffer.seek(0)

                        trigger_auto_download(
                            docx_buffer.getvalue(),
                            file_name=f"batch_rag_proofread_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx",
                            key="batch_final_download_rag_word",
                            mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                        )

                        st.success("✅ Wordファイルの作成が完了しました！")

                    except Exception as e:
                        st.error(f"❌ Wordファイルの作成中にエラーが発生しました: {e}")
                        import traceback
                        st.error(traceback.format_exc())


def video_to_audio_cutter_app():
    st.title("動画から音声を切り出しMP3で保存")
    st.write("動画ファイルをアップロードし、切り出したい開始時間と終了時間を指定してください。複数の区間を切り出すことができます。")

    uploaded_video = st.file_uploader("動画ファイルを選択", type=["wav","mp3","mp4", "mov", "avi", "mkv", "webm"])

    if uploaded_video is not None:
        st.video(uploaded_video)

        st.subheader("切り出し区間の設定")
        # Use st.data_editor for multiple time range inputs
        # Default for the first row includes segment_1
        default_data = pd.DataFrame([
            {"開始時間": "00:00:00", "終了時間": "00:00:30", "出力ファイル名": f"{os.path.splitext(uploaded_video.name)[0]}_"}
        ])
        edited_df = st.data_editor(
            default_data,
            num_rows="dynamic",
            use_container_width=True,
            column_config={
                "開始時間": st.column_config.TextColumn(
                    "開始時間 (HH:MM:SS or seconds)",
                    help="切り出し開始時間 (例: 00:00:10 または 10)",
                    default="00:00:00"
                ),
                "終了時間": st.column_config.TextColumn(
                    "終了時間 (HH:MM:SS or seconds)",
                    help="切り出し終了時間 (例: 00:00:30 または 30)",
                    default="00:00:30"
                ),
                "出力ファイル名": st.column_config.TextColumn(
                    "出力ファイル名 (.mp3)",
                    help="この区間のMP3出力ファイル名を入力してください (例: my_audio_segment.mp3)。'AUTO_GENERATE'と入力するか空欄の場合、自動で連番が振られます。",
                    default=f"{os.path.splitext(uploaded_video.name)[0]}_" # Explicit placeholder for new rows
                )
            }
        )

        if st.button("音声を切り出してMP3で保存"):
            if edited_df.empty:
                st.warning("切り出し区間が設定されていません。")
                return

            temp_video_path = ""
            output_audio_paths = [] # List to store paths of all generated MP3s
            zip_buffer = BytesIO()

            try:
                # Save uploaded video to a temporary file
                with tempfile.NamedTemporaryFile(delete=False, suffix=f".{uploaded_video.name.split('.')[-1]}") as temp_video_file:
                    temp_video_file.write(uploaded_video.read())
                    temp_video_path = temp_video_file.name

                with st.spinner("音声の切り出しとMP3への変換中..."):
                    for index, row in edited_df.iterrows():
                        start_time_str = str(row["開始時間"])
                        end_time_str = str(row["終了時間"])
                        output_filename_raw = str(row["出力ファイル名"]).strip()

                        try:
                            start_seconds = parse_time_to_seconds(start_time_str)
                            end_seconds = parse_time_to_seconds(end_time_str)

                            if start_seconds >= end_seconds:
                                st.error(f"区間 {index+1}: 開始時間 ({start_time_str}) は終了時間 ({end_time_str}) より前に設定してください。この区間はスキップされます。")
                                continue

                            # If output filename is empty or matches the explicit placeholder, generate one with index
                            base_name_from_video = os.path.splitext(uploaded_video.name)[0]

                            if not output_filename_raw or output_filename_raw.upper() == "AUTO_GENERATE":
                                output_filename_to_use = f"{base_name_from_video}_segment_{index+1}.mp3"
                            else:
                                output_filename_to_use = output_filename_raw

                            # Ensure the output filename ends with .mp3
                            if not output_filename_to_use.lower().endswith(".mp3"):
                                output_filename_to_use += ".mp3"

                            output_audio_path = os.path.join(tempfile.gettempdir(), output_filename_to_use)

                            command = [
                                "ffmpeg",
                                "-i", temp_video_path,
                                "-ss", format_time(start_seconds),
                                "-to", format_time(end_seconds),
                                "-vn",  # No video
                                "-ab", "192k", # Audio bitrate
                                "-map_metadata", "-1", # Remove metadata
                                "-y", # Overwrite output files without asking
                                output_audio_path
                            ]

                            st.info(f"区間 {index+1} FFmpegコマンドを実行中: {' '.join(command)}")

                            process = subprocess.run(command, capture_output=True, text=True, encoding="utf-8", check=True)
                            st.success(f"区間 {index+1} の音声切り出しとMP3への変換が完了しました！")
                            st.code(process.stdout)
                            st.code(process.stderr)
                            output_audio_paths.append(output_audio_path)

                        except subprocess.CalledProcessError as e:
                            st.error(f"区間 {index+1} FFmpegの実行中にエラーが発生しました: {e}")
                            st.code(e.stdout)
                            st.code(e.stderr)
                            st.warning("FFmpegがシステムにインストールされ、PATHが通っていることを確認してください。")
                        except ValueError as e:
                            st.error(f"区間 {index+1} 時間形式エラー: {e}")
                        except Exception as e:
                            st.error(f"区間 {index+1} 処理中にエラーが発生しました: {e}")

                if output_audio_paths:
                    st.subheader("生成されたMP3ファイル")
                    with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zf:
                        for audio_path in output_audio_paths:
                            if os.path.exists(audio_path):
                                zf.write(audio_path, os.path.basename(audio_path))
                                st.write(f"- {os.path.basename(audio_path)}")
                    zip_buffer.seek(0)

                    st.download_button(
                        label="全てのMP3ファイルをまとめてダウンロード (ZIP)",
                        data=zip_buffer,
                        file_name=f"{os.path.splitext(uploaded_video.name)[0]}_cut_audios.zip",
                        mime="application/zip"
                    )
                else:
                    st.warning("切り出されたMP3ファイルはありませんでした。")

            except Exception as e:
                st.error(f"動画ファイルの処理中にエラーが発生しました: {e}")
            finally:
                # Clean up temporary files
                if os.path.exists(temp_video_path):
                    os.remove(temp_video_path)
                for audio_path in output_audio_paths:
                    if os.path.exists(audio_path):
                        os.remove(audio_path)

def _init_session_state(defaults):
    """セッション状態を初期化"""
    for key, value in defaults.items():
        if key not in st.session_state:
            st.session_state[key] = value

def _create_row_labels(df):
    """データフレームから行選択用ラベルを作成"""
    labels = []
    for idx, row in df.iterrows():
        text = row.get('text', '')
        text_display = text[:30] + "..." if len(text) > 30 else text
        label = f"{idx}: {row.get('start', '')} | {row.get('end', '')} | {row.get('speaker', '')} | {text_display}"
        labels.append(label)
    return labels

def prepare_embedding_preview_audio(audio_io, df, selected_indices):
    """プレビュー用の音声を準備

    Args:
        audio_io: BytesIO - 音声ファイル
        df: DataFrame - タイムスタンプ含む
        selected_indices: list - 選択行インデックスリスト

    Returns:
        preview_bytes: bytes - WAV形式の音声バイト列
        duration: float - セグメント長（秒）
    """
    if not selected_indices:
        raise ValueError("選択行がありません")

    selected_rows = df.iloc[sorted(selected_indices)]
    start_sec = selected_rows['start'].min()
    end_sec = selected_rows['end'].max()

    audio_io.seek(0)
    with temp_file_path(audio_io.getvalue(), ".mp3") as audio_path:
        audio = AudioSegment.from_file(audio_path)

        start_ms = int(start_sec * 1000)
        end_ms = int(end_sec * 1000)
        audio_segment = audio[start_ms:end_ms]

        preview_bytes = audio_segment.export(format="wav").read()

    return preview_bytes, (end_sec - start_sec)

def extract_audio_segment_for_embedding(audio_io, df, selected_indices):
    """選択行から音声セグメントを抽出して埋め込みを作成

    Args:
        audio_io: BytesIO - 音声ファイル
        df: DataFrame - タイムスタンプ含む
        selected_indices: list - 選択行インデックスリスト

    Returns:
        embedding: np.ndarray - 生成された埋め込みベクトル
        duration: float - セグメント長（秒）
    """
    if not selected_indices:
        raise ValueError("選択行がありません")

    selected_rows = df.iloc[sorted(selected_indices)]
    start_sec = selected_rows['start'].min()
    end_sec = selected_rows['end'].max()

    audio_io.seek(0)
    with temp_file_path(audio_io.getvalue(), ".mp3") as audio_path:
        audio = AudioSegment.from_file(audio_path)

        start_ms = int(start_sec * 1000)
        end_ms = int(end_sec * 1000)
        audio_segment = audio[start_ms:end_ms]

        wav_bytes = audio_segment.export(format="wav").read()

    with temp_file_path(wav_bytes, ".wav") as segment_path:
        wav = preprocess_wav(segment_path)
        encoder = load_voice_encoder()
        embedding = encoder.embed_utterance(wav)

    return embedding, (end_sec - start_sec)

def get_selection_summary(df, selected_indices):
    """選択行の概要情報を取得

    Args:
        df: DataFrame
        selected_indices: set or list - 選択行インデックス

    Returns:
        dict - 概要情報
    """
    if not selected_indices:
        return {
            'count': 0,
            'start_time': None,
            'end_time': None,
            'duration': 0,
            'speakers': []
        }

    sorted_indices = sorted(selected_indices)
    selected_rows = df.iloc[sorted_indices]

    start_time = selected_rows['start'].min()
    end_time = selected_rows['end'].max()
    duration = end_time - start_time
    speakers = selected_rows['speaker'].dropna().unique().tolist() if 'speaker' in selected_rows.columns else []

    return {
        'count': len(sorted_indices),
        'start_time': start_time,
        'end_time': end_time,
        'duration': duration,
        'speakers': speakers
    }

def get_default_ragdb_path():
    """デフォルトRAGDBファイルのフルパスを取得

    Returns:
        str - RAGDBファイルのフルパス
    """
    if DEFAULT_RAGDB_FOLDER:
        # フォルダが存在しない場合は作成
        os.makedirs(DEFAULT_RAGDB_FOLDER, exist_ok=True)
        return os.path.join(DEFAULT_RAGDB_FOLDER, "default_knowledge_base.ragdb")
    else:
        return "default_knowledge_base.ragdb"

def load_meeting_type_config():
    """会議タイプ設定を取得

    Returns:
        list - 会議タイプのリスト
    """
    try:
        return [dict(item) for item in DEFAULT_MEETING_TYPES]
    except Exception as e:
        st.error(f"❌ 会議タイプ設定の初期化エラー: {e}")
        return []

def load_embeddings_from_folder(folder_path):
    """指定されたフォルダから話者埋め込みファイル（.npy）を一括読み込み

    Args:
        folder_path (str): 話者埋め込みファイルが格納されているフォルダパス

    Returns:
        list - ファイル情報のリスト [{'name': filename, 'data': file_content_bytes}, ...]
                読み込めない場合は空リスト
    """
    embeddings = []

    try:
        # 相対パスの場合は絶対パスに変換
        if not os.path.isabs(folder_path):
            script_dir = os.path.dirname(os.path.abspath(__file__))
            folder_path = os.path.join(script_dir, folder_path)

        if not os.path.exists(folder_path):
            st.warning(f"⚠️ 話者埋め込みフォルダが見つかりません: {folder_path}")
            return []

        # フォルダ内の.npyファイルを検索
        npy_files = [f for f in os.listdir(folder_path) if f.endswith('.npy')]

        if not npy_files:
            st.info(f"📁 フォルダ内に.npyファイルが見つかりませんでした: {folder_path}")
            return []

        # 各.npyファイルを読み込み
        for filename in npy_files:
            file_path = os.path.join(folder_path, filename)
            with open(file_path, 'rb') as f:
                file_bytes = f.read()
                embeddings.append({
                    'name': filename,
                    'data': file_bytes
                })

        return embeddings

    except Exception as e:
        st.error(f"❌ 話者埋め込みファイルの読み込みエラー: {e}")
        return []

def video_transcribe_and_identify():
    st.title("文字起こし")
    st.write("音声・動画ファイルから自動で文字起こしを行い、議事録を作成します。")

    st.sidebar.markdown("""
    ### 📝 文字起こし

    **概要**
    文字起こしと話者識別に対応。モデルを選択可能。

    **主な機能**
    - AI文字起こし: タイムスタンプ付き
    - モデル選択可能:
      - gpt-4o-transcribe-diarize（自動話者識別付き）
      - Whisper（参考資料対応、後から話者識別可能）
    - 議事録出力: Word・Excel形式

    **対応形式:** MP3, WAV, M4A, MP4, WebM, OGG, MOV, AVI, MKV
    """)

    # セッション状態の初期化
    _init_session_state({
        'video_combined_step': 1,
        'video_combined_audio_io': None,
        'video_combined_df': None,
        'video_combined_identified_df': pd.DataFrame(),
        'video_combined_uploaded_file_name': "",
        # 話者埋め込み作成用の状態変数
        'embedding_selected_rows': set(),
        'embedding_preview_audio': None,
        'embedding_show_preview': False
    })

    # Step 1: メディアファイルのアップロード
    st.subheader("Step 1: メディアファイルのアップロード")
    uploaded_media = st.file_uploader(
        "メディアファイルを選択してください",
        type=["mp3", "wav", "m4a", "mp4", "webm", "ogg", "mov", "avi", "mkv"],
        key="video_combined_media_upload"
    )

    if uploaded_media is not None:
        st.session_state.video_combined_uploaded_file_name = uploaded_media.name
        st.success(f"ファイル '{uploaded_media.name}' がアップロードされました")

        # メディアプレビュー（動画または音声）
        file_extension = uploaded_media.name.split(".")[-1].lower()
        if file_extension in ["mp4", "webm", "ogg", "mov", "avi", "mkv"]:
            st.video(uploaded_media)
        else:
            st.audio(uploaded_media)

        # Step 2: 文字起こし
        st.subheader("Step 2: 文字起こし")

        # モデル選択
        transcribe_model = st.selectbox(
            "文字起こしモデルを選択してください",
            options=["gpt-4o-transcribe-diarize", "whisper"],
            index=1,  # デフォルト: whisper
            key="video_combined_model_selection",
            help="gpt-4o-transcribe-diarize: 話者識別付き（参考資料非対応）\nwhisper: 話者識別なし（参考資料対応）"
        )

        # モデルの説明
        if transcribe_model == "gpt-4o-transcribe-diarize":
            st.info("🎯 **gpt-4o-transcribe-diarize**: 自動話者識別付き文字起こし。参考資料（prompt）は使用できません。")
        else:
            st.info("🎤 **Whisper**: 汎用音声認識モデル。参考資料を使用可能。Step 4で話者識別を後から実行できます。")

        # 参考資料ファイルアップロード
        reference_file = st.file_uploader(
            "参考資料ファイルを選択（オプション）",
            type=["pdf", "docx", "pptx", "txt", "msg"],
            key="video_combined_reference_file",
            help="文字起こし精度向上のための参考資料（Whisperのみサポート）"
        )

        if st.button("文字起こしを実行", key="video_combined_transcribe"):
            with st.spinner("文字起こし中..."):
                try:
                    # アップロードされたメディアファイルを直接使用
                    uploaded_media.seek(0)
                    audio_file_io = BytesIO(uploaded_media.read())
                    audio_file_io.name = uploaded_media.name

                    seg_df = transcribe_audio_to_dataframe(audio_file_io, reference_file=reference_file, model=transcribe_model)

                    st.session_state.video_combined_df = seg_df
                    st.session_state.video_combined_step = 3

                    # 音声データをセッションに保存（話者識別用）
                    uploaded_media.seek(0)
                    st.session_state.video_combined_audio_io = BytesIO(uploaded_media.read())
                    st.session_state.video_combined_audio_io.name = uploaded_media.name

                    st.success(f"文字起こしが完了しました（{len(seg_df)}行のデータ）")
                    st.rerun()
                except Exception as e:
                    st.error(f"文字起こし中にエラーが発生しました: {e}")
                    import traceback
                    st.error(traceback.format_exc())

        # Step 3: 結果の確認と編集
        if st.session_state.video_combined_step >= 3 and st.session_state.video_combined_df is not None:
            st.subheader("Step 3: 結果の確認と編集")

            if len(st.session_state.video_combined_df) > 0:
                st.write("文字起こし結果を確認・編集してください:")
                base_df = st.session_state.video_combined_df.copy()
                editable_columns = [col for col in base_df.columns if col in ("speaker",)]
                disabled_cols = [col for col in base_df.columns if col not in editable_columns]

                edited_df = st.data_editor(
                    base_df,
                    num_rows="dynamic",
                    use_container_width=True,
                    disabled=disabled_cols,
                    column_config={
                        "speaker": st.column_config.TextColumn(
                            "話者",
                            help="必要に応じて話者名を手動で調整してください"
                        )
                    },
                    key="video_combined_editor"
                )

                if isinstance(edited_df, pd.DataFrame):
                    updated_df = edited_df.copy()
                else:
                    updated_df = pd.DataFrame(edited_df)

                updated_df = updated_df.reset_index(drop=True)
                base_df = base_df.reset_index(drop=True)

                # Ensure all original columns remain
                for column in base_df.columns:
                    if column not in updated_df.columns:
                        updated_df[column] = base_df[column]
                updated_df = updated_df[base_df.columns]

                original_speaker = base_df['speaker'].astype(str).fillna("") if 'speaker' in base_df.columns else None
                updated_speaker = updated_df['speaker'].astype(str).fillna("") if 'speaker' in updated_df.columns else None

                if original_speaker is None or not updated_speaker.equals(original_speaker):
                    st.session_state.video_combined_df = updated_df
                    st.session_state.video_combined_identified_df = updated_df.copy()
                    st.info("✏️ 話者ラベルの変更を保存しました。必要に応じてStep 4以降を再実行してください。")
            else:
                st.warning("データフレームが空です。文字起こし処理でエラーが発生した可能性があります。")

        # Step 4: 話者識別
        if st.session_state.video_combined_step >= 3 and st.session_state.video_combined_df is not None:
            st.subheader("Step 4: 話者識別")

            # タブで話者識別と埋め込み作成を分離
            tab1, tab2 = st.tabs(["話者識別", "話者埋め込み作成"])

            with tab1:
                st.write("話者識別を行う範囲を選択してください:")

                # 最新の編集内容を使用
                edited_df = st.session_state.video_combined_df.copy()

                # 識別済みの話者情報を反映
                display_df = edited_df.copy()
                if not st.session_state.video_combined_identified_df.empty:
                    if 'speaker' in st.session_state.video_combined_identified_df.columns:
                        display_df['speaker'] = st.session_state.video_combined_identified_df['speaker']

                if len(display_df) == 0:
                    st.warning("文字起こし結果がありません。まず文字起こしを実行してください。")
                else:
                    row_labels = _create_row_labels(display_df)

                    col1, col2, col3 = st.columns([2, 2, 1])
                    with col1:
                        start_row = st.selectbox(
                            "開始行を選択",
                            options=range(len(display_df)),
                            format_func=lambda x: row_labels[x],
                            key="video_combined_start_row"
                        )
                    with col2:
                        end_row = st.selectbox(
                            "終了行を選択",
                            options=range(len(display_df)),
                            format_func=lambda x: row_labels[x],
                            index=len(display_df)-1 if len(display_df) > 0 else 0,
                            key="video_combined_end_row"
                        )
                    with col3:
                        similarity_threshold = st.number_input(
                            "類似度閾値",
                            min_value=0.0,
                            max_value=1.0,
                            value=0.7,
                            step=0.01,
                            key="video_combined_similarity_threshold"
                        )

                    if start_row > end_row:
                        st.error("開始行は終了行以前を選択してください")
                    else:
                        # 話者埋め込みファイルのアップロード
                        uploaded_embeddings = st.file_uploader(
                            "話者埋め込みファイルを選択してください（複数選択可）",
                            type=["npy"],
                            accept_multiple_files=True,
                            key="video_combined_embeddings"
                        )

                        if st.button("話者識別を実行", key="video_combined_identify"):
                            if not uploaded_embeddings:
                                st.error("話者埋め込みファイルをアップロードしてください")
                            else:
                                with st.spinner("話者識別中..."):
                                    try:
                                        # 選択範囲のデータフレームを抽出
                                        df_to_identify = edited_df.iloc[start_row:end_row+1].copy()

                                        # 音声ファイルのシーク位置をリセット
                                        st.session_state.video_combined_audio_io.seek(0)

                                        # 話者識別実行
                                        identified_df = identify_speakers_in_dataframe(
                                            st.session_state.video_combined_audio_io,
                                            df_to_identify,
                                            uploaded_embeddings,
                                            similarity_threshold
                                        )

                                        # 選択範囲のみ更新
                                        full_identified_df = edited_df.copy()
                                        for col in identified_df.columns:
                                            full_identified_df.loc[start_row:end_row, col] = identified_df[col].values

                                        st.session_state.video_combined_identified_df = full_identified_df
                                        st.session_state.video_combined_df = full_identified_df
                                        st.session_state.video_combined_step = 5
                                        st.success("話者識別が完了しました")
                                        st.rerun()

                                    except Exception as e:
                                        st.error(f"話者識別中にエラーが発生しました: {e}")
                                        import traceback
                                        st.error("詳細なエラー情報:")
                                        st.code(traceback.format_exc())

            with tab2:
                st.write("音声セグメントを選択し、プレビュー確認してから話者埋め込みファイルを作成できます。")

                if len(edited_df) == 0:
                    st.warning("文字起こし結果がありません。まず文字起こしを実行してください。")
                else:
                    # === 選択方式の切り替え ===
                    selection_mode = st.radio(
                        "選択方式を選んでください",
                        options=["範囲指定モード", "チェックボックスモード"],
                        horizontal=True,
                        help="範囲指定: 連続した行を素早く選択 | チェックボックス: 飛び飛びの行を自由に選択"
                    )

                    st.divider()

                    # === Part 1: 選択UI（方式に応じて変更） ===
                    st.subheader("1️⃣ 音声セグメントを選択")

                    if selection_mode == "範囲指定モード":
                        # 方式A: selectboxによる範囲指定
                        embedding_row_labels = _create_row_labels(edited_df)

                        col1, col2 = st.columns(2)
                        with col1:
                            embedding_start_row = st.selectbox(
                                "開始行を選択",
                                options=range(len(edited_df)),
                                format_func=lambda x: embedding_row_labels[x],
                                key="video_combined_embedding_start_row"
                            )
                        with col2:
                            embedding_end_row = st.selectbox(
                                "終了行を選択",
                                options=range(len(edited_df)),
                                format_func=lambda x: embedding_row_labels[x],
                                index=0,  # デフォルト: 0行目
                                key="video_combined_embedding_end_row"
                            )

                        # 範囲指定の検証とセッション状態への反映
                        if embedding_start_row > embedding_end_row:
                            st.error("⚠️ 開始行は終了行以前を選択してください")
                            st.session_state.embedding_selected_rows = set()
                        else:
                            # 範囲をsetに変換
                            st.session_state.embedding_selected_rows = set(range(embedding_start_row, embedding_end_row + 1))
                            st.success(f"✅ 行 {embedding_start_row} ～ {embedding_end_row} を選択しました")

                    else:
                        # 方式B: data_editorによるチェックボックス選択
                        display_df = edited_df.copy()
                        display_df.insert(0, "選択", False)

                        # セッションに保存されている選択行を反映
                        if st.session_state.embedding_selected_rows:
                            for idx in st.session_state.embedding_selected_rows:
                                if idx < len(display_df):
                                    display_df.at[idx, "選択"] = True

                        edited_display = st.data_editor(
                            display_df,
                            column_config={
                                "選択": st.column_config.CheckboxColumn(
                                    "選択",
                                    help="埋め込み作成対象の行をチェック",
                                    default=False
                                )
                            },
                            disabled=["start", "end", "speaker", "text"],
                            use_container_width=True,
                            hide_index=True,
                            key="embedding_data_editor"
                        )

                        # 選択状態を更新
                        st.session_state.embedding_selected_rows = set(
                            edited_display[edited_display["選択"] == True].index.tolist()
                        )

                    # === Part 2以降: 共通処理（選択がある場合のみ表示） ===
                    if st.session_state.embedding_selected_rows:
                        st.subheader("2️⃣ 選択情報の確認")

                        selection_summary = get_selection_summary(edited_df, st.session_state.embedding_selected_rows)

                        # 情報表示
                        info_cols = st.columns(4)
                        with info_cols[0]:
                            st.metric("選択行数", selection_summary['count'])
                        with info_cols[1]:
                            st.metric("開始時刻", format_time(selection_summary['start_time']))
                        with info_cols[2]:
                            st.metric("終了時刻", format_time(selection_summary['end_time']))
                        with info_cols[3]:
                            st.metric("音声長", f"{selection_summary['duration']:.1f}秒")

                        # 話者情報の表示
                        if selection_summary['speakers']:
                            speakers_text = "、".join([s if s else "不明" for s in selection_summary['speakers']])
                            st.info(f"🎤 選択範囲に含まれる話者: {speakers_text}")

                        # === Part 3: 音声プレビュー ===
                        st.subheader("3️⃣ 音声プレビュー")

                        col_preview, col_clear = st.columns([3, 1])

                        with col_preview:
                            if st.button("🔊 プレビュー音声を生成", key="generate_preview"):
                                try:
                                    with st.spinner("プレビュー音声を生成中..."):
                                        preview_bytes, duration = prepare_embedding_preview_audio(
                                            st.session_state.video_combined_audio_io,
                                            edited_df,
                                            list(st.session_state.embedding_selected_rows)
                                        )
                                        st.session_state.embedding_preview_audio = preview_bytes
                                        st.session_state.embedding_show_preview = True
                                    st.success(f"プレビュー音声を生成しました（{duration:.1f}秒）")
                                except Exception as e:
                                    st.error(f"プレビュー生成エラー: {e}")

                        with col_clear:
                            if st.session_state.embedding_show_preview:
                                if st.button("❌ プレビューをクリア", key="clear_preview"):
                                    st.session_state.embedding_preview_audio = None
                                    st.session_state.embedding_show_preview = False
                                    st.rerun()

                        # プレビュー再生
                        if st.session_state.embedding_show_preview and st.session_state.embedding_preview_audio:
                            st.audio(st.session_state.embedding_preview_audio, format="audio/wav")

                        # === Part 4: ファイル名設定 ===
                        st.subheader("4️⃣ ファイル名設定")

                        embedding_filename = st.text_input(
                            "ファイル名（.npy拡張子は自動追加）",
                            value="speaker_embedding",
                            key="video_combined_embedding_filename",
                            help="作成する話者埋め込みファイルの名前を指定してください"
                        )

                        # === Part 5: ワンクリック埋め込み作成 ===
                        st.subheader("5️⃣ 埋め込み作成")

                        if st.button("✨ 話者埋め込みを作成してダウンロード", key="create_embedding_oneclick", type="primary"):
                            with st.spinner("話者埋め込みを作成中..."):
                                try:
                                    selected_indices = sorted(list(st.session_state.embedding_selected_rows))

                                    # 埋め込みベクトルを生成
                                    embedding, duration = extract_audio_segment_for_embedding(
                                        st.session_state.video_combined_audio_io,
                                        edited_df,
                                        selected_indices
                                    )

                                    # ファイル名処理
                                    filename_with_ext = embedding_filename if embedding_filename.endswith('.npy') else f"{embedding_filename}.npy"

                                    # 埋め込みをバイト列に変換
                                    embedding_io = BytesIO()
                                    np.save(embedding_io, embedding)
                                    embedding_io.seek(0)

                                    embedding_bytes = embedding_io.getvalue()

                                    st.success(f"✅ 話者埋め込みの作成が完了しました（音声長: {duration:.1f}秒）")
                                    trigger_auto_download(
                                        embedding_bytes,
                                        filename_with_ext,
                                        key="video_combined_download_embedding",
                                        mime="application/octet-stream"
                                    )

                                except Exception as e:
                                    st.error(f"❌ 話者埋め込みの作成中にエラーが発生しました: {e}")
                                    import traceback
                                    with st.expander("🔍 詳細なエラー情報"):
                                        st.code(traceback.format_exc())

                    else:
                        if selection_mode == "範囲指定モード":
                            st.info("💡 上記のselectboxで開始行と終了行を選択してください")
                        else:
                            st.info("💡 データエディタで埋め込み作成対象の行をチェックしてください")

        # Step 5: 結果のダウンロード
        if st.session_state.video_combined_step >= 5 and not st.session_state.video_combined_identified_df.empty:
            st.subheader("Step 5: 結果のダウンロード")

            # 議事録形式のテキスト生成（連続する同じ話者の発言を結合）
            df_for_transcript = st.session_state.video_combined_identified_df.copy()

            # 話者列の前処理: 空欄を前後の話者で埋める
            df_for_transcript['speaker_filled'] = df_for_transcript['speaker'].replace('', pd.NA)
            df_for_transcript['speaker_filled'] = df_for_transcript['speaker_filled'].ffill()

            # 話者が変わるごとに新しいグループIDを割り当てる
            df_for_transcript['group_id'] = (df_for_transcript['speaker_filled'] != df_for_transcript['speaker_filled'].shift()).cumsum()

            # グループごとにテキストを結合
            df_merged = df_for_transcript.groupby('group_id').agg(
                speaker=('speaker_filled', 'first'),
                text=('text', ' '.join)
            ).reset_index(drop=True)

            # 議事録テキスト生成
            transcript_lines = []
            for idx, row in df_merged.iterrows():
                speaker = row.get('speaker', '')
                text = row.get('text', '')

                # Format: （話者）テキスト
                speaker_str = speaker if speaker else '不明'
                transcript_lines.append(f"（{speaker_str}）{text}")

            transcript_text = "\n".join(transcript_lines)

            # Word文書として生成
            doc = DocxDocument()
            doc.add_heading('議事録', 0)

            for line in transcript_text.split('\n'):
                if line.strip():
                    doc.add_paragraph(line)
                else:
                    doc.add_paragraph('')  # 空行を保持

            docx_buffer = BytesIO()
            doc.save(docx_buffer)
            docx_buffer.seek(0)

            col1, col2 = st.columns(2)
            with col1:
                st.download_button(
                    label="議事録テキストをダウンロード",
                    data=docx_buffer.getvalue(),
                    file_name=f"transcript_{st.session_state.video_combined_uploaded_file_name.split('.')[0]}.docx",
                    mime='application/vnd.openxmlformats-officedocument.wordprocessingml.document',
                    key="video_combined_download_transcript"
                )

            with col2:
                # Excelファイルとして生成
                excel_buffer = BytesIO()
                with pd.ExcelWriter(excel_buffer, engine='openpyxl') as writer:
                    st.session_state.video_combined_identified_df.to_excel(writer, index=False, sheet_name='文字起こし')
                excel_buffer.seek(0)

                st.download_button(
                    label="Excelファイルをダウンロード",
                    data=excel_buffer.getvalue(),
                    file_name=f"transcript_{st.session_state.video_combined_uploaded_file_name.split('.')[0]}.xlsx",
                    mime='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
                    key="video_combined_download_excel"
                )

def main():
    st.set_page_config(
        page_title="文字起こしアプリ",
        layout="wide",
        page_icon="🎙️",
        initial_sidebar_state="expanded"
    )

    # カスタムCSSスタイル
    st.markdown("""
        <style>
        /* メインコンテナのスタイリング */
        .main {
            background-color: #f8f9fa;
        }

        /* コンテンツエリアの背景 */
        .block-container {
            background: white;
            border-radius: 8px;
            padding: 2rem 3rem;
            margin-top: 1rem;
            box-shadow: 0 1px 3px rgba(0, 0, 0, 0.08);
        }

        /* サイドバーのスタイリング */
        section[data-testid="stSidebar"] {
            background-color: #f8f9fa;
            border-right: 1px solid #e9ecef;
        }

        /* ボタンのスタイリング */
        .stButton > button {
            background-color: #4a5568;
            color: white;
            border: none;
            border-radius: 6px;
            padding: 0.5rem 1.5rem;
            font-weight: 500;
            font-size: 0.95rem;
            transition: background-color 0.2s ease;
        }

        .stButton > button:hover {
            background-color: #2d3748;
        }

        /* プライマリボタン */
        .stButton > button[kind="primary"] {
            background-color: #3182ce;
        }

        .stButton > button[kind="primary"]:hover {
            background-color: #2c5282;
        }

        /* タイトルのスタイリング */
        h1 {
            color: #1a202c;
            font-weight: 700;
            font-size: 2rem;
            padding: 0.5rem 0;
            margin-bottom: 1rem;
            border-bottom: 2px solid #e2e8f0;
        }

        /* サブヘッダーのスタイリング */
        h2, h3 {
            color: #2d3748;
            font-weight: 600;
            margin-top: 1.5rem;
            padding-bottom: 0.5rem;
            border-bottom: 1px solid #e2e8f0;
        }

        /* 入力フィールドのスタイリング */
        .stTextInput > div > div > input,
        .stNumberInput > div > div > input,
        .stTextArea > div > div > textarea {
            border-radius: 6px;
            border: 1px solid #cbd5e0;
            padding: 0.5rem;
            transition: border-color 0.2s ease;
        }

        .stTextInput > div > div > input:focus,
        .stNumberInput > div > div > input:focus,
        .stTextArea > div > div > textarea:focus {
            border-color: #3182ce;
            box-shadow: 0 0 0 3px rgba(49, 130, 206, 0.1);
        }

        /* セレクトボックスのスタイリング */
        .stSelectbox > div > div {
            border-radius: 6px;
            border: 1px solid #cbd5e0;
        }

        /* ファイルアップローダーのスタイリング */
        .stFileUploader {
            background-color: #f7fafc;
            border-radius: 8px;
            padding: 1.5rem;
            border: 2px dashed #cbd5e0;
        }

        /* データフレームのスタイリング */
        .stDataFrame {
            border-radius: 6px;
            overflow: hidden;
            box-shadow: 0 1px 3px rgba(0, 0, 0, 0.08);
        }

        /* ダウンロードボタン */
        .stDownloadButton > button {
            background-color: #38a169;
            color: white;
            border-radius: 6px;
            font-weight: 500;
            transition: background-color 0.2s ease;
        }

        .stDownloadButton > button:hover {
            background-color: #2f855a;
        }

        /* プログレスバー */
        .stProgress > div > div > div {
            background-color: #3182ce;
        }

        /* カラムの間隔調整 */
        [data-testid="column"] {
            padding: 0.5rem;
        }

        /* エクスパンダー */
        .streamlit-expanderHeader {
            background-color: #f7fafc;
            border-radius: 6px;
            font-weight: 500;
            color: #2d3748;
        }

        /* スライダー */
        .stSlider > div > div > div {
            background-color: #3182ce;
        }

        /* タブ */
        .stTabs [data-baseweb="tab-list"] {
            gap: 4px;
        }

        .stTabs [data-baseweb="tab"] {
            background-color: #f7fafc;
            border-radius: 6px 6px 0 0;
            padding: 8px 16px;
            font-weight: 500;
            color: #4a5568;
        }

        .stTabs [aria-selected="true"] {
            background-color: #3182ce;
            color: white;
        }
        </style>
    """, unsafe_allow_html=True)

    st.sidebar.title("メニュー")
    menu_options = [
        "文字起こし",
        "ナレッジベース管理",
        "議事録校正（RAG）",
        "🚀 一括処理パイプライン",
        "動画から音声を切り出しMP3で保存"
    ]
    choice = st.sidebar.selectbox("機能を選択してください", menu_options)

    if choice == "文字起こし":
        video_transcribe_and_identify()
    elif choice == "ナレッジベース管理":
        knowledge_base_management()
    elif choice == "議事録校正（RAG）":
        proofread_meeting_minutes()
    elif choice == "🚀 一括処理パイプライン":
        batch_processing_pipeline()
    elif choice == "動画から音声を切り出しMP3で保存":
        video_to_audio_cutter_app()

if __name__ == "__main__":
    main()
