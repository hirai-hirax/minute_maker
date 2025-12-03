import os
os.environ["KMP_DUPLICATE_LIB_OK"] = "TRUE"
import streamlit as st
from openai import AzureOpenAI
import tempfile
from io import BytesIO
from pydub import AudioSegment
import fitz
import torch
import pandas as pd
from dotenv import load_dotenv
import torchaudio
from speechbrain.pretrained import EncoderClassifier
import numpy as np
import zipfile
from datetime import timedelta, datetime
from concurrent.futures import ThreadPoolExecutor, as_completed
import threading
import json
import importlib
import base64
import uuid
import subprocess
from docx import Document as DocxDocument
from pptx import Presentation
from pathlib import Path
from streamlit.components.v1 import html as components_html

torch.classes.__path__ = []

# .envファイルから環境変数を読み込む
load_dotenv()

# 環境変数から設定を取得（Azure OpenAI のエンドポイント・API キーを設定してください）
AZURE_OPENAI_ENDPOINT = os.environ.get("AZURE_OPENAI_ENDPOINT")
AZURE_OPENAI_API_KEY = os.environ.get("AZURE_OPENAI_API_KEY")
API_VERSION = "2025-03-01-preview"  # gpt-4o-transcribe モデルに対応したバージョン

# RAGデータベースの保存先フォルダ設定
DEFAULT_RAGDB_FOLDER = ""  # 空文字列の場合はカレントディレクトリ（デフォルト）
# 例: DEFAULT_RAGDB_FOLDER = "C:/Users/username/Documents/ragdb"
# 例: DEFAULT_RAGDB_FOLDER = "./data/ragdb"

DEFAULT_MEETING_TYPES = [
    {
        "name": "経営会議",
        "id": "executive_meeting",
        "embeddings_folder": "speaker_embeddings/executive",
        "description": "月例経営会議（役員メンバー固定）",
    },
    {
        "name": "開発チーム定例",
        "id": "dev_team_meeting",
        "embeddings_folder": "speaker_embeddings/dev_team",
        "description": "週次開発チーム会議",
    },
    {
        "name": "営業定例",
        "id": "sales_meeting",
        "embeddings_folder": "speaker_embeddings/sales",
        "description": "月次営業会議",
    },
    {
        "name": "カスタム（手動選択）",
        "id": "custom",
        "embeddings_folder": "",
        "description": "会議タイプを指定せず、手動で話者埋め込みを選択",
    },
]


def _extract_pdf(file: BytesIO) -> str:
    file_bytes = file.read()
    file.seek(0)
    pdf_document = fitz.open(stream=file_bytes, filetype='pdf')
    text = "".join(page.get_text() for page in pdf_document)
    pdf_document.close()
    return text

def _extract_docx(file: BytesIO) -> str:
    doc = DocxDocument(file)
    text = "\n".join(p.text for p in doc.paragraphs)
    for table in doc.tables:
        text += "\n" + "\n".join(" ".join(cell.text for cell in row.cells) for row in table.rows)
    return text.strip()

def _extract_pptx(file: BytesIO) -> str:
    presentation = Presentation(file)
    text_parts = []
    for slide_num, slide in enumerate(presentation.slides, 1):
        text_parts.append(f"\n--- スライド {slide_num} ---\n")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text_parts.append(shape.text + "\n")
            if shape.has_table:
                for row in shape.table.rows:
                    text_parts.append(" ".join(cell.text for cell in row.cells) + "\n")
    return "".join(text_parts).strip()

def _extract_msg(file: BytesIO) -> str:
    try:
        import extract_msg
    except ImportError:
        raise Exception("MSGファイルの処理には extract-msg ライブラリが必要です。pip install extract-msg を実行してください。")

    with tempfile.NamedTemporaryFile(delete=False, suffix=".msg") as temp_file:
        temp_file.write(file.read())
        temp_path = temp_file.name

    try:
        msg = extract_msg.Message(temp_path)
        parts = []
        if msg.subject:
            parts.append(f"件名: {msg.subject}\n")
        if msg.sender:
            parts.append(f"送信者: {msg.sender}\n")
        if msg.to:
            parts.append(f"宛先: {msg.to}\n")
        if msg.body:
            parts.append(f"\n{msg.body}")
        msg.close()
        return "\n".join(parts).strip()
    finally:
        if os.path.exists(temp_path):
            os.unlink(temp_path)

def _extract_txt(file: BytesIO) -> str:
    file_content = file.read()
    for encoding in ['utf-8', 'cp932', 'shift_jis', 'utf-16', 'iso-2022-jp']:
        try:
            return file_content.decode(encoding).strip()
        except UnicodeDecodeError:
            continue
    raise Exception("テキストファイルのエンコーディングを判定できませんでした")

FILE_EXTRACTORS = {
    'pdf': _extract_pdf,
    'docx': _extract_docx,
    'pptx': _extract_pptx,
    'txt': _extract_txt,
    'msg': _extract_msg,
}

def extract_text_from_file(file, file_extension):
    """ファイル形式に応じてテキストを抽出"""
    extractor = FILE_EXTRACTORS.get(file_extension.lower())
    if not extractor:
        raise Exception(f"サポートされていないファイル形式: {file_extension}")
    try:
        return extractor(file)
    except Exception as e:
        raise Exception(f"{file_extension.upper()}読み込みエラー: {e}")

def _create_azure_client():
    """Azure OpenAIクライアントを作成"""
    return AzureOpenAI(
        azure_endpoint=AZURE_OPENAI_ENDPOINT,
        api_key=AZURE_OPENAI_API_KEY,
        api_version=API_VERSION,
    )

def generate_summary(model, prompt, text):
    """テキストを要約"""
    client = _create_azure_client()
    response = client.chat.completions.create(
        model=model,
        messages=[
            {"role": "system", "content": prompt},
            {"role": "user", "content": text},
        ],
    )
    print(f"Response: {response}")
    return response.choices[0].message.content

from contextlib import contextmanager

@contextmanager
def temp_file_path(data: bytes, suffix: str):
    """一時ファイルを作成し、パスを返すコンテキストマネージャ"""
    with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
        tmp.write(data)
        tmp_path = tmp.name
    try:
        yield tmp_path
    finally:
        if os.path.exists(tmp_path):
            os.unlink(tmp_path)


def trigger_auto_download(data: bytes, file_name: str, key: str | None, mime: str = "application/octet-stream"):
    """Streamlitコンポーネントを使って即時ダウンロードを開始"""
    if not data:
        return

    encoded = base64.b64encode(data).decode()
    mime_js = json.dumps(mime)
    file_name_js = json.dumps(file_name)
    element_id_js = json.dumps(key or f"dl_{uuid.uuid4().hex}")

    components_html(
        f"""
        <div id={element_id_js}></div>
        <script>
            (function() {{
                const mimeType = {mime_js};
                const fileName = {file_name_js};
                const link = document.createElement('a');
                link.href = "data:" + mimeType + ";base64,{encoded}";
                link.download = fileName;
                link.style.display = 'none';
                document.body.appendChild(link);
                link.click();
                document.body.removeChild(link);
                const element = document.getElementById({element_id_js});
                if (element) {{
                    element.remove();
                }}
            }})();
        </script>
        """,
        height=0,
    )


def transcribe_audio_to_dataframe(uploaded_file: BytesIO, reference_file: BytesIO = None, model: str = "gpt-4o-transcribe-diarize"):
    """音声ファイル全体を文字起こし（モデル選択可能、25MB制限対応）

    Args:
        uploaded_file: 音声ファイル
        reference_file: 参考資料（Whisperのみサポート）
        model: 使用するモデル ("gpt-4o-transcribe-diarize" または "whisper")
    """
    # ファイルサイズをチェック（25MB = 26,214,400 bytes）
    MAX_FILE_SIZE = 25 * 1024 * 1024
    uploaded_file.seek(0)
    file_size = len(uploaded_file.getvalue())
    uploaded_file.seek(0)

    # 25MB以上の場合は分割処理
    if file_size > MAX_FILE_SIZE:
        st.warning(f"⚠️ ファイルサイズが {file_size / (1024*1024):.1f}MB です。25MB制限のため、音声を分割して処理します。")
        return _transcribe_large_audio_chunked(uploaded_file, reference_file, model)

    # 25MB未満の場合は通常処理
    return _transcribe_audio_single(uploaded_file, reference_file, model)

def _transcribe_large_audio_chunked(uploaded_file: BytesIO, reference_file: BytesIO = None, model: str = "gpt-4o-transcribe-diarize"):
    """大きな音声ファイルを分割して文字起こし

    Args:
        uploaded_file: 音声ファイル
        reference_file: 参考資料
        model: 使用するモデル

    Returns:
        pd.DataFrame: 文字起こし結果
    """
    try:
        suffix = f".{uploaded_file.name.split('.')[-1]}"

        # 音声ファイルを読み込み
        with temp_file_path(uploaded_file.getvalue(), suffix) as tmp_path:
            audio = AudioSegment.from_file(tmp_path)

        # 音声の長さ（ミリ秒）
        total_duration_ms = len(audio)
        total_duration_sec = total_duration_ms / 1000

        # チャンクサイズを決定（10分 = 600秒 = 600,000ミリ秒）
        CHUNK_DURATION_MS = 10 * 60 * 1000  # 10分

        # チャンク数を計算
        num_chunks = (total_duration_ms + CHUNK_DURATION_MS - 1) // CHUNK_DURATION_MS

        st.info(f"📊 音声長: {total_duration_sec/60:.1f}分、{num_chunks}個のチャンクに分割して処理します")

        # 進捗表示
        progress_bar = st.progress(0)
        status_text = st.empty()

        all_segments = []

        for i in range(num_chunks):
            start_ms = i * CHUNK_DURATION_MS
            end_ms = min((i + 1) * CHUNK_DURATION_MS, total_duration_ms)

            status_text.write(f"🎤 チャンク {i+1}/{num_chunks} を処理中... ({start_ms/1000:.1f}秒 ～ {end_ms/1000:.1f}秒)")

            # チャンクを抽出
            chunk = audio[start_ms:end_ms]

            # 一時ファイルとして保存
            chunk_io = BytesIO()
            chunk.export(chunk_io, format="mp3", bitrate="192k")
            chunk_io.seek(0)
            chunk_io.name = f"chunk_{i}.mp3"

            # チャンクを文字起こし
            try:
                chunk_df = _transcribe_audio_single(chunk_io, reference_file, model)

                # タイムスタンプをオフセット調整
                if not chunk_df.empty and 'start' in chunk_df.columns and 'end' in chunk_df.columns:
                    offset_sec = start_ms / 1000
                    chunk_df['start'] = chunk_df['start'] + offset_sec
                    chunk_df['end'] = chunk_df['end'] + offset_sec

                all_segments.append(chunk_df)

            except Exception as e:
                st.error(f"❌ チャンク {i+1} の処理中にエラー: {e}")

            progress_bar.progress((i + 1) / num_chunks)

        # すべてのチャンクを結合
        if all_segments:
            result_df = pd.concat(all_segments, ignore_index=True)
            status_text.write(f"✅ 分割処理完了！合計 {len(result_df)} 個のセグメント")
            return result_df
        else:
            st.error("❌ すべてのチャンクの処理に失敗しました")
            return pd.DataFrame(columns=["start", "end", "speaker", "text"])

    except Exception as e:
        st.error(f"❌ 分割処理中にエラーが発生しました: {e}")
        import traceback
        st.error(traceback.format_exc())
        return pd.DataFrame(columns=["start", "end", "speaker", "text"])

def _transcribe_audio_single(uploaded_file: BytesIO, reference_file: BytesIO = None, model: str = "gpt-4o-transcribe-diarize"):
    """単一の音声ファイルを文字起こし（内部用）

    Args:
        uploaded_file: 音声ファイル
        reference_file: 参考資料（Whisperのみサポート）
        model: 使用するモデル ("gpt-4o-transcribe-diarize" または "whisper")
    """
    # モデルに応じたAPIバージョンを選択
    if model == "whisper":
        api_version = "2024-06-01"  # Whisper用の安定版APIバージョン
        st.info(f"🔧 Whisperモデル用にAPIバージョン {api_version} を使用します")
    else:
        api_version = "2025-03-01-preview"  # gpt-4o-transcribe-diarize用

    # モデル専用のクライアントを作成
    client = AzureOpenAI(
        azure_endpoint=AZURE_OPENAI_ENDPOINT,
        api_key=AZURE_OPENAI_API_KEY,
        api_version=api_version,
    )

    try:
        suffix = f".{uploaded_file.name.split('.')[-1]}"

        if model == "gpt-4o-transcribe-diarize":
            # gpt-4o-transcribe-diarizeモデル（話者識別付き）
            if reference_file:
                st.warning("⚠️ gpt-4o-transcribe-diarizeモデルは参考資料（prompt）をサポートしていません。参考資料は無視されます。")

            st.info("文字起こしを開始します（gpt-4o-transcribe-diarize、話者識別付き）...")

            with temp_file_path(uploaded_file.getvalue(), suffix) as tmp_path:
                with open(tmp_path, "rb") as audio_file:
                    transcript = client.audio.transcriptions.create(
                        model="gpt-4o-transcribe-diarize",
                        file=(uploaded_file.name, audio_file, f"audio/{uploaded_file.name.split('.')[-1]}"),
                        response_format="diarized_json",
                        chunking_strategy="auto"
                    )

            # レスポンスからセグメントを取得
            transcript_dict = transcript.model_dump()
            segments = transcript_dict.get("segments", [])

            if segments:
                # セグメントをデータフレームに変換
                seg_list = []
                for seg in segments:
                    seg_list.append({
                        "start": seg.get("start", 0),
                        "end": seg.get("end", 0),
                        "speaker": seg.get("speaker", ""),
                        "text": seg.get("text", "")
                    })

                st.success(f"文字起こしが完了しました！（{len(seg_list)}個のセグメント、話者識別付き）")
                seg_df = pd.DataFrame(seg_list)
                return seg_df
            else:
                # segmentsがない場合はテキスト全体を取得
                text = transcript_dict.get("text", "")
                if text:
                    st.warning("⚠️ セグメント情報が取得できませんでした。全体を1つのセグメントとして扱います。")
                    seg_df = pd.DataFrame([{
                        "start": 0,
                        "end": 0,
                        "speaker": "",
                        "text": text
                    }])
                    return seg_df
                else:
                    st.error("❌ 文字起こし結果が空でした")
                    return pd.DataFrame(columns=["start", "end", "speaker", "text"])

        elif model == "whisper":
            # Whisperモデル（話者識別なし）
            st.info("文字起こしを開始します（Whisper、話者識別なし）...")

            with temp_file_path(uploaded_file.getvalue(), suffix) as tmp_path:
                with open(tmp_path, "rb") as audio_file:
                    # promptパラメータの準備
                    kwargs = {
                        "model": "whisper",
                        "file": (uploaded_file.name, audio_file, f"audio/{uploaded_file.name.split('.')[-1]}"),
                        "response_format": "verbose_json"
                    }

                    # 参考資料がある場合はpromptとして使用
                    if reference_file:
                        try:
                            file_extension = reference_file.name.split(".")[-1].lower()
                            reference_text = extract_text_from_file(BytesIO(reference_file.read()), file_extension)
                            if reference_text:
                                # promptは最大224トークン程度に制限（約1000文字）
                                kwargs["prompt"] = reference_text[:1000]
                                st.info("✅ 参考資料をpromptとして使用します")
                        except Exception as e:
                            st.warning(f"⚠️ 参考資料の読み込みに失敗しました: {e}")

                    transcript = client.audio.transcriptions.create(**kwargs)

            # レスポンスからセグメントを取得
            transcript_dict = transcript.model_dump()
            segments = transcript_dict.get("segments", [])

            if segments:
                # セグメントをデータフレームに変換（話者情報は空）
                seg_list = []
                for seg in segments:
                    seg_list.append({
                        "start": seg.get("start", 0),
                        "end": seg.get("end", 0),
                        "speaker": "",  # Whisperは話者識別をサポートしていない
                        "text": seg.get("text", "")
                    })

                st.success(f"文字起こしが完了しました！（{len(seg_list)}個のセグメント、Whisper使用）")
                seg_df = pd.DataFrame(seg_list)
                return seg_df
            else:
                # segmentsがない場合はテキスト全体を取得
                text = transcript_dict.get("text", "")
                if text:
                    st.warning("⚠️ セグメント情報が取得できませんでした。全体を1つのセグメントとして扱います。")
                    seg_df = pd.DataFrame([{
                        "start": 0,
                        "end": 0,
                        "speaker": "",
                        "text": text
                    }])
                    return seg_df
                else:
                    st.error("❌ 文字起こし結果が空でした")
                    return pd.DataFrame(columns=["start", "end", "speaker", "text"])
        else:
            st.error(f"❌ サポートされていないモデル: {model}")
            return pd.DataFrame(columns=["start", "end", "speaker", "text"])

    except Exception as e:
        st.error(f"❌ 文字起こし中にエラーが発生しました")
        st.error(f"**エラー詳細**: {str(e)}")
        st.error(f"**使用モデル**: {model}")
        st.error(f"**APIバージョン**: {api_version}")

        import traceback
        error_details = traceback.format_exc()

        with st.expander("🔍 詳細なエラー情報（デバッグ用）"):
            st.code(error_details, language="python")

        # エラーをコンソールにも出力（デバッグ用）
        print(f"===== Transcription Error =====")
        print(f"Model: {model}")
        print(f"API Version: {api_version}")
        print(f"Error: {e}")
        print(error_details)
        print(f"==============================")

        return pd.DataFrame(columns=["start", "end", "speaker", "text"])

@st.cache_resource
def load_speaker_encoder():
    """Caches the SpeechBrain speaker encoder model."""
    return EncoderClassifier.from_hparams(
        source="speechbrain/spkrec-ecapa-voxceleb",
        run_opts={"device": "cpu"}
    )

def _compute_embedding_from_wav_bytes(wav_bytes: bytes) -> np.ndarray:
    """Compute a speaker embedding from WAV bytes using SpeechBrain."""
    with temp_file_path(wav_bytes, ".wav") as wav_path:
        waveform, sample_rate = torchaudio.load(wav_path)

    if waveform.shape[0] > 1:
        waveform = waveform.mean(dim=0, keepdim=True)

    target_sr = 16000
    if sample_rate != target_sr:
        waveform = torchaudio.transforms.Resample(orig_freq=sample_rate, new_freq=target_sr)(waveform)

    encoder = load_speaker_encoder()
    waveform = waveform.to(dtype=torch.float32)

    with torch.no_grad():
        embedding = encoder.encode_batch(waveform)

    return embedding.squeeze().cpu().numpy()

def extract_embedding(audio_content):
    """Extracts embedding from audio content using SpeechBrain."""
    audio_bytes = audio_content.read()
    audio_content.seek(0)

    try:
        audio_segment = AudioSegment.from_file(BytesIO(audio_bytes))
        wav_bytes = audio_segment.export(format="wav").read()
    except Exception:
        wav_bytes = audio_bytes

    return _compute_embedding_from_wav_bytes(wav_bytes)

def load_speaker_embeddings_from_files(uploaded_files):
    """Loads known speaker embeddings from uploaded files.

    Args:
        uploaded_files: リストまたはタプル。各要素は以下のいずれか：
            - Streamlit UploadedFile オブジェクト（file_uploaderから）
            - 辞書 {'name': filename, 'data': file_bytes}（フォルダから読み込んだ場合）

    Returns:
        dict: {speaker_name: embedding_array}
    """
    if not uploaded_files:
        st.warning("話者埋め込みファイルがアップロードされていません。")
        return {}

    speaker_embeddings = {}
    for uploaded_file in uploaded_files:
        try:
            # 辞書形式（フォルダから読み込んだ場合）かUploadedFileオブジェクトかを判定
            if isinstance(uploaded_file, dict):
                # 辞書形式: {'name': filename, 'data': file_bytes}
                filename = uploaded_file['name']
                file_data = BytesIO(uploaded_file['data'])
            else:
                # UploadedFile オブジェクト
                filename = uploaded_file.name
                file_data = uploaded_file

            filename_stem = Path(filename).stem
            # ファイル名の区切り文字で話者名を抽出
            for delimiter in ['‗', '_']:
                if delimiter in filename_stem:
                    filename_stem = filename_stem.split(delimiter)[0]
                    break
            speaker_name = filename_stem.strip()

            if not speaker_name:
                speaker_name = Path(filename).stem

            speaker_embeddings[speaker_name] = np.load(file_data)
        except Exception as e:
            filename_str = uploaded_file.get('name') if isinstance(uploaded_file, dict) else getattr(uploaded_file, 'name', 'unknown')
            st.error(f"埋め込みファイルの読み込み中にエラーが発生しました {filename_str}: {e}")
    return speaker_embeddings

def _calculate_similarity(embedding1, embedding2):
    """2つの埋め込みベクトルの類似度を計算"""
    return np.dot(embedding1, embedding2) / (np.linalg.norm(embedding1) * np.linalg.norm(embedding2))

def _identify_speaker(segment_embedding, known_embeddings, threshold):
    """セグメント埋め込みから話者を識別"""
    best_speaker, best_similarity = None, -1
    for speaker_name, known_embedding in known_embeddings.items():
        similarity = _calculate_similarity(segment_embedding, known_embedding)
        print(f"Comparing segment with {speaker_name}: similarity = {similarity:.4f}")
        if similarity > best_similarity:
            best_similarity, best_speaker = similarity, speaker_name
    return best_speaker if best_similarity >= threshold else ""

def identify_speakers_in_dataframe(audio_file, df: pd.DataFrame, uploaded_embedding_files, similarity_threshold: float) -> pd.DataFrame:
    known_embeddings = load_speaker_embeddings_from_files(uploaded_embedding_files)
    if not known_embeddings:
        st.warning("既知の話者埋め込みが見つかりませんでした。識別を実行できません。")
        df['speaker'] = None
        return df

    st.info(f"Loaded embeddings for speakers: {list(known_embeddings.keys())}")

    with temp_file_path(audio_file.getvalue(), ".wav") as audio_path:
        try:
            audio = AudioSegment.from_file(audio_path)
            df['speaker'] = None
            progress_bar, status_text = st.progress(0), st.empty()

            for index, row in df.iterrows():
                segment = audio[row['start'] * 1000:row['end'] * 1000]

                try:
                    segment_wav_bytes = segment.export(format="wav").read()
                    segment_embedding = _compute_embedding_from_wav_bytes(segment_wav_bytes)
                    speaker = _identify_speaker(segment_embedding, known_embeddings, similarity_threshold)
                    df.at[index, 'speaker'] = speaker
                    status = f"Identified as {speaker}" if speaker else "Similarity below threshold"
                    status_text.text(f"Processed segment {index + 1}/{len(df)}: {status}")
                except Exception as e:
                    st.error(f"Error processing segment {row['start']}-{row['end']}s: {e}")
                    df.at[index, 'speaker'] = "Error"

                progress_bar.progress((index + 1) / len(df))

            status_text.text("Speaker identification complete.")
            return df

        except Exception as e:
            st.error(f"Error loading or processing audio file: {e}")
            return df

def build_meeting_text_from_dataframe(df: pd.DataFrame) -> str:
    """Generate combined meeting text in (speaker) utterance format from transcription data."""
    if df is None or df.empty or 'text' not in df.columns:
        return ""

    if 'speaker' in df.columns:
        df_copy = df.copy()
        df_copy['speaker_filled'] = df_copy['speaker'].replace('', pd.NA)
        df_copy['speaker_filled'] = df_copy['speaker_filled'].ffill()
        df_copy['group_id'] = (df_copy['speaker_filled'] != df_copy['speaker_filled'].shift()).cumsum()
        df_merged = df_copy.groupby('group_id').agg(
            speaker=('speaker_filled', 'first'),
            text=('text', lambda values: ' '.join(map(str, values)))
        ).reset_index(drop=True)

        lines = []
        for _, row in df_merged.iterrows():
            speaker = row.get('speaker')
            text = row.get('text', '')
            speaker_str = speaker if isinstance(speaker, str) and speaker else '不明'
            lines.append(f"（{speaker_str}）{text}")
        return "\n".join(lines)

    # speaker列がない場合は不明話者として扱う
    lines = []
    for text in df['text'].astype(str):
        lines.append(f"（不明）{text}")
    return "\n".join(lines)

def format_time(seconds):
    """Formats seconds into HH:MM:SS."""
    td = timedelta(seconds=seconds)
    hours, remainder = divmod(td.seconds, 3600)
    minutes, seconds = divmod(remainder, 60)
    return f"{hours:02}:{minutes:02}:{seconds:02}"

def parse_time_to_seconds(time_str):
    """Converts HH:MM:SS or seconds string to total seconds."""
    if ':' in time_str:
        parts = list(map(int, time_str.split(':')))
        if len(parts) == 3:
            return parts[0] * 3600 + parts[1] * 60 + parts[2]
        elif len(parts) == 2:
            return parts[0] * 60 + parts[1]
        else:
            raise ValueError("Invalid time format. Use HH:MM:SS or MM:SS.")
    else:
        return int(time_str)

def split_text_by_lines(text, n_parts):
    """文字列を改行位置で適切にn分割"""
    if n_parts <= 0:
        raise ValueError("分割数は1以上である必要があります")
    if n_parts == 1:
        return [text]

    lines = text.split('\n')
    total_lines = len(lines)

    if total_lines <= n_parts:
        return lines + [""] * (n_parts - total_lines)

    lines_per_part = total_lines // n_parts
    remainder = total_lines % n_parts

    result, start = [], 0
    for i in range(n_parts):
        count = lines_per_part + (1 if i < remainder else 0)
        result.append('\n'.join(lines[start:start + count]))
        start += count

    return result

class RAGProofreadingSystem:
    """RAG機能削除後のダミー校正システム"""

    DEFAULT_TEMPERATURE = 0.3

    def __init__(self, azure_endpoint, azure_api_key, api_version):
        self.azure_endpoint = azure_endpoint
        self.azure_api_key = azure_api_key
        self.api_version = api_version
        self.client = AzureOpenAI(
            azure_endpoint=azure_endpoint,
            api_key=azure_api_key,
            api_version=api_version
        )
        self.documents = []

    def create_knowledge_base(self, documents_text_list, mode="add", documents_metadata=None):
        st.info("RAG機能は削除されたため、ナレッジベースは利用できません。")
        self.documents = []
        return False

    def retrieve_relevant_context(self, query, search_type="similarity", top_k=None):
        return ""

    def rag_enhanced_proofread(self, text, model="gpt-4o", search_type="similarity", top_k=None, prompt_preset="standard"):
        """RAGなしの簡易校正を実行"""
        try:
            system_prompt = (
                "あなたは議事録校正の専門家です。関連資料を用いずに、入力された議事録を基本的に校正してください。"
                "主に誤字脱字と明確さの改善に集中してください。"
            )
            st.info("RAG機能は削除されたため、外部文脈を参照しない簡易校正を実行します。")
            response = self.client.chat.completions.create(
                model=model,
                messages=[
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": text}
                ],
                temperature=self.DEFAULT_TEMPERATURE,
            )
            return response.choices[0].message.content
        except Exception as e:
            st.error(f"校正実行エラー: {e}")
            return None

    def save_knowledge_base(self, output_path):
        return False, "RAG機能削除により保存できません"

    def load_knowledge_base(self, input_path):
        return False, "RAG機能削除により読み込めません", {}

    def get_database_info(self):
        return {
            "has_data": False,
            "documents_count": 0,
            "is_indexed": False,
            "total_chunks": 0,
            "vector_files": 0,
            "output_files": 0,
            "search_types": []
        }

    def get_chunks_detail(self):
        return []

    def clear_knowledge_base(self):
        self.documents = []

# ========================================
# 共通ヘルパー関数
# ========================================

def _init_rag_system():
    """RAGシステムの初期化（共通処理）"""
    if 'global_rag_system' not in st.session_state:
        st.session_state.global_rag_system = RAGProofreadingSystem(
            azure_endpoint=AZURE_OPENAI_ENDPOINT,
            azure_api_key=AZURE_OPENAI_API_KEY,
            api_version=API_VERSION
        )

    if 'global_db_info' not in st.session_state:
        st.session_state.global_db_info = st.session_state.global_rag_system.get_database_info()

    return st.session_state.global_rag_system

def _render_database_status(db_status, show_output_files=False):
    """データベース状態表示（共通処理）"""
    metrics = [
        ("チャンク数", db_status.get('total_chunks', 0)),
        ("インデックス状態", "✅ 構築済み" if db_status.get('is_indexed', False) else "❌ 未構築")
    ]

    if show_output_files:
        metrics.insert(1, ("出力ファイル数", db_status.get('output_files', 0)))

    for col, (label, value) in zip(st.columns(len(metrics)), metrics):
        with col:
            st.metric(label, value)

def _render_database_operations(rag_system, key_prefix="", show_save=True):
    """データベース操作UI（共通処理）"""
    st.subheader("🔧 データベース操作")

    has_data = st.session_state.global_db_info.get('has_data', False)
    cols = st.columns(3)

    # === 保存グループ ===
    with cols[0]:
        with st.container():
            st.markdown("##### 💾 保存")
            st.markdown("---")

            if show_save and has_data:
                st.caption("ナレッジベースをファイルとして保存")

                # デフォルトDBとして保存
                if st.button("📌 デフォルトDBに保存", key=f"{key_prefix}_save_default_db", use_container_width=True, type="primary"):
                    output_path = get_default_ragdb_path()
                    success, message = rag_system.save_knowledge_base(output_path)
                    if success:
                        st.success(f"✅ 保存完了")
                    else:
                        st.error(f"❌ {message}")

                st.markdown("")  # スペーサー

                # 別名で保存
                with st.expander("📥 別名でダウンロード", expanded=False):
                    custom_filename = st.text_input(
                        "ファイル名",
                        value="custom_knowledge_base",
                        key=f"{key_prefix}_custom_filename",
                        placeholder="ファイル名を入力"
                    )
                    if st.button("作成", key=f"{key_prefix}_save_custom_db", use_container_width=True):
                        # 拡張子を追加
                        filename_with_ext = custom_filename if custom_filename.endswith('.ragdb') else f"{custom_filename}.ragdb"

                        # 一時ファイルに保存
                        with tempfile.NamedTemporaryFile(delete=False, suffix='.ragdb') as tmp_file:
                            tmp_path = tmp_file.name

                        try:
                            success, message = rag_system.save_knowledge_base(tmp_path)
                            if success:
                                # ファイルを読み込んでダウンロードボタン表示
                                with open(tmp_path, 'rb') as f:
                                    db_bytes = f.read()

                                st.success(f"✅ 作成完了")
                                st.download_button(
                                    label="📥 ダウンロード",
                                    data=db_bytes,
                                    file_name=filename_with_ext,
                                    mime="application/octet-stream",
                                    key=f"{key_prefix}_download_custom_db",
                                    use_container_width=True
                                )
                            else:
                                st.error(f"❌ {message}")
                        finally:
                            # 一時ファイルを削除
                            if os.path.exists(tmp_path):
                                os.remove(tmp_path)
            else:
                st.caption("保存するデータがありません")
                st.markdown("<br>", unsafe_allow_html=True)
                st.button("📌 デフォルトDBに保存", key=f"{key_prefix}_save_default_db_disabled", use_container_width=True, disabled=True)

    # === 読み込みグループ ===
    with cols[1]:
        with st.container():
            st.markdown("##### 📂 読み込み")
            st.markdown("---")
            st.caption("保存済みのナレッジベースを読み込み")

            uploaded_db = st.file_uploader(
                "RAGDBファイルを選択",
                type=['ragdb'],
                key=f"{key_prefix}_load_db",
                label_visibility="collapsed"
            )
            if st.button("📂 読み込み実行", key=f"{key_prefix}_load_btn", use_container_width=True, type="primary", disabled=uploaded_db is None):
                with temp_file_path(uploaded_db.getvalue(), '.ragdb') as tmp_path:
                    success, message, metadata = rag_system.load_knowledge_base(tmp_path)
                    if success:
                        st.success("✅ 読み込み完了")
                        st.session_state.global_db_info = rag_system.get_database_info()
                        st.rerun()
                    else:
                        st.error(message)

    # === クリアグループ ===
    with cols[2]:
        with st.container():
            st.markdown("##### 🗑️ クリア")
            st.markdown("---")

            if has_data:
                st.caption("現在のナレッジベースを削除")
                st.markdown("<br>", unsafe_allow_html=True)
                if st.button("🗑️ データベースをクリア", key=f"{key_prefix}_clear_btn", use_container_width=True, type="primary"):
                    rag_system.clear_knowledge_base()
                    st.session_state.global_db_info = rag_system.get_database_info()
                    st.success("✅ クリア完了")
                    st.rerun()
            else:
                st.caption("クリアするデータがありません")
                st.markdown("<br>", unsafe_allow_html=True)
                st.button("🗑️ データベースをクリア", key=f"{key_prefix}_clear_btn_disabled", use_container_width=True, disabled=True)

# ========================================
# ページ関数
# ========================================

def knowledge_base_management():
    st.title("📚 ナレッジベース管理")
    st.write("RAG機能の削除に伴い、ナレッジベース管理機能は利用できません。")
    st.info("既存のRAGデータベース操作は無効化されています。")


def proofread_meeting_minutes():
    st.title("📝 議事録校正（RAG）")
    st.write("RAG機能は削除されたため、このセクションは利用できません。")
    st.info("校正が必要な場合は、文字起こし結果を手動で編集するか、今後追加される代替機能をご利用ください。")


def dspy_minutes_app():
    st.title("🪄 dspy議事録メイカー")
    st.write("dspyを活用して文字起こしを整理し、簡潔で読みやすい議事録に整形します。")

    st.sidebar.markdown("""
    ### 🪄 dspy議事録メイカー

    **できること**
    - 文字起こしを貼り付けて即座に議事録化
    - スタイルや注目ポイントを指定してリライト
    - 生成結果をコピーやダウンロード

    **ヒント**
    - Azure OpenAIの設定が必要です。
    - dspyが未インストールの場合は `pip install dspy-ai` を実行してください。
    """)

    _init_session_state({
        'dspy_minutes_input_text': "",
        'dspy_minutes_output': "",
        'dspy_minutes_backend': "",
        'dspy_minutes_uploaded_name': "",
        'dspy_minutes_focus': "",
        'dspy_minutes_prompt': DEFAULT_MINUTES_PROMPT,
        'dspy_minutes_dataset': [],
    })

    dspy_module, dspy_error = _load_dspy_module()
    dspy_status = "✅ dspyを利用できます" if dspy_module else f"⚠️ {dspy_error}"

    st.markdown(
        f"<div style='padding:0.5rem 0; color:#4a5568'>接続状態: {dspy_status}</div>",
        unsafe_allow_html=True,
    )

    if not AZURE_OPENAI_ENDPOINT or not AZURE_OPENAI_API_KEY:
        st.warning("Azure OpenAIの設定が不足しています。環境変数 AZURE_OPENAI_ENDPOINT / AZURE_OPENAI_API_KEY を確認してください。")

    with st.container():
        uploaded_text = st.file_uploader(
            "文字起こしテキストファイルを読み込む (txt/md)",
            type=["txt", "md"],
            key="dspy_minutes_file_upload"
        )
        if uploaded_text is not None:
            raw_bytes = uploaded_text.read()
            decoded = raw_bytes.decode("utf-8", errors="replace")
            st.session_state.dspy_minutes_input_text = decoded
            st.session_state.dspy_minutes_uploaded_name = uploaded_text.name
            st.info(f"{uploaded_text.name} を読み込みました。下のテキストを確認してください。")

    st.text_area(
        "文字起こしを貼り付け",
        key="dspy_minutes_input_text",
        height=260,
        placeholder="ここに文字起こし結果を貼り付けてください。不要なタイムスタンプは自動で除去されます。",
    )

    focus_points = st.text_area(
        "強調したい観点 (任意)",
        key="dspy_minutes_focus",
        placeholder="例: 決定事項、宿題、論点、リスク、次回までのTODO など"
    )

    col1, col2 = st.columns(2)
    with col1:
        model_name = st.selectbox(
            "利用するモデル",
            options=["gpt-4o", "gpt-4o-mini"],
            index=0,
            help="dspyおよびAzure OpenAIで利用するモデル名。"
        )
        style_label = st.radio(
            "整形スタイル",
            options=["要点サマリー", "時系列ダイジェスト", "決定事項ファースト"],
            index=0,
        )
    with col2:
        length_hint = st.slider(
            "分量の目安 (段落数)",
            min_value=3,
            max_value=20,
            value=8,
            help="生成する議事録のおおよその長さを指定します。"
        )
        include_todo = st.checkbox("決定事項とTODOを強調する", value=True)

    st.markdown("### プロンプト設定")
    st.caption("dspyに渡す基礎プロンプト。MIPROv2で更新できます。")
    st.text_area(
        "基礎プロンプト",
        key="dspy_minutes_prompt",
        height=140,
    )

    with st.expander("MIPROv2でプロンプトを改善する", expanded=False):
        st.markdown(
            "- transcript（文字起こし）と minutes（熟練者が作成した議事録）のペアをJSON/JSONLでアップロードしてください。\n"
            "- JSONLは1行1サンプル。JSONは配列、または `{'data': [...]} の形式をサポートします。"
        )

        dataset_file = st.file_uploader(
            "プロンプト改善用データセット (json / jsonl)",
            type=["json", "jsonl"],
            key="dspy_minutes_dataset_upload",
            help="各サンプルに transcript と minutes のキーが必要です。",
        )

        if dataset_file is not None:
            dataset, dataset_error = _parse_minutes_dataset(dataset_file)
            if dataset_error:
                st.error(dataset_error)
            else:
                st.session_state.dspy_minutes_dataset = dataset
                st.success(f"{len(dataset)}件のサンプルを読み込みました。")
                preview = dataset[0] if dataset else {}
                if preview:
                    st.caption("サンプルプレビュー")
                    st.json({
                        "transcript": (preview.get("transcript", "")[:80] + "...") if preview.get("transcript") else "",
                        "minutes": (preview.get("minutes", "")[:80] + "...") if preview.get("minutes") else "",
                    })

        opt_col1, opt_col2 = st.columns(2)
        with opt_col1:
            max_iters = st.slider("最適化イテレーション", min_value=1, max_value=8, value=3)
        with opt_col2:
            num_candidates = st.slider("候補プロンプト数", min_value=2, max_value=10, value=4)

        if st.button("MIPROv2でプロンプト最適化", key="dspy_minutes_optimize", use_container_width=True):
            if not st.session_state.dspy_minutes_dataset:
                st.error("最適化用データセットを読み込んでください。")
            else:
                with st.spinner("MIPROv2でプロンプトを改善しています..."):
                    optimized_prompt, opt_error = _optimize_minutes_prompt(
                        st.session_state.dspy_minutes_dataset,
                        st.session_state.dspy_minutes_prompt,
                        model_name,
                        max_iters=max_iters,
                        num_candidates=num_candidates,
                    )
                if opt_error:
                    st.error(opt_error)
                else:
                    st.session_state.dspy_minutes_prompt = optimized_prompt
                    st.success("最適化済みプロンプトを更新しました。")

    base_prompt_text = st.session_state.dspy_minutes_prompt.strip() or DEFAULT_MINUTES_PROMPT

    if st.button("dspyで議事録を生成", type="primary"):
        transcript_text = st.session_state.dspy_minutes_input_text.strip()
        if not transcript_text:
            st.error("文字起こしを入力してください。")
        else:
            directives = base_prompt_text + "\n" + _build_minutes_directives(style_label, focus_points, length_hint, include_todo)
            with st.spinner("dspyで議事録化しています..."):
                minutes_text, error_message = _generate_minutes_with_dspy(transcript_text, directives, model_name)
                backend = "dspy"

                if not minutes_text:
                    backend = "Azure OpenAI"
                    if error_message:
                        st.warning(f"dspy経由の生成に失敗したためAzure OpenAIで再実行します: {error_message}")
                    minutes_text = _generate_minutes_with_fallback(transcript_text, directives, model_name)

                st.session_state.dspy_minutes_output = minutes_text
                st.session_state.dspy_minutes_backend = backend

            st.success(f"{st.session_state.dspy_minutes_backend}で議事録を生成しました。")

    if st.session_state.dspy_minutes_output:
        st.subheader("生成された議事録")
        st.caption(f"出力元: {st.session_state.dspy_minutes_backend}")
        st.text_area(
            "議事録プレビュー",
            value=st.session_state.dspy_minutes_output,
            height=320,
            key="dspy_minutes_output_preview",
        )

        st.download_button(
            label="議事録をテキストでダウンロード",
            data=st.session_state.dspy_minutes_output.encode("utf-8"),
            file_name="dspy_minutes.txt",
            mime="text/plain",
            key="dspy_minutes_download",
        )


def batch_processing_pipeline():
    """一括処理パイプライン: 1本のメディアから複数セグメントを段階的に処理"""
    st.title("🚀 一括処理パイプライン")
    st.write("アップロードした動画/音声ファイルを区間ごとに切り出し、文字起こしから話者識別までをまとめて実行します。")
    st.write("文字起こし → 個別話者識別 → 出力のパイプラインを順番に進めるだけで完了します。")

    st.sidebar.markdown("""
    ### 🚀 一括処理パイプライン

    **概要**
    動画・音声ファイルから必要な区間を切り出し、一括で文字起こし・話者識別を実行します。

    **処理フロー**
    1. メディアファイルのアップロードと区間設定
    2. モデル選択と文字起こし実行
    3. 各ファイルごとの話者識別（個別設定可能）
    4. 処理結果の確認とダウンロード

    **対応形式:** MP4, MOV, AVI, MKV, WebM, MP3, WAV, M4A等
    """)

    # セッション状態の初期化
    _init_session_state({
        'batch_uploaded_video': None,
        'batch_extracted_files': [],
        'batch_processing_results': {},
        'batch_processing_status': {},
        'batch_current_step': 1,
        'batch_rag_system': None,
        'batch_transcribe_model': 'whisper',
        'batch_reference_file': None,
        'batch_current_speaker_file_index': 0,
        'batch_db_info': None,
        'batch_meeting_type': None,
        'batch_default_embeddings': [],
        'batch_file_embeddings': {},  # 各ファイルごとの埋め込み設定
        'batch_embedding_states': {},  # 話者埋め込み作成UI用の状態
        'batch_temp_dir': None,  # 一時ファイルディレクトリ
        'batch_cut_settings_df': None  # 切り出し設定DataFrame
    })

    # Step 1: 動画アップロードと音声切り出し設定
    st.subheader("Step 1: メディアアップロードと音声切り出し")
    st.write("動画・音声ファイルをアップロードし、処理したい区間を設定してください。複数のセグメントを指定できます。")
    st.caption("※ 切り出しは任意です。ファイル全体を使う場合はそのまま次のステップへ進めます。")

    uploaded_video = st.file_uploader(
        "メディアファイルを選択してください",
        type=["wav", "mp3", "mp4", "mov", "avi", "mkv", "webm", "m4a"],
        key="batch_video_upload"
    )

    if uploaded_video is not None:
        st.success(f"✅ メディアファイル '{uploaded_video.name}' を読み込みました")

        # メディアプレビュー
        video_ext = os.path.splitext(uploaded_video.name)[1].lower()
        if video_ext in ['.mp4', '.mov', '.webm', '.avi', '.mkv']:
            st.video(uploaded_video)
        elif video_ext in ['.mp3', '.wav', '.m4a']:
            st.audio(uploaded_video)

        st.subheader("切り出し区間の設定")

        # デフォルトのdata_editorデータ
        if st.session_state.batch_cut_settings_df is None:
            default_data = pd.DataFrame([
                {"開始時間": "00:00:00", "終了時間": "00:00:30", "出力ファイル名": f"{os.path.splitext(uploaded_video.name)[0]}_"}
            ])
        else:
            default_data = st.session_state.batch_cut_settings_df

        edited_df = st.data_editor(
            default_data,
            num_rows="dynamic",
            use_container_width=True,
            column_config={
                "開始時間": st.column_config.TextColumn(
                    "開始時間 (HH:MM:SS or seconds)",
                    help="切り出し開始時間 (例: 00:00:10 または 10)",
                    default="00:00:00"
                ),
                "終了時間": st.column_config.TextColumn(
                    "終了時間 (HH:MM:SS or seconds)",
                    help="切り出し終了時間 (例: 00:00:30 または 30)",
                    default="00:00:30"
                ),
                "出力ファイル名": st.column_config.TextColumn(
                    "出力ファイル名 (.mp3)",
                    help="この区間のMP3出力ファイル名を入力してください (例: meeting1.mp3)。空欄の場合、自動で連番が振られます。",
                    default=f"{os.path.splitext(uploaded_video.name)[0]}_"
                )
            },
            key="batch_cut_settings_editor"
        )

        # 設定を保存
        st.session_state.batch_cut_settings_df = edited_df

        # 音声切り出し実行
        if st.button("🎵 音声を切り出してStep 2へ進む", key="batch_execute_cut", type="primary"):
            if edited_df.empty:
                st.warning("切り出し区間が設定されていません。")
            else:
                with st.spinner("音声の切り出しとMP3への変換中..."):
                    try:
                        # 一時ディレクトリを作成
                        if st.session_state.batch_temp_dir is None:
                            st.session_state.batch_temp_dir = tempfile.mkdtemp()

                        temp_dir = st.session_state.batch_temp_dir
                        temp_video_path = None
                        extracted_files = []

                        # 動画ファイルを一時保存
                        uploaded_video.seek(0)
                        temp_video_path = os.path.join(temp_dir, f"source_video{os.path.splitext(uploaded_video.name)[1]}")
                        with open(temp_video_path, 'wb') as f:
                            f.write(uploaded_video.read())

                        # 各区間を切り出し
                        for index, row in edited_df.iterrows():
                            start_time_str = str(row["開始時間"])
                            end_time_str = str(row["終了時間"])
                            output_filename_raw = str(row["出力ファイル名"]).strip()

                            try:
                                start_seconds = parse_time_to_seconds(start_time_str)
                                end_seconds = parse_time_to_seconds(end_time_str)

                                if start_seconds >= end_seconds:
                                    st.error(f"区間 {index+1}: 開始時間 ({start_time_str}) は終了時間 ({end_time_str}) より前に設定してください。この区間はスキップされます。")
                                    continue

                                # 出力ファイル名を決定
                                base_name_from_video = os.path.splitext(uploaded_video.name)[0]
                                if not output_filename_raw or output_filename_raw.upper() == "AUTO_GENERATE":
                                    output_filename = f"{base_name_from_video}_segment_{index+1}.mp3"
                                else:
                                    output_filename = output_filename_raw

                                # .mp3拡張子を確保
                                if not output_filename.lower().endswith(".mp3"):
                                    output_filename += ".mp3"

                                output_audio_path = os.path.join(temp_dir, output_filename)

                                # FFmpegで切り出し
                                command = [
                                    "ffmpeg",
                                    "-i", temp_video_path,
                                    "-ss", format_time(start_seconds),
                                    "-to", format_time(end_seconds),
                                    "-vn",  # No video
                                    "-ab", "192k",  # Audio bitrate
                                    "-map_metadata", "-1",  # Remove metadata
                                    "-y",  # Overwrite output files without asking
                                    output_audio_path
                                ]

                                process = subprocess.run(command, capture_output=True, text=True, encoding="utf-8", check=True)

                                # 生成されたファイルをメモリに読み込み
                                with open(output_audio_path, 'rb') as f:
                                    file_data = f.read()
                                    file_io = BytesIO(file_data)
                                    file_io.name = output_filename
                                    extracted_files.append({
                                        'name': output_filename,
                                        'data': file_io,
                                        'size': len(file_data)
                                    })

                                st.success(f"✅ 区間 {index+1}: {output_filename} の切り出しが完了しました")

                            except subprocess.CalledProcessError as e:
                                st.error(f"❌ 区間 {index+1}: FFmpegの実行中にエラーが発生しました: {e}")
                                st.code(e.stderr)
                            except ValueError as e:
                                st.error(f"❌ 区間 {index+1}: 時間形式エラー: {e}")
                            except Exception as e:
                                st.error(f"❌ 区間 {index+1}: 処理中にエラーが発生しました: {e}")

                        if extracted_files:
                            # 抽出されたファイルをsession stateに保存
                            st.session_state.batch_extracted_files = extracted_files
                            st.session_state.batch_current_step = 2

                            # 各ファイルの処理状態を初期化
                            for file_info in extracted_files:
                                st.session_state.batch_processing_status[file_info['name']] = {
                                    'transcription': 'pending',
                                    'speaker_id': 'pending',
                                    'rag_proofread': 'pending'
                                }
                                st.session_state.batch_processing_results[file_info['name']] = {}

                            st.success(f"✅ {len(extracted_files)} 個の音声ファイルの切り出しが完了しました")
                            st.rerun()
                        else:
                            st.warning("⚠️ 音声ファイルが1つも生成されませんでした。設定を確認してください。")

                    except Exception as e:
                        st.error(f"❌ 音声切り出し処理中にエラーが発生しました: {e}")
                        import traceback
                        st.error(traceback.format_exc())

        if st.button("⏭ 切り出しをスキップしてStep 2へ進む", key="batch_skip_cut"):
            uploaded_video.seek(0)
            file_bytes = uploaded_video.read()

            if not file_bytes:
                st.error("❌ ファイルの読み込みに失敗しました。もう一度アップロードしてください。")
            else:
                file_buffer = BytesIO(file_bytes)
                file_buffer.name = uploaded_video.name

                st.session_state.batch_extracted_files = [{
                    'name': uploaded_video.name,
                    'data': file_buffer,
                    'size': len(file_bytes)
                }]

                st.session_state.batch_processing_status = {
                    uploaded_video.name: {
                        'transcription': 'pending',
                        'speaker_id': 'pending',
                        'rag_proofread': 'pending'
                    }
                }
                st.session_state.batch_processing_results = {
                    uploaded_video.name: {}
                }
                st.session_state.batch_current_step = 2
                st.success("✅ 切り出しをスキップしました。元のメディアファイルをStep 2でそのまま文字起こしします。")
                st.rerun()

    # Step 2: モデル選択と文字起こし
    if st.session_state.batch_current_step >= 2 and len(st.session_state.batch_extracted_files) > 0:
        st.subheader("Step 2: モデル選択と文字起こし")

        # ファイル一覧を表示
        st.write(f"**検出されたファイル数:** {len(st.session_state.batch_extracted_files)}")

        file_df = pd.DataFrame([
            {
                'ファイル名': f['name'],
                'サイズ (KB)': f'{f["size"] / 1024:.1f}'
            }
            for f in st.session_state.batch_extracted_files
        ])
        st.dataframe(file_df, use_container_width=True, hide_index=True)

        st.divider()

        reference_file = None
        meeting_types = load_meeting_type_config()
        col_settings, col_meeting = st.columns(2, gap="large")

        with col_settings:
            st.write("**文字起こし設定**")

            transcribe_model = st.selectbox(
                "文字起こしモデル",
                options=["whisper", "gpt-4o-transcribe-diarize"],
                index=0,
                key="batch_transcribe_model_select",
                help="whisper: 参考資料対応 | gpt-4o-transcribe-diarize: 自動話者識別付き"
            )

            if transcribe_model == "whisper":
                reference_file = st.file_uploader(
                    "参考資料（オプション）",
                    type=["pdf", "docx", "pptx", "txt", "msg"],
                    key="batch_reference_file_upload",
                    help="全ファイルの文字起こし精度向上に使用"
                )

        with col_meeting:
            st.write("**会議タイプ設定（オプション）**")
            st.write("会議タイプを選択すると、事前に設定された話者埋め込みファイルを自動的に読み込みます。")

            if meeting_types:
                meeting_type_options = {mt['id']: f"{mt['name']} - {mt['description']}" for mt in meeting_types}

                selected_meeting_type_id = st.selectbox(
                    "会議タイプ",
                    options=list(meeting_type_options.keys()),
                    format_func=lambda x: meeting_type_options[x],
                    key="batch_meeting_type_select"
                )

                selected_meeting_type = next((mt for mt in meeting_types if mt['id'] == selected_meeting_type_id), None)

                if selected_meeting_type and selected_meeting_type['embeddings_folder']:
                    if st.session_state.batch_meeting_type != selected_meeting_type_id:
                        st.session_state.batch_meeting_type = selected_meeting_type_id
                        st.session_state.batch_default_embeddings = load_embeddings_from_folder(
                            selected_meeting_type['embeddings_folder']
                        )

                        if st.session_state.batch_default_embeddings:
                            st.success(f"✅ {len(st.session_state.batch_default_embeddings)}個の話者埋め込みファイルを読み込みました")
                            embedding_names = [emb['name'] for emb in st.session_state.batch_default_embeddings]
                            st.info(f"📁 読み込まれたファイル: {', '.join(embedding_names)}")
                elif selected_meeting_type_id == 'custom':
                    st.session_state.batch_meeting_type = 'custom'
                    st.session_state.batch_default_embeddings = []
                    st.info("💡 カスタムモード: Step 3で個別に話者埋め込みファイルを指定してください")
            else:
                st.warning("⚠️ 会議タイプのプリセットが定義されていません。手動で話者埋め込みを指定してください。")

        st.divider()

        # 文字起こし実行
        if st.button("📝 文字起こしを開始", key="batch_start_transcription", type="primary"):
            # 進捗表示用のプレースホルダー
            progress_bar = st.progress(0)
            status_text = st.empty()

            total_files = len(st.session_state.batch_extracted_files)

            for idx, file_info in enumerate(st.session_state.batch_extracted_files):
                file_name = file_info['name']
                file_data = file_info['data']

                status_text.write(f"**文字起こし中: {file_name}** ({idx + 1}/{total_files})")

                try:
                    # 文字起こし実行
                    st.session_state.batch_processing_status[file_name]['transcription'] = 'processing'

                    file_data.seek(0)
                    seg_df = transcribe_audio_to_dataframe(
                        file_data,
                        reference_file=reference_file,
                        model=transcribe_model
                    )

                    st.session_state.batch_processing_results[file_name]['transcription_df'] = seg_df
                    st.session_state.batch_processing_status[file_name]['transcription'] = 'completed'
                    st.success(f"✅ 文字起こし完了: {file_name} ({len(seg_df)}行)")

                except Exception as e:
                    st.error(f"❌ エラー発生: {file_name} - {e}")
                    st.session_state.batch_processing_status[file_name]['transcription'] = 'error'
                    import traceback
                    st.error(traceback.format_exc())

                # 進捗バー更新
                progress_bar.progress((idx + 1) / total_files)

            status_text.write("✅ **すべてのファイルの文字起こしが完了しました！**")
            st.session_state.batch_current_step = 3
            st.balloons()
            st.rerun()

    # Step 3: 個別ファイルの話者識別
    if st.session_state.batch_current_step >= 3 and len(st.session_state.batch_extracted_files) > 0:
        st.subheader("Step 3: 個別ファイルの話者識別")
        st.write("各ファイルごとに話者識別を実行できます。必要に応じてファイルごとに異なる話者埋め込みを使用できます。")

        # 処理状況サマリー
        transcription_completed = sum(1 for status in st.session_state.batch_processing_status.values() if status.get('transcription') == 'completed')
        speaker_id_completed = sum(1 for status in st.session_state.batch_processing_status.values() if status.get('speaker_id') == 'completed')

        col1, col2, col3 = st.columns(3)
        with col1:
            st.metric("文字起こし完了", f"{transcription_completed}/{len(st.session_state.batch_extracted_files)}")
        with col2:
            st.metric("話者識別完了", f"{speaker_id_completed}/{len(st.session_state.batch_extracted_files)}")
        with col3:
            if st.button("🔄 話者識別をスキップして次へ", key="skip_speaker_id"):
                st.session_state.batch_current_step = 4
                st.rerun()

        st.divider()

        # ファイルごとにタブを作成
        file_names = [f['name'] for f in st.session_state.batch_extracted_files]
        file_tabs = st.tabs(file_names)

        for tab_idx, (file_tab, file_info) in enumerate(zip(file_tabs, st.session_state.batch_extracted_files)):
            with file_tab:
                selected_file_name = file_info['name']
                selected_file_info = file_info

                if selected_file_name in st.session_state.batch_processing_results:
                    result = st.session_state.batch_processing_results[selected_file_name]

                    # 文字起こし結果のプレビュー
                    if 'transcription_df' in result:
                        transcription_df = result['transcription_df']
                        with st.expander("📝 文字起こし結果プレビュー", expanded=True):
                            if 'speaker' in transcription_df.columns:
                                disabled_columns = [col for col in transcription_df.columns if col != 'speaker']
                                edited_transcription = st.data_editor(
                                    transcription_df,
                                    hide_index=True,
                                    use_container_width=True,
                                    num_rows="dynamic",
                                    disabled=disabled_columns,
                                    column_config={
                                        "speaker": st.column_config.TextColumn(
                                            "話者",
                                            help="必要に応じて話者名を直接編集してください"
                                        )
                                    },
                                    key=f"transcription_editor_{selected_file_name}_{tab_idx}"
                                )

                                if isinstance(edited_transcription, pd.DataFrame):
                                    updated_transcription_df = edited_transcription.copy()
                                else:
                                    updated_transcription_df = pd.DataFrame(edited_transcription)

                                updated_transcription_df = updated_transcription_df.reset_index(drop=True)
                                transcription_df_normalized = transcription_df.reset_index(drop=True)

                                # Ensure original columns exist in editor output
                                for column in transcription_df_normalized.columns:
                                    if column not in updated_transcription_df.columns:
                                        updated_transcription_df[column] = transcription_df_normalized[column]
                                updated_transcription_df = updated_transcription_df[transcription_df_normalized.columns]

                                original_speaker = transcription_df_normalized['speaker'].astype(str).fillna("")
                                updated_speaker = updated_transcription_df['speaker'].astype(str).fillna("")

                                if not updated_speaker.equals(original_speaker):
                                    st.session_state.batch_processing_results[selected_file_name]['transcription_df'] = updated_transcription_df
                                    transcription_df = updated_transcription_df
                                    st.session_state.batch_processing_results[selected_file_name].pop('identified_df', None)
                                    st.session_state.batch_processing_results[selected_file_name].pop('meeting_text', None)
                                    st.session_state.batch_processing_results[selected_file_name].pop('proofread_text', None)

                                    status_entry = st.session_state.batch_processing_status.get(selected_file_name, {})
                                    status_entry['speaker_id'] = 'pending'
                                    status_entry['rag_proofread'] = 'pending'
                                    st.session_state.batch_processing_status[selected_file_name] = status_entry

                                    st.success("✏️ 話者列の変更を保存しました。必要に応じて話者識別を再実行してください。")
                            else:
                                st.dataframe(transcription_df.head(10), use_container_width=True)

                        tab_identify, tab_embed = st.tabs(["話者識別", "話者埋め込み作成"])

                        with tab_identify:
                            st.write("**話者識別設定**")

                            # デフォルトの埋め込みファイルがあれば表示
                            if st.session_state.batch_default_embeddings:
                                st.info(f"✅ 会議タイプから{len(st.session_state.batch_default_embeddings)}個の話者埋め込みを読み込み済み")
                                default_names = [emb['name'] for emb in st.session_state.batch_default_embeddings]
                                with st.expander("読み込み済み話者埋め込みファイル", expanded=True):
                                    st.write(", ".join(default_names))

                            col1, col2 = st.columns(2)

                            with col1:
                                # 追加の埋め込みファイルをアップロード
                                additional_embeddings = st.file_uploader(
                                    "追加の話者埋め込みファイル（.npy）",
                                    type=["npy"],
                                    accept_multiple_files=True,
                                    key=f"batch_additional_embeddings_{selected_file_name}_{tab_idx}",
                                    help="会議タイプで読み込んだファイルに加えて、追加で埋め込みファイルを指定できます"
                                )

                            with col2:
                                similarity_threshold = st.slider(
                                    "類似度閾値",
                                    min_value=0.0,
                                    max_value=1.0,
                                    value=0.7,
                                    step=0.01,
                                    key=f"batch_similarity_threshold_{selected_file_name}_{tab_idx}"
                                )

                            # 使用する埋め込みファイルを結合
                            # デフォルト埋め込み + 追加アップロード
                            all_embeddings = list(st.session_state.batch_default_embeddings)
                            if additional_embeddings:
                                all_embeddings.extend(additional_embeddings)

                            # 埋め込みファイル数を表示
                            total_embeddings = len(all_embeddings)
                            if total_embeddings > 0:
                                st.success(f"📊 合計 {total_embeddings}個の話者埋め込みファイルを使用します")
                            else:
                                st.warning("⚠️ 話者埋め込みファイルがありません。Step 2で会議タイプを選択するか、上記で追加してください。")

                            # 話者識別実行
                            if all_embeddings and st.button(f"🎤 {selected_file_name} の話者識別を実行", key=f"execute_speaker_id_{selected_file_name}_{tab_idx}", type="primary"):
                                with st.spinner(f"🎤 話者識別中: {selected_file_name}"):
                                    try:
                                        file_data = selected_file_info['data']
                                        file_data.seek(0)

                                        identified_df = identify_speakers_in_dataframe(
                                            file_data,
                                            transcription_df,
                                            all_embeddings,
                                            similarity_threshold
                                        )

                                        st.session_state.batch_processing_results[selected_file_name]['identified_df'] = identified_df
                                        st.session_state.batch_processing_status[selected_file_name]['speaker_id'] = 'completed'
                                        st.success(f"✅ 話者識別完了: {selected_file_name}")
                                        st.rerun()

                                    except Exception as e:
                                        st.error(f"❌ エラー発生: {e}")
                                        st.session_state.batch_processing_status[selected_file_name]['speaker_id'] = 'error'
                                        import traceback
                                        st.error(traceback.format_exc())

                            if 'identified_df' in result and not result['identified_df'].empty:
                                st.divider()
                                st.write("**話者識別結果の手動修正**")
                                st.caption("話者列を直接編集してラベルを調整できます。修正後は必要に応じて再度エクスポートを実行してください。")

                                current_identified_df = result['identified_df'].copy()
                                if 'speaker' not in current_identified_df.columns:
                                    st.warning("話者列が見つかりません。識別結果を再実行してください。")
                                else:
                                    disabled_cols = [col for col in current_identified_df.columns if col != 'speaker']
                                    edited_df = st.data_editor(
                                        current_identified_df,
                                        column_config={
                                            "speaker": st.column_config.TextColumn(
                                                "話者",
                                                help="必要に応じて話者名を直接入力してください"
                                            )
                                        },
                                        disabled=disabled_cols,
                                        hide_index=True,
                                        use_container_width=True,
                                        key=f"identified_editor_{selected_file_name}_{tab_idx}"
                                    )

                                    if isinstance(edited_df, pd.DataFrame):
                                        updated_df = edited_df.copy()
                                    else:
                                        updated_df = pd.DataFrame(edited_df)

                                    updated_df = updated_df.reset_index(drop=True)
                                    current_normalized = current_identified_df.reset_index(drop=True)

                                    # Ensure column order matches original
                                    updated_df = updated_df[current_normalized.columns]

                                    if not updated_df.equals(current_normalized):
                                        st.session_state.batch_processing_results[selected_file_name]['identified_df'] = updated_df
                                        st.session_state.batch_processing_results[selected_file_name].pop('meeting_text', None)
                                        st.session_state.batch_processing_results[selected_file_name].pop('proofread_text', None)

                                        status_entry = st.session_state.batch_processing_status.get(selected_file_name, {})
                                        status_entry['speaker_id'] = 'completed'
                                        status_entry['rag_proofread'] = 'pending'
                                        st.session_state.batch_processing_status[selected_file_name] = status_entry

                                        st.success("✏️ 話者ラベルの変更を保存しました。")
                                        st.info("変更内容を反映するには、必要に応じてStep 4以降を再実行してください。")

                        with tab_embed:
                            st.write("音声セグメントを選択し、プレビュー確認してから話者埋め込みファイルを作成できます。")

                            target_df = result.get('identified_df', transcription_df).copy()

                            if len(target_df) == 0:
                                st.warning("文字起こし結果がありません。まず文字起こしを実行してください。")
                            else:
                                embedding_states = st.session_state.batch_embedding_states
                                if selected_file_name not in embedding_states:
                                    embedding_states[selected_file_name] = {
                                        'selected_rows': set(),
                                        'preview_audio': None,
                                        'show_preview': False
                                    }
                                embedding_state = embedding_states[selected_file_name]

                                selection_mode = st.radio(
                                    "選択方式を選んでください",
                                    options=["範囲指定モード", "チェックボックスモード"],
                                    horizontal=True,
                                    help="範囲指定: 連続した行を素早く選択 | チェックボックス: 飛び飛びの行を自由に選択",
                                    key=f"batch_embedding_selection_mode_{selected_file_name}_{tab_idx}"
                                )

                                st.divider()

                                st.subheader("1️⃣ 音声セグメントを選択")

                                if selection_mode == "範囲指定モード":
                                    embedding_row_labels = _create_row_labels(target_df)
                                    col1, col2 = st.columns(2)
                                    with col1:
                                        start_row = st.selectbox(
                                            "開始行を選択",
                                            options=range(len(target_df)),
                                            format_func=lambda x: embedding_row_labels[x],
                                            key=f"batch_embedding_start_row_{selected_file_name}_{tab_idx}"
                                        )
                                    with col2:
                                        end_row = st.selectbox(
                                            "終了行を選択",
                                            options=range(len(target_df)),
                                            format_func=lambda x: embedding_row_labels[x],
                                            index=len(target_df) - 1 if len(target_df) > 0 else 0,
                                            key=f"batch_embedding_end_row_{selected_file_name}_{tab_idx}"
                                        )

                                    if start_row > end_row:
                                        st.error("⚠️ 開始行は終了行以前を選択してください")
                                        embedding_state['selected_rows'] = set()
                                    else:
                                        embedding_state['selected_rows'] = set(range(start_row, end_row + 1))
                                        st.success(f"✅ 行 {start_row} ～ {end_row} を選択しました")
                                else:
                                    display_df = target_df.copy()
                                    display_df.insert(0, "選択", False)

                                    if embedding_state['selected_rows']:
                                        for idx in embedding_state['selected_rows']:
                                            if idx in display_df.index:
                                                display_df.at[idx, "選択"] = True

                                    disabled_columns = [col for col in ["start", "end", "speaker", "text"] if col in display_df.columns]

                                    edited_display = st.data_editor(
                                        display_df,
                                        column_config={
                                            "選択": st.column_config.CheckboxColumn(
                                                "選択",
                                                help="埋め込み作成対象の行をチェック",
                                                default=False
                                            )
                                        },
                                        disabled=disabled_columns,
                                        use_container_width=True,
                                        hide_index=True,
                                        key=f"batch_embedding_data_editor_{selected_file_name}_{tab_idx}"
                                    )

                                    embedding_state['selected_rows'] = set(
                                        edited_display[edited_display["選択"] == True].index.tolist()
                                    )

                                selected_rows = embedding_state['selected_rows']

                                if selected_rows:
                                    st.subheader("2️⃣ 選択情報の確認")

                                    selection_summary = get_selection_summary(target_df, selected_rows)

                                    info_cols = st.columns(4)
                                    with info_cols[0]:
                                        st.metric("選択行数", selection_summary['count'])
                                    with info_cols[1]:
                                        st.metric("開始時刻", format_time(selection_summary['start_time']))
                                    with info_cols[2]:
                                        st.metric("終了時刻", format_time(selection_summary['end_time']))
                                    with info_cols[3]:
                                        st.metric("音声長", f"{selection_summary['duration']:.1f}秒")

                                    if selection_summary['speakers']:
                                        speakers_text = "、".join([s if s else "不明" for s in selection_summary['speakers']])
                                        st.info(f"🎤 選択範囲に含まれる話者: {speakers_text}")

                                    st.subheader("3️⃣ 音声プレビュー")

                                    col_preview, col_clear = st.columns([3, 1])

                                    with col_preview:
                                        if st.button("🔊 プレビュー音声を生成", key=f"batch_generate_preview_{selected_file_name}_{tab_idx}"):
                                            try:
                                                with st.spinner("プレビュー音声を生成中..."):
                                                    audio_io = BytesIO(selected_file_info['data'].getvalue())
                                                    preview_bytes, duration = prepare_embedding_preview_audio(
                                                        audio_io,
                                                        target_df,
                                                        sorted(selected_rows)
                                                    )
                                                    embedding_state['preview_audio'] = preview_bytes
                                                    embedding_state['show_preview'] = True
                                                st.success(f"プレビュー音声を生成しました（{duration:.1f}秒）")
                                            except Exception as e:
                                                st.error(f"プレビュー生成エラー: {e}")

                                    with col_clear:
                                        if embedding_state.get('show_preview'):
                                            if st.button("❌ プレビューをクリア", key=f"batch_clear_preview_{selected_file_name}_{tab_idx}"):
                                                embedding_state['preview_audio'] = None
                                                embedding_state['show_preview'] = False
                                                st.rerun()

                                    if embedding_state.get('show_preview') and embedding_state.get('preview_audio'):
                                        st.audio(embedding_state['preview_audio'], format="audio/wav")

                                    st.subheader("4️⃣ ファイル名設定")

                                    embedding_filename = st.text_input(
                                        "ファイル名（.npy拡張子は自動追加）",
                                        value="speaker_embedding",
                                        key=f"batch_embedding_filename_{selected_file_name}_{tab_idx}",
                                        help="作成する話者埋め込みファイルの名前を指定してください"
                                    )

                                    st.subheader("5️⃣ 埋め込み作成")

                                    if st.button("✨ 話者埋め込みを作成してダウンロード", key=f"batch_create_embedding_{selected_file_name}_{tab_idx}", type="primary"):
                                        with st.spinner("話者埋め込みを作成中..."):
                                            try:
                                                audio_io = BytesIO(selected_file_info['data'].getvalue())
                                                embedding, duration = extract_audio_segment_for_embedding(
                                                    audio_io,
                                                    target_df,
                                                    sorted(selected_rows)
                                                )

                                                filename_with_ext = embedding_filename if embedding_filename.endswith('.npy') else f"{embedding_filename}.npy"

                                                embedding_io = BytesIO()
                                                np.save(embedding_io, embedding)
                                                embedding_io.seek(0)

                                                embedding_bytes = embedding_io.getvalue()

                                                st.success(f"✅ 話者埋め込みの作成が完了しました（音声長: {duration:.1f}秒）")
                                                trigger_auto_download(
                                                    embedding_bytes,
                                                    filename_with_ext,
                                                    key=f"batch_download_embedding_{selected_file_name}_{tab_idx}",
                                                    mime="application/octet-stream"
                                                )

                                            except Exception as e:
                                                st.error(f"❌ 話者埋め込みの作成中にエラーが発生しました: {e}")
                                                import traceback
                                                with st.expander("🔍 詳細なエラー情報"):
                                                    st.code(traceback.format_exc())
                                else:
                                    if selection_mode == "範囲指定モード":
                                        st.info("💡 上記のselectboxで開始行と終了行を選択してください")
                                    else:
                                        st.info("💡 データエディタで埋め込み作成対象の行をチェックしてください")

                    # 話者識別結果の確認
                    if 'identified_df' in result:
                        st.success("✅ このファイルの話者識別は完了しています")
                else:
                    st.warning("⚠️ このファイルはまだ文字起こしが完了していません")

        st.divider()

        # 全ファイルの処理状況一覧
        st.write("**全ファイルの処理状況**")
        status_df = pd.DataFrame([
            {
                'ファイル名': f['name'],
                '文字起こし': '✅' if st.session_state.batch_processing_status.get(f['name'], {}).get('transcription') == 'completed' else '❌',
                '話者識別': '✅' if st.session_state.batch_processing_status.get(f['name'], {}).get('speaker_id') == 'completed' else '⏭️'
            }
            for f in st.session_state.batch_extracted_files
        ])
        st.dataframe(status_df, use_container_width=True, hide_index=True)

        st.divider()

        # 次のステップへ
        if st.button("➡️ 次のステップへ進む", key="proceed_to_rag", type="primary"):
            st.session_state.batch_current_step = 4
            st.rerun()

    # Step 4: 校正（RAGなし）
    if st.session_state.batch_current_step >= 4:
        st.subheader("Step 4: 校正")
        st.write("RAG機能は削除されたため、このステップでは自動校正をスキップします。")
        st.session_state.batch_processing_status = {name: {**status, 'rag_proofread': 'skipped'} for name, status in st.session_state.batch_processing_status.items()}
        if st.button("➡️ 最終確認へ進む", key="skip_rag_step", type="primary"):
            st.session_state.batch_current_step = 5
            st.rerun()

    if st.session_state.batch_current_step >= 5:
        st.subheader("Step 5: 処理結果の確認とダウンロード")

        # 処理状態サマリー
        st.write("**処理状態サマリー**")
        status_summary = []
        for file_name, status in st.session_state.batch_processing_status.items():
            status_summary.append({
                'ファイル名': file_name,
                '文字起こし': '✅' if status['transcription'] == 'completed' else '❌' if status['transcription'] == 'error' else '⏭️',
                '話者識別': '✅' if status['speaker_id'] == 'completed' else '❌' if status['speaker_id'] == 'error' else '⏭️',
                '校正': '✅' if status['rag_proofread'] == 'completed' else '❌' if status['rag_proofread'] == 'error' else '⏭️'
            })

        st.dataframe(pd.DataFrame(status_summary), use_container_width=True, hide_index=True)

        st.divider()

        # 個別ファイルの結果表示
        st.write("**個別ファイルの結果**")

        # ファイルごとにタブを作成
        file_names = [f['name'] for f in st.session_state.batch_extracted_files]
        file_tabs = st.tabs(file_names)

        for tab_idx, (file_tab, file_info) in enumerate(zip(file_tabs, st.session_state.batch_extracted_files)):
            with file_tab:
                selected_file = file_info['name']

                if selected_file in st.session_state.batch_processing_results:
                    result = st.session_state.batch_processing_results[selected_file]

                    content_tab1, content_tab2, content_tab3 = st.tabs(["📝 文字起こし結果", "🎤 話者識別結果", "📚 校正結果"])

                    with content_tab1:
                        if 'transcription_df' in result:
                            st.dataframe(result['transcription_df'], use_container_width=True)

                            # Excel形式でダウンロード
                            excel_buffer = BytesIO()
                            with pd.ExcelWriter(excel_buffer, engine='openpyxl') as writer:
                                result['transcription_df'].to_excel(writer, index=False, sheet_name='文字起こし')
                            excel_buffer.seek(0)

                            st.download_button(
                                label="📥 Excel形式でダウンロード",
                                data=excel_buffer,
                                file_name=f"{os.path.splitext(selected_file)[0]}_transcription.xlsx",
                                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                                key=f"download_excel_{selected_file}_{tab_idx}"
                            )
                        else:
                            st.info("文字起こし結果がありません")

                    with content_tab2:
                        meeting_text = result.get('meeting_text')
                        if not meeting_text:
                            source_df = result.get('identified_df') or result.get('transcription_df')
                            meeting_text = build_meeting_text_from_dataframe(source_df) if source_df is not None else ""
                            if meeting_text:
                                st.session_state.batch_processing_results[selected_file]['meeting_text'] = meeting_text

                        if meeting_text:
                            st.text_area(
                                "議事録形式テキスト",
                                meeting_text,
                                height=400,
                                key=f"meeting_text_view_{selected_file}_{tab_idx}"
                            )

                            doc = DocxDocument()
                            doc.add_heading(f'議事録: {selected_file}', level=1)
                            for line in meeting_text.splitlines():
                                doc.add_paragraph(line)

                            docx_buffer = BytesIO()
                            doc.save(docx_buffer)
                            docx_buffer.seek(0)

                            st.download_button(
                                label="📥 議事録（Word）をダウンロード",
                                data=docx_buffer,
                                file_name=f"{os.path.splitext(selected_file)[0]}_minutes.docx",
                                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                                key=f"download_docx_{selected_file}_{tab_idx}"
                            )
                        else:
                            st.info("議事録形式テキストがありません（話者識別をスキップした場合はStep 2の結果を確認してください）")

                    with content_tab3:
                        if 'proofread_text' in result:
                            st.text_area(
                                "校正後テキスト",
                                value=result['proofread_text'],
                                height=400,
                                key=f"proofread_view_{selected_file}_{tab_idx}"
                            )

                            # Word形式でダウンロード
                            doc = DocxDocument()
                            doc.add_heading(f'校正済み議事録: {selected_file}', level=1)
                            doc.add_paragraph(result['proofread_text'])

                            docx_buffer = BytesIO()
                            doc.save(docx_buffer)
                            docx_buffer.seek(0)

                            st.download_button(
                                label="📥 校正済み議事録（Word）をダウンロード",
                                data=docx_buffer,
                                file_name=f"{os.path.splitext(selected_file)[0]}_proofread.docx",
                                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                                key=f"download_proofread_{selected_file}_{tab_idx}"
                            )
                        else:
                            st.info("校正結果がありません（スキップされたか、エラーが発生しました）")
                else:
                    st.info("このファイルの処理結果がありません")

        st.divider()

        # 一括ダウンロード機能
        st.write("**一括ダウンロード**")
        st.write("すべての処理結果をまとめてダウンロードできます。")

        col1, col2, col3 = st.columns(3)

        with col1:
            if st.button("📦 ZIPを作成してダウンロード", key="batch_download_all", use_container_width=True):
                with st.spinner("ZIPファイルを作成中..."):
                    try:
                        zip_buffer = BytesIO()

                        with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zip_file:
                            for file_name, result in st.session_state.batch_processing_results.items():
                                base_name = os.path.splitext(file_name)[0]

                                # 文字起こし結果（Excel）
                                if 'transcription_df' in result:
                                    excel_buffer = BytesIO()
                                    with pd.ExcelWriter(excel_buffer, engine='openpyxl') as writer:
                                        result['transcription_df'].to_excel(writer, index=False, sheet_name='文字起こし')
                                    excel_buffer.seek(0)
                                    zip_file.writestr(f"{base_name}_transcription.xlsx", excel_buffer.read())

                                # 話者識別結果（Word）
                                if 'identified_df' in result:
                                    doc = DocxDocument()
                                    doc.add_heading(f'議事録: {file_name}', level=1)

                                    # 連続する同じ話者の発言を結合
                                    df = result['identified_df'].copy()
                                    df['speaker_filled'] = df['speaker'].replace('', pd.NA)
                                    df['speaker_filled'] = df['speaker_filled'].ffill()
                                    df['group_id'] = (df['speaker_filled'] != df['speaker_filled'].shift()).cumsum()
                                    df_merged = df.groupby('group_id').agg(
                                        speaker=('speaker_filled', 'first'),
                                        text=('text', ' '.join)
                                    ).reset_index(drop=True)

                                    for _, row in df_merged.iterrows():
                                        speaker = row.get('speaker', '不明')
                                        text = row.get('text', '')
                                        speaker_str = speaker if speaker else '不明'
                                        doc.add_paragraph(f"（{speaker_str}）{text}")

                                    docx_buffer = BytesIO()
                                    doc.save(docx_buffer)
                                    docx_buffer.seek(0)
                                    zip_file.writestr(f"{base_name}_minutes.docx", docx_buffer.read())

                                # RAG校正結果（Word）
                                if 'proofread_text' in result:
                                    doc = DocxDocument()
                                    doc.add_heading(f'校正済み議事録: {file_name}', level=1)
                                    doc.add_paragraph(result['proofread_text'])
                                    docx_buffer = BytesIO()
                                    doc.save(docx_buffer)
                                    docx_buffer.seek(0)
                                    zip_file.writestr(f"{base_name}_proofread.docx", docx_buffer.read())

                        zip_buffer.seek(0)

                        trigger_auto_download(
                            zip_buffer.getvalue(),
                            file_name=f"batch_results_{datetime.now().strftime('%Y%m%d_%H%M%S')}.zip",
                            key="batch_final_download",
                            mime="application/zip"
                        )

                        st.success("✅ ZIPファイルの作成が完了しました！")

                    except Exception as e:
                        st.error(f"❌ ZIPファイルの作成中にエラーが発生しました: {e}")
                        import traceback
                        st.error(traceback.format_exc())

        with col2:
            if st.button("📝 文字起こし結果Wordを作成してダウンロード", key="batch_download_transcription_word", use_container_width=True):
                with st.spinner("Wordファイルを作成中..."):
                    try:
                        # 1つのWordドキュメントを作成
                        doc = DocxDocument()
                        doc.add_heading('文字起こし結果（一括）', level=0)
                        doc.add_paragraph(f'作成日時: {datetime.now().strftime("%Y年%m月%d日 %H:%M:%S")}')
                        doc.add_paragraph('')

                        # 各ファイルの結果を追加
                        for idx, (file_name, result) in enumerate(st.session_state.batch_processing_results.items(), 1):
                            # ファイル見出し
                            doc.add_heading(f'{idx}. {file_name}', level=1)
                            doc.add_paragraph('')

                            base_df = result.get('identified_df')
                            if base_df is None or base_df.empty:
                                fallback_df = result.get('transcription_df')
                                base_df = fallback_df if fallback_df is not None else None
                            meeting_text = result.get('meeting_text')
                            if meeting_text is None and base_df is not None and not base_df.empty:
                                meeting_text = build_meeting_text_from_dataframe(base_df)
                                if meeting_text:
                                    st.session_state.batch_processing_results[file_name]['meeting_text'] = meeting_text

                            if meeting_text:
                                doc.add_heading('議事録（話者識別結果）', level=2)
                                for line in meeting_text.splitlines():
                                    if line.strip():
                                        doc.add_paragraph(line)
                            else:
                                doc.add_paragraph('（文字起こし結果なし）')

                            doc.add_paragraph('')

                            # ファイル間の区切り
                            if idx < len(st.session_state.batch_processing_results):
                                doc.add_page_break()

                        # Wordファイルを保存
                        docx_buffer = BytesIO()
                        doc.save(docx_buffer)
                        docx_buffer.seek(0)

                        trigger_auto_download(
                            docx_buffer.getvalue(),
                            file_name=f"batch_transcription_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx",
                            key="batch_final_download_transcription_word",
                            mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                        )

                        st.success("✅ Wordファイルの作成が完了しました！")

                    except Exception as e:
                        st.error(f"❌ Wordファイルの作成中にエラーが発生しました: {e}")
                        import traceback
                        st.error(traceback.format_exc())

        with col3:
            if st.button("📚 校正結果Wordを作成してダウンロード", key="batch_download_rag_word", use_container_width=True):
                with st.spinner("Wordファイルを作成中..."):
                    try:
                        # 1つのWordドキュメントを作成
                        doc = DocxDocument()
                        doc.add_heading('校正結果（一括）', level=0)
                        doc.add_paragraph(f'作成日時: {datetime.now().strftime("%Y年%m月%d日 %H:%M:%S")}')
                        doc.add_paragraph('')

                        # 校正結果があるファイルのみ処理
                        rag_count = 0
                        for idx, (file_name, result) in enumerate(st.session_state.batch_processing_results.items(), 1):
                            if 'proofread_text' in result:
                                rag_count += 1
                                # ファイル見出し
                                doc.add_heading(f'{rag_count}. {file_name}', level=1)
                                doc.add_paragraph('')

                                # 校正結果
                                doc.add_heading('校正済み議事録', level=2)
                                # 段落ごとに分割して追加
                                for paragraph in result['proofread_text'].split('\n'):
                                    if paragraph.strip():
                                        doc.add_paragraph(paragraph)
                                doc.add_paragraph('')

                                # ファイル間の区切り（最後のファイル以外）
                                if rag_count < sum(1 for r in st.session_state.batch_processing_results.values() if 'proofread_text' in r):
                                    doc.add_page_break()

                        if rag_count == 0:
                            doc.add_paragraph('（校正結果がありません。このステップはスキップされています）')

                        # Wordファイルを保存
                        docx_buffer = BytesIO()
                        doc.save(docx_buffer)
                        docx_buffer.seek(0)

                        trigger_auto_download(
                            docx_buffer.getvalue(),
                            file_name=f"batch_rag_proofread_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx",
                            key="batch_final_download_rag_word",
                            mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                        )

                        st.success("✅ Wordファイルの作成が完了しました！")

                    except Exception as e:
                        st.error(f"❌ Wordファイルの作成中にエラーが発生しました: {e}")
                        import traceback
                        st.error(traceback.format_exc())


def video_to_audio_cutter_app():
    st.title("動画から音声を切り出しMP3で保存")
    st.write("動画ファイルをアップロードし、切り出したい開始時間と終了時間を指定してください。複数の区間を切り出すことができます。")

    uploaded_video = st.file_uploader("動画ファイルを選択", type=["wav","mp3","mp4", "mov", "avi", "mkv", "webm"])

    if uploaded_video is not None:
        st.video(uploaded_video)

        st.subheader("切り出し区間の設定")
        # Use st.data_editor for multiple time range inputs
        # Default for the first row includes segment_1
        default_data = pd.DataFrame([
            {"開始時間": "00:00:00", "終了時間": "00:00:30", "出力ファイル名": f"{os.path.splitext(uploaded_video.name)[0]}_"}
        ])
        edited_df = st.data_editor(
            default_data,
            num_rows="dynamic",
            use_container_width=True,
            column_config={
                "開始時間": st.column_config.TextColumn(
                    "開始時間 (HH:MM:SS or seconds)",
                    help="切り出し開始時間 (例: 00:00:10 または 10)",
                    default="00:00:00"
                ),
                "終了時間": st.column_config.TextColumn(
                    "終了時間 (HH:MM:SS or seconds)",
                    help="切り出し終了時間 (例: 00:00:30 または 30)",
                    default="00:00:30"
                ),
                "出力ファイル名": st.column_config.TextColumn(
                    "出力ファイル名 (.mp3)",
                    help="この区間のMP3出力ファイル名を入力してください (例: my_audio_segment.mp3)。'AUTO_GENERATE'と入力するか空欄の場合、自動で連番が振られます。",
                    default=f"{os.path.splitext(uploaded_video.name)[0]}_" # Explicit placeholder for new rows
                )
            }
        )

        if st.button("音声を切り出してMP3で保存"):
            if edited_df.empty:
                st.warning("切り出し区間が設定されていません。")
                return

            temp_video_path = ""
            output_audio_paths = [] # List to store paths of all generated MP3s
            zip_buffer = BytesIO()

            try:
                # Save uploaded video to a temporary file
                with tempfile.NamedTemporaryFile(delete=False, suffix=f".{uploaded_video.name.split('.')[-1]}") as temp_video_file:
                    temp_video_file.write(uploaded_video.read())
                    temp_video_path = temp_video_file.name

                with st.spinner("音声の切り出しとMP3への変換中..."):
                    for index, row in edited_df.iterrows():
                        start_time_str = str(row["開始時間"])
                        end_time_str = str(row["終了時間"])
                        output_filename_raw = str(row["出力ファイル名"]).strip()

                        try:
                            start_seconds = parse_time_to_seconds(start_time_str)
                            end_seconds = parse_time_to_seconds(end_time_str)

                            if start_seconds >= end_seconds:
                                st.error(f"区間 {index+1}: 開始時間 ({start_time_str}) は終了時間 ({end_time_str}) より前に設定してください。この区間はスキップされます。")
                                continue

                            # If output filename is empty or matches the explicit placeholder, generate one with index
                            base_name_from_video = os.path.splitext(uploaded_video.name)[0]

                            if not output_filename_raw or output_filename_raw.upper() == "AUTO_GENERATE":
                                output_filename_to_use = f"{base_name_from_video}_segment_{index+1}.mp3"
                            else:
                                output_filename_to_use = output_filename_raw

                            # Ensure the output filename ends with .mp3
                            if not output_filename_to_use.lower().endswith(".mp3"):
                                output_filename_to_use += ".mp3"

                            output_audio_path = os.path.join(tempfile.gettempdir(), output_filename_to_use)

                            command = [
                                "ffmpeg",
                                "-i", temp_video_path,
                                "-ss", format_time(start_seconds),
                                "-to", format_time(end_seconds),
                                "-vn",  # No video
                                "-ab", "192k", # Audio bitrate
                                "-map_metadata", "-1", # Remove metadata
                                "-y", # Overwrite output files without asking
                                output_audio_path
                            ]

                            st.info(f"区間 {index+1} FFmpegコマンドを実行中: {' '.join(command)}")

                            process = subprocess.run(command, capture_output=True, text=True, encoding="utf-8", check=True)
                            st.success(f"区間 {index+1} の音声切り出しとMP3への変換が完了しました！")
                            st.code(process.stdout)
                            st.code(process.stderr)
                            output_audio_paths.append(output_audio_path)

                        except subprocess.CalledProcessError as e:
                            st.error(f"区間 {index+1} FFmpegの実行中にエラーが発生しました: {e}")
                            st.code(e.stdout)
                            st.code(e.stderr)
                            st.warning("FFmpegがシステムにインストールされ、PATHが通っていることを確認してください。")
                        except ValueError as e:
                            st.error(f"区間 {index+1} 時間形式エラー: {e}")
                        except Exception as e:
                            st.error(f"区間 {index+1} 処理中にエラーが発生しました: {e}")

                if output_audio_paths:
                    st.subheader("生成されたMP3ファイル")
                    with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zf:
                        for audio_path in output_audio_paths:
                            if os.path.exists(audio_path):
                                zf.write(audio_path, os.path.basename(audio_path))
                                st.write(f"- {os.path.basename(audio_path)}")
                    zip_buffer.seek(0)

                    st.download_button(
                        label="全てのMP3ファイルをまとめてダウンロード (ZIP)",
                        data=zip_buffer,
                        file_name=f"{os.path.splitext(uploaded_video.name)[0]}_cut_audios.zip",
                        mime="application/zip"
                    )
                else:
                    st.warning("切り出されたMP3ファイルはありませんでした。")

            except Exception as e:
                st.error(f"動画ファイルの処理中にエラーが発生しました: {e}")
            finally:
                # Clean up temporary files
                if os.path.exists(temp_video_path):
                    os.remove(temp_video_path)
                for audio_path in output_audio_paths:
                    if os.path.exists(audio_path):
                        os.remove(audio_path)

def _init_session_state(defaults):
    """セッション状態を初期化"""
    for key, value in defaults.items():
        if key not in st.session_state:
            st.session_state[key] = value


def _load_dspy_module():
    """dspyの読み込みを試み、モジュールと状態を返却"""
    spec = importlib.util.find_spec("dspy")
    if spec is None:
        return None, "dspyがインストールされていません。`pip install dspy-ai` を実行してください。"

    dspy = importlib.import_module("dspy")
    return dspy, None


DEFAULT_MINUTES_PROMPT = (
    "あなたは熟練の議事録専門家です。発言の意図を汲み取り、決定事項・TODO・論点を中心に、"
    "簡潔で読みやすい日本語の段落に整形してください。不要なノイズやタイムスタンプは除去します。"
)


def _build_minutes_directives(style_label: str, focus_points: str, length_hint: int, include_todo: bool) -> str:
    """議事録整形のディレクティブ文字列を生成"""
    focus_text = focus_points.strip() if focus_points else "決定事項・TODO・論点を中心に整理してください。"
    style_templates = {
        "要点サマリー": "重要な意思決定・論点・TODOを見出し付きで箇条書き。1項目につき1-2文で端的に。",
        "時系列ダイジェスト": "議事進行の順に、発言のまとまりごとに短い段落でまとめる。流れが追いやすいよう接続詞を適度に配置。",
        "決定事項ファースト": "決定事項・合意事項を先頭にまとめ、続けて根拠や懸念点を簡潔に列挙。",
    }
    todo_line = "決定事項とTODOは太字の見出しでまとめ、箇条書きで簡潔に書いてください。" if include_todo else "重要箇条書きは簡潔にまとめてください。"

    return (
        f"整理スタイル: {style_templates.get(style_label, style_label)}\n"
        f"フォーカス: {focus_text}\n"
        f"長さの目安: {length_hint} 段落程度で簡潔に\n"
        f"{todo_line}\n"
        "時刻表現やノイズは除去し、日本語で読みやすく編集します。"
    )


def _parse_minutes_dataset(uploaded_file):
    """プロンプト改善用のデータセットを読み込む"""
    try:
        content = uploaded_file.read().decode("utf-8")
    except Exception as e:
        return [], f"データセットの読み込みに失敗しました: {e}"

    records = []
    try:
        if uploaded_file.name.lower().endswith(".jsonl"):
            for line in content.splitlines():
                if not line.strip():
                    continue
                data = json.loads(line)
                transcript = data.get("transcript")
                minutes = data.get("minutes")
                if transcript and minutes:
                    records.append({"transcript": transcript, "minutes": minutes})
        else:
            data = json.loads(content)
            if isinstance(data, list):
                for item in data:
                    transcript = item.get("transcript")
                    minutes = item.get("minutes")
                    if transcript and minutes:
                        records.append({"transcript": transcript, "minutes": minutes})
            elif isinstance(data, dict):
                dataset_items = data.get("data") or []
                for item in dataset_items:
                    transcript = item.get("transcript")
                    minutes = item.get("minutes")
                    if transcript and minutes:
                        records.append({"transcript": transcript, "minutes": minutes})
    except json.JSONDecodeError as e:
        return [], f"JSONのパースに失敗しました: {e}"

    if not records:
        return [], "有効なtranscript・minutesのペアが見つかりませんでした。"

    return records, None


def _minutes_similarity_metric(example, prediction, trace=None):
    """MIPROv2用の簡易類似度メトリクス"""
    target_minutes = getattr(example, "minutes", "") or ""
    predicted_minutes = getattr(prediction, "minutes", "") or ""

    if not target_minutes or not predicted_minutes:
        return 0.0

    target_tokens = set(target_minutes.split())
    predicted_tokens = set(predicted_minutes.split())
    if not target_tokens:
        return 0.0

    overlap = len(target_tokens & predicted_tokens)
    return overlap / len(target_tokens)


def _optimize_minutes_prompt(dataset, base_prompt, model_name, max_iters=3, num_candidates=4):
    """MIPROv2を用いて議事録プロンプトを最適化"""
    dspy, error_message = _load_dspy_module()
    if dspy is None:
        return None, error_message

    teleprompt_spec = importlib.util.find_spec("dspy.teleprompt")
    if teleprompt_spec is None:
        return None, "dspy.teleprompt モジュールが見つかりません。dspyのバージョンを確認してください。"

    teleprompt_module = importlib.import_module("dspy.teleprompt")
    if not hasattr(teleprompt_module, "MIPROv2"):
        return None, "MIPROv2 が利用できません。dspyをアップデートしてください。"

    MIPROv2 = getattr(teleprompt_module, "MIPROv2")

    azure_lm = dspy.AzureOpenAI(
        model=model_name,
        api_base=AZURE_OPENAI_ENDPOINT,
        api_key=AZURE_OPENAI_API_KEY,
        api_version=API_VERSION,
        max_tokens=1200,
        temperature=0.3,
    )
    dspy.settings.configure(lm=azure_lm)

    class MinutesRewrite(dspy.Signature):
        """MIPRO用の議事録リライトシグネチャ"""

        transcript: str = dspy.InputField(desc="元の文字起こし")
        refinement_directives: str = dspy.InputField(desc="整形方針")
        minutes: str = dspy.OutputField(desc="整形済み議事録")

    program = dspy.Predict(MinutesRewrite)

    trainset = []
    for item in dataset:
        transcript = item.get("transcript")
        minutes = item.get("minutes")
        if not transcript or not minutes:
            continue
        example = dspy.Example(
            transcript=transcript,
            refinement_directives=base_prompt,
            minutes=minutes,
        ).with_inputs("transcript", "refinement_directives")
        trainset.append(example)

    if not trainset:
        return None, "学習に使えるサンプルがありません。"

    teleprompter = MIPROv2(
        metric=_minutes_similarity_metric,
        init_prompt=base_prompt,
        num_candidates=num_candidates,
        max_iters=max_iters,
        allow_refusal=False,
        verbose=False,
    )

    optimized_program = teleprompter.compile(program, trainset=trainset)

    optimized_prompt = getattr(teleprompter, "best_prompt", None)
    if not optimized_prompt:
        optimized_prompt = getattr(optimized_program, "prompt", None)
    if not optimized_prompt and hasattr(optimized_program, "signature"):
        optimized_prompt = getattr(optimized_program.signature, "instructions", None)

    return optimized_prompt or base_prompt, None


def _generate_minutes_with_dspy(transcript_text: str, directives: str, model_name: str):
    """dspyを利用して議事録を生成"""
    dspy, error_message = _load_dspy_module()
    if dspy is None:
        return None, error_message

    try:
        azure_lm = dspy.AzureOpenAI(
            model=model_name,
            api_base=AZURE_OPENAI_ENDPOINT,
            api_key=AZURE_OPENAI_API_KEY,
            api_version=API_VERSION,
            max_tokens=1200,
            temperature=0.3,
        )
        dspy.settings.configure(lm=azure_lm)

        class MinutesRewrite(dspy.Signature):
            """文字起こしを議事録形式に整える"""

            transcript: str = dspy.InputField(desc="元の文字起こし")
            refinement_directives: str = dspy.InputField(desc="整形方針")
            minutes: str = dspy.OutputField(desc="整形済み議事録")

        predictor = dspy.Predict(MinutesRewrite)
        result = predictor(
            transcript=transcript_text,
            refinement_directives=directives
        )

        minutes_text = getattr(result, "minutes", None)
        if not minutes_text and hasattr(result, "response"):
            minutes_text = getattr(result, "response")

        return minutes_text, None
    except Exception as e:
        return None, str(e)


def _generate_minutes_with_fallback(transcript_text: str, directives: str, model_name: str):
    """Azure OpenAIを使ったフォールバックの議事録生成"""
    fallback_system_prompt = (
        "あなたは議事録要約の専門家です。文字起こしを読みやすく整理し、参加者がすぐに振り返れる形でまとめてください。\n"
        "不要な相槌やノイズは除去し、文体を整えます。決定事項・TODO・論点が分かるように短い見出しを付けてください。\n"
        f"{directives}\n"
        "出力は日本語で、過度に長くならないようにしてください。"
    )

    return generate_summary(model_name, fallback_system_prompt, transcript_text)


def _create_row_labels(df):
    """データフレームから行選択用ラベルを作成"""
    labels = []
    for idx, row in df.iterrows():
        text = row.get('text', '')
        text_display = text[:30] + "..." if len(text) > 30 else text
        label = f"{idx}: {row.get('start', '')} | {row.get('end', '')} | {row.get('speaker', '')} | {text_display}"
        labels.append(label)
    return labels

def prepare_embedding_preview_audio(audio_io, df, selected_indices):
    """プレビュー用の音声を準備

    Args:
        audio_io: BytesIO - 音声ファイル
        df: DataFrame - タイムスタンプ含む
        selected_indices: list - 選択行インデックスリスト

    Returns:
        preview_bytes: bytes - WAV形式の音声バイト列
        duration: float - セグメント長（秒）
    """
    if not selected_indices:
        raise ValueError("選択行がありません")

    selected_rows = df.iloc[sorted(selected_indices)]
    start_sec = selected_rows['start'].min()
    end_sec = selected_rows['end'].max()

    audio_io.seek(0)
    with temp_file_path(audio_io.getvalue(), ".mp3") as audio_path:
        audio = AudioSegment.from_file(audio_path)

        start_ms = int(start_sec * 1000)
        end_ms = int(end_sec * 1000)
        audio_segment = audio[start_ms:end_ms]

        preview_bytes = audio_segment.export(format="wav").read()

    return preview_bytes, (end_sec - start_sec)

def extract_audio_segment_for_embedding(audio_io, df, selected_indices):
    """選択行から音声セグメントを抽出して埋め込みを作成

    Args:
        audio_io: BytesIO - 音声ファイル
        df: DataFrame - タイムスタンプ含む
        selected_indices: list - 選択行インデックスリスト

    Returns:
        embedding: np.ndarray - 生成された埋め込みベクトル
        duration: float - セグメント長（秒）
    """
    if not selected_indices:
        raise ValueError("選択行がありません")

    selected_rows = df.iloc[sorted(selected_indices)]
    start_sec = selected_rows['start'].min()
    end_sec = selected_rows['end'].max()

    audio_io.seek(0)
    with temp_file_path(audio_io.getvalue(), ".mp3") as audio_path:
        audio = AudioSegment.from_file(audio_path)

        start_ms = int(start_sec * 1000)
        end_ms = int(end_sec * 1000)
        audio_segment = audio[start_ms:end_ms]

        wav_bytes = audio_segment.export(format="wav").read()

    embedding = _compute_embedding_from_wav_bytes(wav_bytes)

    return embedding, (end_sec - start_sec)

def get_selection_summary(df, selected_indices):
    """選択行の概要情報を取得

    Args:
        df: DataFrame
        selected_indices: set or list - 選択行インデックス

    Returns:
        dict - 概要情報
    """
    if not selected_indices:
        return {
            'count': 0,
            'start_time': None,
            'end_time': None,
            'duration': 0,
            'speakers': []
        }

    sorted_indices = sorted(selected_indices)
    selected_rows = df.iloc[sorted_indices]

    start_time = selected_rows['start'].min()
    end_time = selected_rows['end'].max()
    duration = end_time - start_time
    speakers = selected_rows['speaker'].dropna().unique().tolist() if 'speaker' in selected_rows.columns else []

    return {
        'count': len(sorted_indices),
        'start_time': start_time,
        'end_time': end_time,
        'duration': duration,
        'speakers': speakers
    }

def get_default_ragdb_path():
    """デフォルトRAGDBファイルのフルパスを取得

    Returns:
        str - RAGDBファイルのフルパス
    """
    if DEFAULT_RAGDB_FOLDER:
        # フォルダが存在しない場合は作成
        os.makedirs(DEFAULT_RAGDB_FOLDER, exist_ok=True)
        return os.path.join(DEFAULT_RAGDB_FOLDER, "default_knowledge_base.ragdb")
    else:
        return "default_knowledge_base.ragdb"

def load_meeting_type_config():
    """会議タイプ設定を取得

    Returns:
        list - 会議タイプのリスト
    """
    try:
        return [dict(item) for item in DEFAULT_MEETING_TYPES]
    except Exception as e:
        st.error(f"❌ 会議タイプ設定の初期化エラー: {e}")
        return []

def load_embeddings_from_folder(folder_path):
    """指定されたフォルダから話者埋め込みファイル（.npy）を一括読み込み

    Args:
        folder_path (str): 話者埋め込みファイルが格納されているフォルダパス

    Returns:
        list - ファイル情報のリスト [{'name': filename, 'data': file_content_bytes}, ...]
                読み込めない場合は空リスト
    """
    embeddings = []

    try:
        # 相対パスの場合は絶対パスに変換
        if not os.path.isabs(folder_path):
            script_dir = os.path.dirname(os.path.abspath(__file__))
            folder_path = os.path.join(script_dir, folder_path)

        if not os.path.exists(folder_path):
            st.warning(f"⚠️ 話者埋め込みフォルダが見つかりません: {folder_path}")
            return []

        # フォルダ内の.npyファイルを検索
        npy_files = [f for f in os.listdir(folder_path) if f.endswith('.npy')]

        if not npy_files:
            st.info(f"📁 フォルダ内に.npyファイルが見つかりませんでした: {folder_path}")
            return []

        # 各.npyファイルを読み込み
        for filename in npy_files:
            file_path = os.path.join(folder_path, filename)
            with open(file_path, 'rb') as f:
                file_bytes = f.read()
                embeddings.append({
                    'name': filename,
                    'data': file_bytes
                })

        return embeddings

    except Exception as e:
        st.error(f"❌ 話者埋め込みファイルの読み込みエラー: {e}")
        return []

def video_transcribe_and_identify():
    st.title("文字起こし")
    st.write("音声・動画ファイルから自動で文字起こしを行い、議事録を作成します。")

    st.sidebar.markdown("""
    ### 📝 文字起こし

    **概要**
    文字起こしと話者識別に対応。モデルを選択可能。

    **主な機能**
    - AI文字起こし: タイムスタンプ付き
    - モデル選択可能:
      - gpt-4o-transcribe-diarize（自動話者識別付き）
      - Whisper（参考資料対応、後から話者識別可能）
    - 議事録出力: Word・Excel形式

    **対応形式:** MP3, WAV, M4A, MP4, WebM, OGG, MOV, AVI, MKV
    """)

    # セッション状態の初期化
    _init_session_state({
        'video_combined_step': 1,
        'video_combined_audio_io': None,
        'video_combined_df': None,
        'video_combined_identified_df': pd.DataFrame(),
        'video_combined_uploaded_file_name': "",
        # 話者埋め込み作成用の状態変数
        'embedding_selected_rows': set(),
        'embedding_preview_audio': None,
        'embedding_show_preview': False
    })

    # Step 1: メディアファイルのアップロード
    st.subheader("Step 1: メディアファイルのアップロード")
    uploaded_media = st.file_uploader(
        "メディアファイルを選択してください",
        type=["mp3", "wav", "m4a", "mp4", "webm", "ogg", "mov", "avi", "mkv"],
        key="video_combined_media_upload"
    )

    if uploaded_media is not None:
        st.session_state.video_combined_uploaded_file_name = uploaded_media.name
        st.success(f"ファイル '{uploaded_media.name}' がアップロードされました")

        # メディアプレビュー（動画または音声）
        file_extension = uploaded_media.name.split(".")[-1].lower()
        if file_extension in ["mp4", "webm", "ogg", "mov", "avi", "mkv"]:
            st.video(uploaded_media)
        else:
            st.audio(uploaded_media)

        # Step 2: 文字起こし
        st.subheader("Step 2: 文字起こし")

        # モデル選択
        transcribe_model = st.selectbox(
            "文字起こしモデルを選択してください",
            options=["gpt-4o-transcribe-diarize", "whisper"],
            index=1,  # デフォルト: whisper
            key="video_combined_model_selection",
            help="gpt-4o-transcribe-diarize: 話者識別付き（参考資料非対応）\nwhisper: 話者識別なし（参考資料対応）"
        )

        # モデルの説明
        if transcribe_model == "gpt-4o-transcribe-diarize":
            st.info("🎯 **gpt-4o-transcribe-diarize**: 自動話者識別付き文字起こし。参考資料（prompt）は使用できません。")
        else:
            st.info("🎤 **Whisper**: 汎用音声認識モデル。参考資料を使用可能。Step 4で話者識別を後から実行できます。")

        # 参考資料ファイルアップロード
        reference_file = st.file_uploader(
            "参考資料ファイルを選択（オプション）",
            type=["pdf", "docx", "pptx", "txt", "msg"],
            key="video_combined_reference_file",
            help="文字起こし精度向上のための参考資料（Whisperのみサポート）"
        )

        if st.button("文字起こしを実行", key="video_combined_transcribe"):
            with st.spinner("文字起こし中..."):
                try:
                    # アップロードされたメディアファイルを直接使用
                    uploaded_media.seek(0)
                    audio_file_io = BytesIO(uploaded_media.read())
                    audio_file_io.name = uploaded_media.name

                    seg_df = transcribe_audio_to_dataframe(audio_file_io, reference_file=reference_file, model=transcribe_model)

                    st.session_state.video_combined_df = seg_df
                    st.session_state.video_combined_step = 3

                    # 音声データをセッションに保存（話者識別用）
                    uploaded_media.seek(0)
                    st.session_state.video_combined_audio_io = BytesIO(uploaded_media.read())
                    st.session_state.video_combined_audio_io.name = uploaded_media.name

                    st.success(f"文字起こしが完了しました（{len(seg_df)}行のデータ）")
                    st.rerun()
                except Exception as e:
                    st.error(f"文字起こし中にエラーが発生しました: {e}")
                    import traceback
                    st.error(traceback.format_exc())

        # Step 3: 結果の確認と編集
        if st.session_state.video_combined_step >= 3 and st.session_state.video_combined_df is not None:
            st.subheader("Step 3: 結果の確認と編集")

            if len(st.session_state.video_combined_df) > 0:
                st.write("文字起こし結果を確認・編集してください:")
                base_df = st.session_state.video_combined_df.copy()
                editable_columns = [col for col in base_df.columns if col in ("speaker",)]
                disabled_cols = [col for col in base_df.columns if col not in editable_columns]

                edited_df = st.data_editor(
                    base_df,
                    num_rows="dynamic",
                    use_container_width=True,
                    disabled=disabled_cols,
                    column_config={
                        "speaker": st.column_config.TextColumn(
                            "話者",
                            help="必要に応じて話者名を手動で調整してください"
                        )
                    },
                    key="video_combined_editor"
                )

                if isinstance(edited_df, pd.DataFrame):
                    updated_df = edited_df.copy()
                else:
                    updated_df = pd.DataFrame(edited_df)

                updated_df = updated_df.reset_index(drop=True)
                base_df = base_df.reset_index(drop=True)

                # Ensure all original columns remain
                for column in base_df.columns:
                    if column not in updated_df.columns:
                        updated_df[column] = base_df[column]
                updated_df = updated_df[base_df.columns]

                original_speaker = base_df['speaker'].astype(str).fillna("") if 'speaker' in base_df.columns else None
                updated_speaker = updated_df['speaker'].astype(str).fillna("") if 'speaker' in updated_df.columns else None

                if original_speaker is None or not updated_speaker.equals(original_speaker):
                    st.session_state.video_combined_df = updated_df
                    st.session_state.video_combined_identified_df = updated_df.copy()
                    st.info("✏️ 話者ラベルの変更を保存しました。必要に応じてStep 4以降を再実行してください。")
            else:
                st.warning("データフレームが空です。文字起こし処理でエラーが発生した可能性があります。")

        # Step 4: 話者識別
        if st.session_state.video_combined_step >= 3 and st.session_state.video_combined_df is not None:
            st.subheader("Step 4: 話者識別")

            # タブで話者識別と埋め込み作成を分離
            tab1, tab2 = st.tabs(["話者識別", "話者埋め込み作成"])

            with tab1:
                st.write("話者識別を行う範囲を選択してください:")

                # 最新の編集内容を使用
                edited_df = st.session_state.video_combined_df.copy()

                # 識別済みの話者情報を反映
                display_df = edited_df.copy()
                if not st.session_state.video_combined_identified_df.empty:
                    if 'speaker' in st.session_state.video_combined_identified_df.columns:
                        display_df['speaker'] = st.session_state.video_combined_identified_df['speaker']

                if len(display_df) == 0:
                    st.warning("文字起こし結果がありません。まず文字起こしを実行してください。")
                else:
                    row_labels = _create_row_labels(display_df)

                    col1, col2, col3 = st.columns([2, 2, 1])
                    with col1:
                        start_row = st.selectbox(
                            "開始行を選択",
                            options=range(len(display_df)),
                            format_func=lambda x: row_labels[x],
                            key="video_combined_start_row"
                        )
                    with col2:
                        end_row = st.selectbox(
                            "終了行を選択",
                            options=range(len(display_df)),
                            format_func=lambda x: row_labels[x],
                            index=len(display_df)-1 if len(display_df) > 0 else 0,
                            key="video_combined_end_row"
                        )
                    with col3:
                        similarity_threshold = st.number_input(
                            "類似度閾値",
                            min_value=0.0,
                            max_value=1.0,
                            value=0.7,
                            step=0.01,
                            key="video_combined_similarity_threshold"
                        )

                    if start_row > end_row:
                        st.error("開始行は終了行以前を選択してください")
                    else:
                        # 話者埋め込みファイルのアップロード
                        uploaded_embeddings = st.file_uploader(
                            "話者埋め込みファイルを選択してください（複数選択可）",
                            type=["npy"],
                            accept_multiple_files=True,
                            key="video_combined_embeddings"
                        )

                        if st.button("話者識別を実行", key="video_combined_identify"):
                            if not uploaded_embeddings:
                                st.error("話者埋め込みファイルをアップロードしてください")
                            else:
                                with st.spinner("話者識別中..."):
                                    try:
                                        # 選択範囲のデータフレームを抽出
                                        df_to_identify = edited_df.iloc[start_row:end_row+1].copy()

                                        # 音声ファイルのシーク位置をリセット
                                        st.session_state.video_combined_audio_io.seek(0)

                                        # 話者識別実行
                                        identified_df = identify_speakers_in_dataframe(
                                            st.session_state.video_combined_audio_io,
                                            df_to_identify,
                                            uploaded_embeddings,
                                            similarity_threshold
                                        )

                                        # 選択範囲のみ更新
                                        full_identified_df = edited_df.copy()
                                        for col in identified_df.columns:
                                            full_identified_df.loc[start_row:end_row, col] = identified_df[col].values

                                        st.session_state.video_combined_identified_df = full_identified_df
                                        st.session_state.video_combined_df = full_identified_df
                                        st.session_state.video_combined_step = 5
                                        st.success("話者識別が完了しました")
                                        st.rerun()

                                    except Exception as e:
                                        st.error(f"話者識別中にエラーが発生しました: {e}")
                                        import traceback
                                        st.error("詳細なエラー情報:")
                                        st.code(traceback.format_exc())

            with tab2:
                st.write("音声セグメントを選択し、プレビュー確認してから話者埋め込みファイルを作成できます。")

                if len(edited_df) == 0:
                    st.warning("文字起こし結果がありません。まず文字起こしを実行してください。")
                else:
                    # === 選択方式の切り替え ===
                    selection_mode = st.radio(
                        "選択方式を選んでください",
                        options=["範囲指定モード", "チェックボックスモード"],
                        horizontal=True,
                        help="範囲指定: 連続した行を素早く選択 | チェックボックス: 飛び飛びの行を自由に選択"
                    )

                    st.divider()

                    # === Part 1: 選択UI（方式に応じて変更） ===
                    st.subheader("1️⃣ 音声セグメントを選択")

                    if selection_mode == "範囲指定モード":
                        # 方式A: selectboxによる範囲指定
                        embedding_row_labels = _create_row_labels(edited_df)

                        col1, col2 = st.columns(2)
                        with col1:
                            embedding_start_row = st.selectbox(
                                "開始行を選択",
                                options=range(len(edited_df)),
                                format_func=lambda x: embedding_row_labels[x],
                                key="video_combined_embedding_start_row"
                            )
                        with col2:
                            embedding_end_row = st.selectbox(
                                "終了行を選択",
                                options=range(len(edited_df)),
                                format_func=lambda x: embedding_row_labels[x],
                                index=0,  # デフォルト: 0行目
                                key="video_combined_embedding_end_row"
                            )

                        # 範囲指定の検証とセッション状態への反映
                        if embedding_start_row > embedding_end_row:
                            st.error("⚠️ 開始行は終了行以前を選択してください")
                            st.session_state.embedding_selected_rows = set()
                        else:
                            # 範囲をsetに変換
                            st.session_state.embedding_selected_rows = set(range(embedding_start_row, embedding_end_row + 1))
                            st.success(f"✅ 行 {embedding_start_row} ～ {embedding_end_row} を選択しました")

                    else:
                        # 方式B: data_editorによるチェックボックス選択
                        display_df = edited_df.copy()
                        display_df.insert(0, "選択", False)

                        # セッションに保存されている選択行を反映
                        if st.session_state.embedding_selected_rows:
                            for idx in st.session_state.embedding_selected_rows:
                                if idx < len(display_df):
                                    display_df.at[idx, "選択"] = True

                        edited_display = st.data_editor(
                            display_df,
                            column_config={
                                "選択": st.column_config.CheckboxColumn(
                                    "選択",
                                    help="埋め込み作成対象の行をチェック",
                                    default=False
                                )
                            },
                            disabled=["start", "end", "speaker", "text"],
                            use_container_width=True,
                            hide_index=True,
                            key="embedding_data_editor"
                        )

                        # 選択状態を更新
                        st.session_state.embedding_selected_rows = set(
                            edited_display[edited_display["選択"] == True].index.tolist()
                        )

                    # === Part 2以降: 共通処理（選択がある場合のみ表示） ===
                    if st.session_state.embedding_selected_rows:
                        st.subheader("2️⃣ 選択情報の確認")

                        selection_summary = get_selection_summary(edited_df, st.session_state.embedding_selected_rows)

                        # 情報表示
                        info_cols = st.columns(4)
                        with info_cols[0]:
                            st.metric("選択行数", selection_summary['count'])
                        with info_cols[1]:
                            st.metric("開始時刻", format_time(selection_summary['start_time']))
                        with info_cols[2]:
                            st.metric("終了時刻", format_time(selection_summary['end_time']))
                        with info_cols[3]:
                            st.metric("音声長", f"{selection_summary['duration']:.1f}秒")

                        # 話者情報の表示
                        if selection_summary['speakers']:
                            speakers_text = "、".join([s if s else "不明" for s in selection_summary['speakers']])
                            st.info(f"🎤 選択範囲に含まれる話者: {speakers_text}")

                        # === Part 3: 音声プレビュー ===
                        st.subheader("3️⃣ 音声プレビュー")

                        col_preview, col_clear = st.columns([3, 1])

                        with col_preview:
                            if st.button("🔊 プレビュー音声を生成", key="generate_preview"):
                                try:
                                    with st.spinner("プレビュー音声を生成中..."):
                                        preview_bytes, duration = prepare_embedding_preview_audio(
                                            st.session_state.video_combined_audio_io,
                                            edited_df,
                                            list(st.session_state.embedding_selected_rows)
                                        )
                                        st.session_state.embedding_preview_audio = preview_bytes
                                        st.session_state.embedding_show_preview = True
                                    st.success(f"プレビュー音声を生成しました（{duration:.1f}秒）")
                                except Exception as e:
                                    st.error(f"プレビュー生成エラー: {e}")

                        with col_clear:
                            if st.session_state.embedding_show_preview:
                                if st.button("❌ プレビューをクリア", key="clear_preview"):
                                    st.session_state.embedding_preview_audio = None
                                    st.session_state.embedding_show_preview = False
                                    st.rerun()

                        # プレビュー再生
                        if st.session_state.embedding_show_preview and st.session_state.embedding_preview_audio:
                            st.audio(st.session_state.embedding_preview_audio, format="audio/wav")

                        # === Part 4: ファイル名設定 ===
                        st.subheader("4️⃣ ファイル名設定")

                        embedding_filename = st.text_input(
                            "ファイル名（.npy拡張子は自動追加）",
                            value="speaker_embedding",
                            key="video_combined_embedding_filename",
                            help="作成する話者埋め込みファイルの名前を指定してください"
                        )

                        # === Part 5: ワンクリック埋め込み作成 ===
                        st.subheader("5️⃣ 埋め込み作成")

                        if st.button("✨ 話者埋め込みを作成してダウンロード", key="create_embedding_oneclick", type="primary"):
                            with st.spinner("話者埋め込みを作成中..."):
                                try:
                                    selected_indices = sorted(list(st.session_state.embedding_selected_rows))

                                    # 埋め込みベクトルを生成
                                    embedding, duration = extract_audio_segment_for_embedding(
                                        st.session_state.video_combined_audio_io,
                                        edited_df,
                                        selected_indices
                                    )

                                    # ファイル名処理
                                    filename_with_ext = embedding_filename if embedding_filename.endswith('.npy') else f"{embedding_filename}.npy"

                                    # 埋め込みをバイト列に変換
                                    embedding_io = BytesIO()
                                    np.save(embedding_io, embedding)
                                    embedding_io.seek(0)

                                    embedding_bytes = embedding_io.getvalue()

                                    st.success(f"✅ 話者埋め込みの作成が完了しました（音声長: {duration:.1f}秒）")
                                    trigger_auto_download(
                                        embedding_bytes,
                                        filename_with_ext,
                                        key="video_combined_download_embedding",
                                        mime="application/octet-stream"
                                    )

                                except Exception as e:
                                    st.error(f"❌ 話者埋め込みの作成中にエラーが発生しました: {e}")
                                    import traceback
                                    with st.expander("🔍 詳細なエラー情報"):
                                        st.code(traceback.format_exc())

                    else:
                        if selection_mode == "範囲指定モード":
                            st.info("💡 上記のselectboxで開始行と終了行を選択してください")
                        else:
                            st.info("💡 データエディタで埋め込み作成対象の行をチェックしてください")

        # Step 5: 結果のダウンロード
        if st.session_state.video_combined_step >= 5 and not st.session_state.video_combined_identified_df.empty:
            st.subheader("Step 5: 結果のダウンロード")

            # 議事録形式のテキスト生成（連続する同じ話者の発言を結合）
            df_for_transcript = st.session_state.video_combined_identified_df.copy()

            # 話者列の前処理: 空欄を前後の話者で埋める
            df_for_transcript['speaker_filled'] = df_for_transcript['speaker'].replace('', pd.NA)
            df_for_transcript['speaker_filled'] = df_for_transcript['speaker_filled'].ffill()

            # 話者が変わるごとに新しいグループIDを割り当てる
            df_for_transcript['group_id'] = (df_for_transcript['speaker_filled'] != df_for_transcript['speaker_filled'].shift()).cumsum()

            # グループごとにテキストを結合
            df_merged = df_for_transcript.groupby('group_id').agg(
                speaker=('speaker_filled', 'first'),
                text=('text', ' '.join)
            ).reset_index(drop=True)

            # 議事録テキスト生成
            transcript_lines = []
            for idx, row in df_merged.iterrows():
                speaker = row.get('speaker', '')
                text = row.get('text', '')

                # Format: （話者）テキスト
                speaker_str = speaker if speaker else '不明'
                transcript_lines.append(f"（{speaker_str}）{text}")

            transcript_text = "\n".join(transcript_lines)

            # Word文書として生成
            doc = DocxDocument()
            doc.add_heading('議事録', 0)

            for line in transcript_text.split('\n'):
                if line.strip():
                    doc.add_paragraph(line)
                else:
                    doc.add_paragraph('')  # 空行を保持

            docx_buffer = BytesIO()
            doc.save(docx_buffer)
            docx_buffer.seek(0)

            col1, col2 = st.columns(2)
            with col1:
                st.download_button(
                    label="議事録テキストをダウンロード",
                    data=docx_buffer.getvalue(),
                    file_name=f"transcript_{st.session_state.video_combined_uploaded_file_name.split('.')[0]}.docx",
                    mime='application/vnd.openxmlformats-officedocument.wordprocessingml.document',
                    key="video_combined_download_transcript"
                )

            with col2:
                # Excelファイルとして生成
                excel_buffer = BytesIO()
                with pd.ExcelWriter(excel_buffer, engine='openpyxl') as writer:
                    st.session_state.video_combined_identified_df.to_excel(writer, index=False, sheet_name='文字起こし')
                excel_buffer.seek(0)

                st.download_button(
                    label="Excelファイルをダウンロード",
                    data=excel_buffer.getvalue(),
                    file_name=f"transcript_{st.session_state.video_combined_uploaded_file_name.split('.')[0]}.xlsx",
                    mime='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
                    key="video_combined_download_excel"
                )

def main():
    st.set_page_config(
        page_title="文字起こしアプリ",
        layout="wide",
        page_icon="🎙️",
        initial_sidebar_state="expanded"
    )

    # カスタムCSSスタイル
    st.markdown("""
        <style>
        /* メインコンテナのスタイリング */
        .main {
            background-color: #f8f9fa;
        }

        /* コンテンツエリアの背景 */
        .block-container {
            background: white;
            border-radius: 8px;
            padding: 2rem 3rem;
            margin-top: 1rem;
            box-shadow: 0 1px 3px rgba(0, 0, 0, 0.08);
        }

        /* サイドバーのスタイリング */
        section[data-testid="stSidebar"] {
            background-color: #f8f9fa;
            border-right: 1px solid #e9ecef;
        }

        /* ボタンのスタイリング */
        .stButton > button {
            background-color: #4a5568;
            color: white;
            border: none;
            border-radius: 6px;
            padding: 0.5rem 1.5rem;
            font-weight: 500;
            font-size: 0.95rem;
            transition: background-color 0.2s ease;
        }

        .stButton > button:hover {
            background-color: #2d3748;
        }

        /* プライマリボタン */
        .stButton > button[kind="primary"] {
            background-color: #3182ce;
        }

        .stButton > button[kind="primary"]:hover {
            background-color: #2c5282;
        }

        /* タイトルのスタイリング */
        h1 {
            color: #1a202c;
            font-weight: 700;
            font-size: 2rem;
            padding: 0.5rem 0;
            margin-bottom: 1rem;
            border-bottom: 2px solid #e2e8f0;
        }

        /* サブヘッダーのスタイリング */
        h2, h3 {
            color: #2d3748;
            font-weight: 600;
            margin-top: 1.5rem;
            padding-bottom: 0.5rem;
            border-bottom: 1px solid #e2e8f0;
        }

        /* 入力フィールドのスタイリング */
        .stTextInput > div > div > input,
        .stNumberInput > div > div > input,
        .stTextArea > div > div > textarea {
            border-radius: 6px;
            border: 1px solid #cbd5e0;
            padding: 0.5rem;
            transition: border-color 0.2s ease;
        }

        .stTextInput > div > div > input:focus,
        .stNumberInput > div > div > input:focus,
        .stTextArea > div > div > textarea:focus {
            border-color: #3182ce;
            box-shadow: 0 0 0 3px rgba(49, 130, 206, 0.1);
        }

        /* セレクトボックスのスタイリング */
        .stSelectbox > div > div {
            border-radius: 6px;
            border: 1px solid #cbd5e0;
        }

        /* ファイルアップローダーのスタイリング */
        .stFileUploader {
            background-color: #f7fafc;
            border-radius: 8px;
            padding: 1.5rem;
            border: 2px dashed #cbd5e0;
        }

        /* データフレームのスタイリング */
        .stDataFrame {
            border-radius: 6px;
            overflow: hidden;
            box-shadow: 0 1px 3px rgba(0, 0, 0, 0.08);
        }

        /* ダウンロードボタン */
        .stDownloadButton > button {
            background-color: #38a169;
            color: white;
            border-radius: 6px;
            font-weight: 500;
            transition: background-color 0.2s ease;
        }

        .stDownloadButton > button:hover {
            background-color: #2f855a;
        }

        /* プログレスバー */
        .stProgress > div > div > div {
            background-color: #3182ce;
        }

        /* カラムの間隔調整 */
        [data-testid="column"] {
            padding: 0.5rem;
        }

        /* エクスパンダー */
        .streamlit-expanderHeader {
            background-color: #f7fafc;
            border-radius: 6px;
            font-weight: 500;
            color: #2d3748;
        }

        /* スライダー */
        .stSlider > div > div > div {
            background-color: #3182ce;
        }

        /* タブ */
        .stTabs [data-baseweb="tab-list"] {
            gap: 4px;
        }

        .stTabs [data-baseweb="tab"] {
            background-color: #f7fafc;
            border-radius: 6px 6px 0 0;
            padding: 8px 16px;
            font-weight: 500;
            color: #4a5568;
        }

        .stTabs [aria-selected="true"] {
            background-color: #3182ce;
            color: white;
        }
        </style>
    """, unsafe_allow_html=True)

    st.sidebar.title("メニュー")
    menu_options = [
        "文字起こし",
        "🪄 dspy議事録メイカー",
        "🚀 一括処理パイプライン",
        "動画から音声を切り出しMP3で保存"
    ]
    choice = st.sidebar.selectbox("機能を選択してください", menu_options)

    if choice == "文字起こし":
        video_transcribe_and_identify()
    elif choice == "🪄 dspy議事録メイカー":
        dspy_minutes_app()
    elif choice == "🚀 一括処理パイプライン":
        batch_processing_pipeline()
    elif choice == "動画から音声を切り出しMP3で保存":
        video_to_audio_cutter_app()

if __name__ == "__main__":
    main()
